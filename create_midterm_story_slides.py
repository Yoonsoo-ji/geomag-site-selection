# -*- coding: utf-8 -*-
"""
중간보고회 LMM 본문 — 스토리 판본 (16장)
==========================================

기존 28장은 검증 기록을 빠짐없이 담은 «자료»에 가깝다. 발표 본문은 그것을
그대로 읽을 수 없으므로, **연구 흐름 하나로 꿰는 16장**을 따로 만든다.
디자인 토큰·도판·수치는 `create_midterm_lmm_slides.py` 것을 그대로 재사용한다.

    Ⅰ 왜 국가 LMM 인가      1 문제제기 · 2 층별 기여(수치) · 3 연구목표
    Ⅱ 무엇을 만들었나        4 4축 구조 · 5 자료 활용현황 · 6 파이프라인 · 7 웹 시연
    Ⅲ 어떻게 검증했나        6 검증체계 · 7 모델 선택 결과
    Ⅳ 어디까지 왔나          8 현재 성능 · 9 중요한 발견
    Ⅴ 왜 그런가              10 편각 원인진단 · 11 지각 벡터
    Ⅵ 기존 연구 대비         12 2020 연구와의 관계
    Ⅶ 중간결론·후반기        13 중간결론 · 14 후반기 3 Track

⚠ **중간보고 기준 버전(Midterm Baseline v1) 을 동결한 상태의 수치**를 쓴다.
   R0 + Crustal F 1:1(α=1) + External E2(provisional) + 지각벡터 OFF.

    python create_midterm_story_slides.py
"""
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import create_midterm_lmm_slides as B      # 디자인·도판·자료 재사용

M, CW, W, H = B.M, B.CW, B.W, B.H
BODY_Y = B.BODY_Y
NAVY, BLUE, GRAY, RED, ORANGE, GREEN = B.NAVY, B.BLUE, B.GRAY, B.RED, B.ORANGE, B.GREEN
PEACH, LBLUE, LGRAY, BAND, WHITE = B.PEACH, B.LBLUE, B.LGRAY, B.BAND, B.WHITE
slide_base, text, rect, stat, band, table = (
    B.slide_base, B.text, B.rect, B.stat, B.band, B.table)

LOO, DIAG, MODEL = B.LOO, B.DIAG, B.MODEL
N_SITES, KPI_D, KPI_F = B.N_SITES, B.KPI_D, B.KPI_F
ASY = B.ASY
D_MIN = LOO["D"] * 60      # 분


PAGES = "https://yoonsoo-ji.github.io/geomag-site-selection"

VAL = {(r["성분"], r["단계"]): r["RMS"] for r in MODEL["validation"]}

# 2020 연구 15점 ↔ ’22~25 측량 15점 — 명단·배치를 하드코딩하지 않는다
import legacy2020_sets as LS                                   # noqa: E402

SET20 = LS.sets()
COV20 = LS.coverage(LS.coords(SET20["s2020"]))
COV22 = LS.coverage(LS.coords(SET20["s2225"], prefer_model=True))


def link(slide, x, y, w, h, label, url, size=11.5, color=BLUE):
    """클릭 가능한 링크 상자 — 발표 중 그대로 눌러 띄운다."""
    box = text(slide, x, y, w, h, label, size, color, True)
    r = box.text_frame.paragraphs[0].runs[0]
    r.hyperlink.address = url
    return box


# ══════════════════════════════════════════════════════ Ⅰ. 왜 국가 LMM 인가
def s01_problem(prs, page):
    s = slide_base(prs, "Ⅰ. 연구 배경",
                   "전지구 모델만으로는 국내 지표 자기장을 설명하지 못한다",
                   "IGRF 는 지구 중심핵이 만드는 장을 기술한다. 지표에서 실제로 재는 "
                   "값과의 차이가 이 연구의 대상이다.", page)
    y = BODY_Y
    rows = [
        ["① 핵 (Core)", "지구 중심핵", "3,000 km 이상", "IGRF 가 담당", "—"],
        ["② 지역 (Regional)", "IGRF 가 놓치는 광역 편차", "50~3,000 km",
         "지표 절대측정 필요", "반복관측점"],
        ["③ 지각 (Crustal)", "암석 자화", "0.05~50 km",
         "항공자력 자료 필요", "KIGAM 자력이상도"],
        ["④ 외부 (External)", "전리층·자기권 전류", "일·시·분",
         "상시관측 필요", "청양·KASA"],
    ]
    y = table(s, M, y, CW, ["층", "기원", "공간규모", "필요 자료", "국내 보유"],
              rows, [0.16, 0.24, 0.15, 0.24, 0.21], row_h=0.52)
    band(s, M, y + 0.30, CW, 1.30,
         "국가기본도의 자침편차 표기는 ② 와 ③ 까지 맞아야 성립한다.\n\n"
         "· 지구물리측량 작업규정 제12조 — 1등 지자기점은 국가기본도의 자편각 "
         "표기를 위한 점\n"
         "· 목표 정확도 편각 0.1°(6′) 는 IGRF 단독으로는 도달하지 못한다")
    return s


def s02_goal(prs, page):
    s = slide_base(prs, "Ⅰ. 연구 배경",
                   "15점으로는 부족하다 — 관측망과 모델을 함께 세운다",
                   "2020 연구가 「15점으로 전국 지자기도를 만들기에는 부족하다」고 "
                   "판단한 지점에서 출발한다.", page)
    y = BODY_Y
    bw = (CW - 0.36) / 2
    for i, (tt, col, lines) in enumerate([
        ("종전 (2020)", GRAY,
         "· 반복측량 15점\n· 전국 지도는 IGRF 30점 계산값에\n  청양 상시관측 보정을 더해 작성\n"
         "· 층별 분리·교차검증 체계 없음"),
        ("이번 연구", BLUE,
         "· 4축(Core·Regional·Crustal·External) 결합\n"
         "· 국내 실측자료로 각 층을 직접 구성\n· 독립 교차검증으로 성능 판정\n"
         "· 관측망 50점 확대 설계")]):
        x = M + i * (bw + 0.36)
        rect(s, x, y, bw, 1.90, LGRAY if i == 0 else LBLUE)
        text(s, x + 0.24, y + 0.16, bw - 0.48, 0.32, tt, 14, col, True)
        text(s, x + 0.24, y + 0.58, bw - 0.48, 1.20, lines, 11.5, NAVY,
             space=1.45)
        if i == 0:
            text(s, x + bw + 0.04, y + 0.80, 0.28, 0.34, "▶", 18, BLUE, True, "c")
    band(s, M, y + 2.16, CW, 1.34,
         "중간보고 시점의 위치\n\n"
         "① 자료 통합 → ② 4축 시범구축 → ③ 독립검증(현재) → "
         "④ 원인별 개선 → ⑤ 신규 50점 최종검증\n"
         "지금은 ③ 의 후반부다. 만들고, 재고, 한계를 특정하는 단계까지 왔다.")
    return s


def s01a_why_lmm(prs, page):
    """왜 층을 쌓는가 — 각 층이 실제로 줄인 몫을 수치로."""
    s = slide_base(prs, "Ⅰ. 연구 배경",
                   "층을 쌓는 이유 — 각 층이 실제로 줄이는 몫이 있다",
                   "국내 실측 성과에 대해 IGRF 단독부터 단계별로 잔차를 쟀다. "
                   "개념이 아니라 측정값이다.", page, rule=GREEN)
    y = BODY_Y
    rows = [
        ["총자력 F", f"{VAL[('F_nT','IGRF')]:.0f} nT",
         f"{VAL[('F_nT','+Crustal')]:.0f} nT",
         f"{VAL[('F_nT','+Crustal+Regional')]:.0f} nT",
         f"−{100*(1-VAL[('F_nT','+Crustal+Regional')]/VAL[('F_nT','IGRF')]):.0f} %"],
        ["편각 D", f"{VAL[('D_deg','IGRF')]*60:.1f}′", "—",
         f"{VAL[('D_deg','+Regional')]*60:.1f}′",
         f"−{100*(1-VAL[('D_deg','+Regional')]/VAL[('D_deg','IGRF')]):.0f} %"],
        ["복각 I", f"{VAL[('I_deg','IGRF')]*60:.1f}′", "—",
         f"{VAL[('I_deg','+Regional')]*60:.1f}′",
         f"−{100*(1-VAL[('I_deg','+Regional')]/VAL[('I_deg','IGRF')]):.0f} %"],
    ]
    y = table(s, M, y, 7.90,
              ["성분", "① Core 단독", "+ ③ Crustal", "+ ② Regional", "감소"],
              rows, [0.18, 0.20, 0.20, 0.22, 0.20], row_h=0.52)
    band(s, 8.35, BODY_Y, CW + M - 8.35, 2.10,
         "읽는 법\n\n"
         "· 총자력은 지각 자기이상이 대부분을\n  걷어낸다 (87 → 59 nT)\n\n"
         "· 편각·복각은 지각층이 스칼라만 주므로\n  ② 지역장이 담당한다\n\n"
         "· 그래도 편각은 목표 6′ 에 못 미친다", "23456B")
    band(s, M, y + 0.30, CW, 1.42,
         "이것이 「왜 LMM 인가」에 대한 답이다.\n\n"
         "· IGRF 는 전지구 핵 자기장 모델이라 국내 지표값과 총자력에서 87 nT, "
         "편각에서 32′ 어긋난다\n"
         "· 층을 쌓으면 그 몫이 실제로 줄어든다 — 다만 남는 편각 오차가 이 연구의 "
         "과제다", GREEN)
    return s


def s05a_web(prs, page):
    """웹 시연 — 링크를 걸어 발표 중 바로 띄운다."""
    s = slide_base(prs, "Ⅱ. 구축 내용",
                   "구축 결과는 웹에서 바로 확인할 수 있다",
                   "설치 없이 브라우저만으로 열린다. 아래 주소를 눌러 시연한다.",
                   page, rule=BLUE)
    items = [
        ("LMM 계산기", "lmm.html",
         "좌표·표고·일자를 넣으면 편각·복각·총자력을 계산한다.\n"
         "층별 기여(Core·Regional·Crustal)와 정확도 고지를 함께 보여준다."),
        ("한국 지자기도", "lmm_map.html",
         "IGRF · LMM · 차이를 나란히 그린다.\n"
         "2020 연구의 지자기도와 같은 형식이라 직접 견줄 수 있다."),
        ("3D 지구 자기장", "geomag_globe.html",
         "전지구 자기장 위에 한국 고해상 패치를 얹는다.\n"
         "선점 후보·기존점 마커에서 상세 지도로 이어진다."),
        ("선점 검토 지도", "survey_review.html",
         "현장조사 103점의 등급·사진·방위표지와\n기존 측정점 30점을 함께 본다."),
    ]
    y = BODY_Y
    for i, (name, page_, desc) in enumerate(items):
        col = 0 if i % 2 == 0 else 1
        row = i // 2
        x = M + col * (CW / 2 + 0.20)
        yy = y + row * 1.62
        w = CW / 2 - 0.20
        rect(s, x, yy, w, 1.42, LGRAY)
        rect(s, x, yy, 0.06, 1.42, BLUE)
        text(s, x + 0.24, yy + 0.14, w - 0.48, 0.30, name, 13.5, NAVY, True)
        link(s, x + 0.24, yy + 0.48, w - 0.48, 0.26,
             f"{PAGES}/{page_}", f"{PAGES}/{page_}", 9.5)
        text(s, x + 0.24, yy + 0.80, w - 0.48, 0.52, desc, 10.5, GRAY,
             space=1.35)
    band(s, M, y + 3.30, CW, 1.10,
         "시연 순서 — ① 계산기에서 한 점을 계산해 층별 기여를 보여주고 "
         "② 지자기도에서 IGRF 와의 차이를 보인 뒤\n"
         "③ 지구본에서 한국 패치를 확대해 선점 후보까지 이어 간다. "
         "모두 외부 서버 없이 동작한다.")
    return s


# ══════════════════════════════════════════════════════ Ⅱ. 무엇을 만들었나
def s03_structure(prs, page):
    s = slide_base(prs, "Ⅱ. 구축 내용",
                   "4축 결합 — 층마다 자료원과 시간·공간 규모가 다르다",
                   "층을 나누는 이유는 편의가 아니라 물리다. 기원이 다르므로 "
                   "단순 내삽으로 분리할 수 없다.", page)
    text(s, M, BODY_Y, CW, 0.46,
         "B_LMM (r, t)  =  B_IGRF  +  B_Regional  +  B_Crustal  +  B_External",
         17, NAVY, True, "c")
    img = B.fig_bands()
    s.shapes.add_picture(str(img), Inches(M), Inches(BODY_Y + 0.62),
                         width=Inches(7.4))
    band(s, 7.95, BODY_Y + 0.62, CW - 7.50, 3.10,
         "층마다 다른 것\n\n"
         "· 시간상수 — 핵은 수십 년,\n  외부장은 분\n\n"
         "· 공간파장 — 핵은 3,000 km 이상,\n  지각은 수 km\n\n"
         "· 자료원 — 전지구 모델 / 지표 절대측정 /\n  항공자력 / 상시관측", "23456B")
    return s


def s04_data(prs, page):
    """핵심 표 — 4축 자료를 어디까지 실제로 썼는가."""
    ext = B.EXT
    s = slide_base(prs, "Ⅱ. 구축 내용",
                   "네 축을 모두 국내 실측자료로 구성했다",
                   "개념도가 아니라 실제 자료를 연결해 시험한 결과다.", page,
                   rule=GREEN)
    rows = [
        ["① Core", "IGRF-14 (13차)", "1900~2030", "구현 완료",
         "전 연도 단일 판본 적용 · 공식 계산기와 자릿수 내 일치"],
        ["② Regional", f"지리원 반복측량 (’19 · ’22~’25)", f"{N_SITES}측점",
         "구현·검증", "R0(상수항) 채택 · 편각 잔차가 크다"],
        ["③ Crustal", "KIGAM 자력이상 (1982~2018)", "2.8 km 격자",
         "F 적용 · D/I 진단", "1:1 결합이 교차검증 최적"],
        ["④ External", "청양 + KASA 3소 + Kp", f"{ext['ok']}세션",
         "시범 적용·검증", "4소 공간보간이 단일 관측소보다 우수"],
    ]
    y = table(s, M, BODY_Y, CW,
              ["축", "사용 자료", "규모", "활용 수준", "중간 결론"],
              rows, [0.10, 0.26, 0.12, 0.16, 0.36], row_h=0.56,
              head_fill=GREEN)
    stat(s, M, y + 0.28, 3.90, 1.10, "4 / 4", "국내 자료로 구성한 축", LBLUE,
         GREEN, 24)
    stat(s, M + 4.05, y + 0.28, 3.90, 1.10, f"{ext['y0']}~{ext['y1']}",
         "관측시각을 복원한 구간", LGRAY, BLUE, 22)
    band(s, M + 8.10, y + 0.28, CW - 8.10, 1.10,
         "성과표에는 연도만 있었다\n야장 원본에서 분 단위 시각을 복원했다", "23456B")
    return s


def s05_pipeline(prs, page):
    s = slide_base(prs, "Ⅱ. 구축 내용",
                   "구축 절차 — 자료 이력 감사부터 시작한다",
                   "모델 알고리즘보다 순서가 중요하다. 측점 이력과 관측시각을 "
                   "먼저 정리해야 뒤 단계가 의미를 갖는다.", page)
    steps = [("0", "이력 감사", "표지 재설·장비·등급"),
             ("1", "시각 통일", "야장 KST → UTC"),
             ("2", "세션 QC", "관측소 변동폭·Kp"),
             ("3", "외부장 제거", "4소 공간보간"),
             ("4", "IGRF 계산", "실제 관측일시"),
             ("5", "LMM 적합", "Regional + Crustal"),
             ("6", "교차검증", "Station-LOSO")]
    bw = (CW - 6 * 0.10) / 7
    y = BODY_Y + 0.10
    for i, (no, tt, sub) in enumerate(steps):
        x = M + i * (bw + 0.10)
        rect(s, x, y, bw, 1.28, LGRAY)
        rect(s, x, y, bw, 0.34, NAVY)
        text(s, x, y + 0.03, bw, 0.28, f"{no}단계", 10.5, WHITE, True, "c")
        text(s, x + 0.08, y + 0.44, bw - 0.16, 0.30, tt, 11.5, NAVY, True, "c")
        text(s, x + 0.08, y + 0.78, bw - 0.16, 0.44, sub, 9.5, GRAY, align="c",
             space=1.25)
        if i < 6:
            text(s, x + bw - 0.02, y + 0.52, 0.14, 0.26, "›", 13, "9FB4C7",
                 True, "c")
    y += 1.56
    rows = [
        ["0단계가 새로 들어갔다",
         "같은 이름이라도 표지가 재설되면 다른 측점이다. 2020 보고서의 신뢰도 "
         "판정(2017 장비고장·2011 D/I 미관측·표석 망실)을 이력표로 만들어 등급을 매겼다"],
        ["시각은 야장에서 복원했다",
         "성과표에는 연도만 있으나 야장에는 분 단위 시각이 있다. 편각·복각 관측 "
         "구간이 따로 기록돼 있어 성분마다 자기 구간의 외부장으로 보정한다"],
        ["층은 따로 손보지 않는다",
         "Regional 은 Core·Crustal 이 남긴 것을 흡수하는 항이므로, 자료가 바뀌면 "
         "전 층을 다시 적합한다"],
    ]
    table(s, M, y, CW, ["요점", "내용"], rows, [0.24, 0.76], row_h=0.62)
    return s


# ══════════════════════════════════════════════════════ Ⅲ. 어떻게 검증했나
def s06_valid_system(prs, page):
    s = slide_base(prs, "Ⅲ. 검증",
                   "만들고 끝내지 않고, 무엇이 맞는지 자료로 골랐다",
                   "구현 검증 · 구조 선택 · 예측 검증의 세 층으로 나누고, 평가 "
                   "측점을 미리 얼렸다.", page, rule=GREEN)
    cols = [
        ("구현 검증", "계산이 설계대로 되는가",
         "· 공식 IGRF 계산기와 대조\n· Python·JS·Excel 세 구현 일치\n"
         "· 표시 자릿수 내 일치 확인"),
        ("구조 선택", "어떤 형태가 옳은가",
         "· Regional 차수 R0/R1/R2\n· 지각 결합 F0/F1/Fα\n· External E0/E1/E2"),
        ("예측 검증", "쓰지 않은 점에서 맞는가",
         "· Station-LOSO\n· 블록 홀드아웃(남중북·동서)\n· 2010 공표성과 대조"),
    ]
    bw = (CW - 2 * 0.24) / 3
    for i, (tt, sub, body) in enumerate(cols):
        x = M + i * (bw + 0.24)
        rect(s, x, BODY_Y, bw, 2.10, WHITE, line=NAVY)
        rect(s, x, BODY_Y, bw, 0.40, NAVY)
        text(s, x, BODY_Y + 0.05, bw, 0.30, tt, 12.5, WHITE, True, "c")
        text(s, x + 0.18, BODY_Y + 0.52, bw - 0.36, 0.26, sub, 10, GRAY, True)
        text(s, x + 0.18, BODY_Y + 0.86, bw - 0.36, 1.10, body, 10.5, NAVY,
             space=1.45)
    y = BODY_Y + 2.36
    band(s, M, y, CW, 1.40,
         "평가 측점을 얼렸다 — 이 연구의 방법론적 핵심\n\n"
         "여러 후보를 같은 교차검증 수치를 보며 고르면 그 수치에도 과적합된다. "
         "그래서 기준 설정에서\n평가 측점을 한 번만 정해 파일로 고정하고"
         f"(편각 14 · 복각 15 · 총자력 13 / {N_SITES}측점), 이후 모든 비교를 "
         "그 위에서만 했다.", GREEN)
    return s


def s07_selection(prs, page):
    s = slide_base(prs, "Ⅲ. 검증",
                   "세 갈래를 자료로 갈랐다 — 복잡한 쪽이 이기지 않았다",
                   "모두 같은 조건(Grade A · 동일 차수 · 평가 측점 고정)에서 처음부터 다시 "
                   "적합했다. 「가장 가까운 관측소와 변화가 같다」는 가정은 성립하지 "
                   "않았다.", page, rule=GREEN)

    img = B.fig_selection()
    s.shapes.add_picture(str(img), Inches(M), Inches(BODY_Y),
                         width=Inches(7.00))
    band(s, 7.80, BODY_Y, CW + M - 7.80, 2.36,
         "선택의 요점\n\n"
         "· 측점 16개로는 «공간 구조»를\n  지지할 수 없다 → 상수항(R0)\n\n"
         "· 지각 결합계수를 자유롭게 두면\n  과적합 → 1:1 (α = 1)\n\n"
         "복잡한 모형이 훈련 오차만 낮추고\n예측을 악화시키는 사례를 확인했다",
         "23456B")

    y = BODY_Y + 2.52
    text(s, M, y, CW, 0.30,
         "External — 단일 관측소로는 안 되고, 공간보간이라야 된다", 12.5, NAVY, True)
    imgE = B.fig_external()
    s.shapes.add_picture(str(imgE), Inches(M), Inches(y + 0.32),
                         width=Inches(4.60))
    rows3 = [
        ["미적용", "—", "0.5123°", "기준선"],
        ["E0", "청양 단독", "0.5495°", "기각"],
        ["E1", "최근접 관측소", "0.5504°", "기각"],
        ["E2", "4소 공간보간", "0.4987°", "채택(잠정)"],
        ["E2b", "4소 역거리가중", "0.5595°", "기각"],
    ]
    table(s, 5.40, y + 0.34, CW + M - 5.40,
          ["모델", "자료원", "LOO 편각", "판정"], rows3,
          [0.14, 0.40, 0.26, 0.20], row_h=0.31, size=10)
    return s


# ══════════════════════════════════════════════════════ Ⅳ. 어디까지 왔나
def s08_perf(prs, page):
    s = slide_base(prs, "Ⅳ. 현재 성능",
                   "총자력은 목표에 근접했고, 편각이 남았다",
                   "적합에 쓰지 않은 측점에서의 예측오차(Station-LOSO)다.",
                   page)
    img = B.fig_loo()
    s.shapes.add_picture(str(img), Inches(8.20), Inches(BODY_Y + 0.05),
                         width=Inches(4.53))
    stat(s, M, BODY_Y, 2.42, 1.34, f"{D_MIN:.0f}′",
         f"편각 D  |  목표 6′", PEACH, RED)
    stat(s, M + 2.55, BODY_Y, 2.42, 1.34, f"{LOO['F']:.0f} nT",
         f"총자력 F  |  목표 {KPI_F:g} nT", PEACH, ORANGE)
    stat(s, M + 5.10, BODY_Y, 2.42, 1.34, f"{LOO['I']*60:.0f}′",
         "복각 I  |  별도 목표 없음", LBLUE, BLUE)
    y = BODY_Y + 1.62
    rows = [
        ["총자력 F", f"{LOO['F']:.1f} nT", f"{KPI_F:g} nT",
         f"목표의 {LOO['F']/KPI_F:.2f}배", "근접"],
        ["편각 D", f"{D_MIN:.1f}′", "6′", f"목표의 {D_MIN/6:.1f}배", "구조적 개선 필요"],
        ["복각 I", f"{LOO['I']*60:.1f}′", "—", "—", "검증 중"],
    ]
    table(s, M, y, 7.45, ["성분", "현재", "목표", "배수", "상태"], rows,
          [0.16, 0.18, 0.14, 0.20, 0.32], row_h=0.46, size=10)
    band(s, M, y + 1.70, 7.45, 1.10,
         "중간보고 기준 버전 (Midterm Baseline v1)\n"
         "Regional R0 · 지각 1:1 결합 · External E2(잠정) · 지각벡터 미적용")
    band(s, 8.20, BODY_Y + 2.62, 4.53, 1.62,
         "현 단계 사용 범위\n\n"
         "· 연구·시범계산 : 가능\n"
         "· 지형도 자침편차 등 정식 편각 산출 : 자료 확충 후", "23456B")
    return s


def s09_finding(prs, page):
    s = slide_base(prs, "Ⅳ. 현재 성능",
                   "중간 시점에서 얻은 것 — 무엇이 되고 무엇이 남았는지 특정했다",
                   "성능 미달 자체보다, 그 원인을 자료·물리·구조로 나눠 지목한 "
                   "것이 이 단계의 성과다.", page, rule=GREEN)
    y = BODY_Y
    bw = (CW - 0.30) / 2
    rect(s, M, y, bw, 2.30, LBLUE)
    text(s, M + 0.24, y + 0.18, bw - 0.48, 0.34, "확인된 것", 14, GREEN, True)
    text(s, M + 0.24, y + 0.62, bw - 0.48, 1.60,
         "· 4축 구조를 국내 자료로 구현했다\n"
         "· Core 구현은 공식 계산기와 일치한다\n"
         "· 지각 자력이상을 총자력에 넣으면 실제로 좋아진다\n"
         "  (72.3 → 63.0 nT)\n"
         "· 외부장은 4소 공간보간에서만 이득이 난다\n"
         "· 총자력은 목표의 1.3배까지 왔다", 11.5, NAVY, space=1.5)
    rect(s, M + bw + 0.30, y, bw, 2.30, PEACH)
    text(s, M + bw + 0.54, y + 0.18, bw - 0.48, 0.34, "남은 것", 14, RED, True)
    text(s, M + bw + 0.54, y + 0.62, bw - 0.48, 1.60,
         f"· 편각이 목표의 {D_MIN/6:.0f}배다 ({D_MIN:.0f}′ vs 6′)\n"
         "· 측점 16개로는 지역장의 공간 구조를\n  결정할 수 없다\n"
         "· 지각 벡터는 자료 결측 39% 로 보류\n"
         "· 상시관측소 축 규약 문서가 없다", 11.5, NAVY, space=1.5)
    band(s, M, y + 2.56, CW, 1.10,
         "이 단계의 결론은 「미달」이 아니라 「특정」이다.\n"
         "무엇을 고치면 무엇이 좋아지는지가 정해졌고, 그것이 후반기 연구계획이 된다.",
         GREEN)
    return s


# ══════════════════════════════════════════════════════ Ⅴ. 왜 그런가
def s10_why_D(prs, page):
    A = DIAG["hypotheses"]["A_vector_D"]
    s = slide_base(prs, "Ⅴ. 원인 진단",
                   "편각 오차의 원인을 좁혔다 — 모형이 아니라 자료다",
                   "모형 쪽 후보는 검정으로 하나씩 지웠다. 남은 것은 관측 자료의 "
                   "계통 오차와 측점 밀도다.", page, rule=RED)
    y = BODY_Y
    rows = [
        ["IGRF 구현 오류", "공식 계산기와 표시 자릿수 내 일치", "배제"],
        ["Regional 차수 부족", "1·2차가 상수항보다 나쁘다", "배제"],
        ["지역 영년변화 항 누락", "시간항 추가 시 오히려 악화", "배제"],
        ["단일 관측소 외부장 보정", "청양 단독·최근접 모두 기준선보다 나쁨", "배제"],
        ["Kp>2 자료 제거", "관측소 실측이 조용한 세션이 대부분", "근거 부족"],
        ["지각 결합계수 오류", "계수 자유추정은 과적합(fold 0.29~0.81)", "배제"],
    ]
    y = table(s, M, y, 7.55, ["검정한 가설", "결과", "판정"], rows,
              [0.34, 0.48, 0.18], row_h=0.42, size=10)
    rect(s, 8.30, BODY_Y, CW + M - 8.30, 2.66, PEACH)
    text(s, 8.50, BODY_Y + 0.16, CW + M - 8.70, 0.32, "남은 요인", 13.5, RED, True)
    text(s, 8.50, BODY_Y + 0.58, CW + M - 8.70, 2.00,
         "① 기존 반복측량의 편각 계통오차\n"
         "   장비 · 방위표지 · 측점 이설 ·\n   관측 캠페인\n\n"
         "② 지각 벡터 성분 미반영\n"
         f"   예상 기여 {DIAG['vector_recovery']['predD_rms_arcmin']:.1f}′ "
         "(KPI 6′ 초과)\n\n"
         "③ 측점 16개라는 공간 밀도", 11, NAVY, space=1.5)
    band(s, M, y + 0.28, 7.55, 1.22,
         "방위 기준이 실제로 흔들린다 — 재방문 감사에서 확인\n\n"
         "같은 측점인데 마크 참방위각이 방문마다 옮겨 간다(예: 159.9° → 330.6°). "
         "재방문 잔여가\n부호까지 뒤집힌 방문 2건은 원인이 특정돼 표본에서 제외했다.",
         RED)
    return s


def s11_vector(prs, page):
    VR = DIAG["vector_recovery"]
    s = slide_base(prs, "Ⅴ. 원인 진단",
                   "개선 여지는 이미 확인했다 — 다만 자료가 아직 못 받친다",
                   "항공자력의 스칼라 이상에서 벡터를 복원하면 편각이 일관되게 "
                   "좋아진다. 자료 결측이 걸림돌이다.", page, rule=BLUE)
    text(s, M, BODY_Y, CW, 0.40,
         "ΔF  =  l · b_north  +  m · b_east  +  n · b_down       "
         "→ 역산하면 편각·복각에 쓸 수 있다", 14, NAVY, True, "c")
    imgV = B.fig_vector_robust()
    s.shapes.add_picture(str(imgV), Inches(M), Inches(BODY_Y + 0.60),
                         width=Inches(6.60))
    text(s, M, BODY_Y + 3.72, 6.60, 0.28,
         "미적용 0.4987° · zero 0.4506 · mean 0.4507 · harmonic 0.4739 · "
         "경계배제 0.4638 · 블록 0.4560", 9.5, GRAY, align="c")
    band(s, M, BODY_Y + 4.06, 6.60, 0.58,
         "요지는 「미완」이 아니라 「개선 경로가 이미 확인됐다」는 것이다.")
    stat(s, 7.15, BODY_Y + 0.56, 2.60, 1.20,
         f"{VR['predD_rms_arcmin']:.1f}′", "복원 벡터의 편각 기여", LBLUE, BLUE, 22)
    stat(s, 9.95, BODY_Y + 0.56, CW + M - 9.95, 1.20,
         f"{VR['grid']['gap_pct']:.0f} %", "격자 결측 — 0 으로 채움", PEACH, RED, 22)
    band(s, 7.15, BODY_Y + 1.92, CW + M - 7.15, 2.34,
         "왜 아직 넣지 않는가\n\n"
         "모든 채움 방식과 블록 홀드아웃에서\n편각이 좋아진다 — 신호는 실재한다.\n\n"
         "그러나 결측 39% 를 0 으로 채운 역산이라\n복원 진폭 자체가 불확실하다.\n\n"
         "승격 조건은 채움 알고리즘 교체가 아니라\n원자료 결측의 해결이다.", "23456B")
    return s


# ══════════════════════════════════════════════════════ Ⅵ. 기존 연구 대비
def s12_vs2020(prs, page):
    s = slide_base(prs, "Ⅵ. 기존 연구와의 관계",
                   "2020 연구가 멈춘 지점에서 이어 간다 — 다만 같은 15점이 아니다",
                   "2020 연구를 대체하는 것이 아니라, 그 판단을 출발점으로 삼아 "
                   "구조를 확장했다. 측점망은 거의 교체됐다.", page)
    y = BODY_Y
    bw = (CW - 0.40) / 2
    rect(s, M, y, bw, 2.30, LGRAY)
    text(s, M + 0.24, y + 0.16, bw - 0.48, 0.32, "2020 연구", 14, GRAY, True)
    text(s, M + 0.24, y + 0.58, bw - 0.48, 1.60,
         "2회 이상 측량된 15점\n        ↓\n전국 모델에는 부족하다고 판단\n        ↓\n"
         "IGRF 30점 계산값 + 청양 상시관측 보정\n        ↓\n"
         "2020 지자기도", 10.5, NAVY, space=1.24)
    rect(s, M + bw + 0.40, y, bw, 2.30, LBLUE)
    text(s, M + bw + 0.64, y + 0.16, bw - 0.48, 0.32, "이번 연구", 14, BLUE, True)
    text(s, M + bw + 0.64, y + 0.58, bw - 0.48, 1.60,
         "IGRF (Core) + 지표 절대측정 (Regional)\n        +\n"
         "KIGAM 자력이상 (Crustal) + 다중 관측소 (External)\n        ↓\n"
         "통합 LMM + 독립 교차검증", 10.5, NAVY, space=1.24)
    text(s, M + bw + 0.08, y + 1.00, 0.30, 0.34, "▶", 20, BLUE, True, "c")
    rows = [
        ["측점 구성", "2012~19 측량 2회 이상 15점",
         f"공통 {len(SET20['common'])}점 · 2020 배제분 {len(SET20['only2225'])}점 "
         f"· 야장 2점 = {N_SITES}점"],
        ["관측 공백", f"최대 {COV20['max']:.0f} km · 50 km 초과 {COV20['over50']:.0f} %",
         f"최대 {COV22['max']:.0f} km · {COV22['over50']:.0f} % — "
         "강원·충청 내륙 축 상실 (50점 배치 설계의 근거)"],
        ["지역장", "별도 층 없음 — IGRF 계산값 사용",
         "지표 절대측정으로 직접 구성(R0)"],
        ["지각장", "반영 없음",
         f"KIGAM 자력이상 1:1 결합 — 총자력 잔차 "
         f"{VAL[('F_nT','IGRF')]:.0f} → {VAL[('F_nT','+Crustal')]:.0f} nT"],
        ["외부장", "청양 단독 보정", "관측소 4소 공간보간 · 성분별 관측구간(잠정)"],
        ["검증", "과거 성과와 비교", "평가집합 동결 + Station-LOSO + 블록 홀드아웃"],
    ]
    table(s, M, y + 2.44, CW, ["구분", "2020 연구", "이번 연구"], rows,
          [0.14, 0.34, 0.52], row_h=0.36, size=9.5)
    return s


# ══════════════════════════════════════════════════════ Ⅶ. 결론·후반기
def s13_conclusion(prs, page):
    s = slide_base(prs, "Ⅶ. 중간결론",
                   "중간결론 — 만들었고, 쟀고, 한계를 특정했다", "", page,
                   rule=GREEN)
    items = [
        ("①", "국내 자료 기반 4축 LMM 시범구축 완료",
         "Core·Regional·Crustal·External 을 모두 국내 실측자료로 구성하고 "
         "세 계산환경(Python·웹·엑셀)에서 같은 값을 확인했다"),
        ("②", "독립 교차검증으로 적용 가능성과 한계를 정량 확인",
         f"평가 측점을 동결한 Station-LOSO 로 편각 {D_MIN:.0f}′ · "
         f"총자력 {LOO['F']:.0f} nT 를 얻었다. 총자력은 목표의 1.3배까지 왔다"),
        ("③", "편각 정확도 제한요인을 특정하여 후반기 방향 확정",
         "모형 쪽 후보는 검정으로 배제했고, 기존 관측성과의 계통오차·측점 밀도·"
         "지각 벡터 자료 결측이 남았다"),
    ]
    y = BODY_Y + 0.10
    for no, tt, body in items:
        rect(s, M, y, CW, 1.16, LGRAY)
        text(s, M + 0.22, y + 0.14, 0.50, 0.44, no, 22, GREEN, True)
        text(s, M + 0.86, y + 0.16, CW - 1.10, 0.34, tt, 14, NAVY, True)
        text(s, M + 0.86, y + 0.58, CW - 1.10, 0.50, body, 11, NAVY, space=1.4)
        y += 1.30
    band(s, M, y + 0.06, CW, 0.92,
         "현 단계 산출물은 연구·시범계산용이다. 지형도 자침편차 등 정식 편각 "
         "산출에는 자료 확충 후 적용한다.", GREEN)
    return s


def s14_next(prs, page):
    s = slide_base(prs, "Ⅶ. 후반기 계획",
                   "후반기는 세 갈래로 나눠 진행한다",
                   "기존자료 고도화 · 모델 고도화 · 신규 관측망을 병렬로 두고, "
                   "마지막에 독립검증으로 합친다.", page, rule=ORANGE)
    tracks = [
        ("Track A", "기존자료 고도화", ORANGE,
         "· 2012~2019 ↔ 2022~25 동일측점 대조\n· 장비·방위표지·측점 이력 추적\n"
         "· 오차 budget 작성\n· 상시관측소 축 규약 문서 확보"),
        ("Track B", "모델 고도화", BLUE,
         "· KIGAM 결측 해결 → 지각 벡터 재판정\n· External E2 확정\n"
         "· 지역장 공간구조 재검토\n· 목표 편각 6′ · 총자력 50 nT"),
        ("Track C", "신규 50점 관측망", GREEN,
         "· 자기환경 조사 기반 선점\n· 35~40점 모델 구축용\n"
         "· 10~15점 검증 전용\n· 최종 독립검증의 근거"),
    ]
    bw = (CW - 2 * 0.24) / 3
    for i, (no, tt, col, body) in enumerate(tracks):
        x = M + i * (bw + 0.24)
        rect(s, x, BODY_Y, bw, 2.30, WHITE, line=col)
        rect(s, x, BODY_Y, bw, 0.44, col)
        text(s, x, BODY_Y + 0.06, bw, 0.32, f"{no}   {tt}", 12, WHITE, True, "c")
        text(s, x + 0.20, BODY_Y + 0.60, bw - 0.40, 1.56, body, 11, NAVY,
             space=1.5)
    band(s, M, BODY_Y + 2.56, CW, 1.34,
         "신규 50점의 설계가 가장 중요하다\n\n"
         "전부를 모델 적합에 쓰면 최종 성능을 객관적으로 말할 수 없다. "
         "일부를 처음부터 검증 전용으로\n남겨 두면 최종 보고서에서 "
         "「모델 구축에 사용하지 않은 신규 국가기준점에서 독립검증하였다」고 "
         "쓸 수 있다.", ORANGE)
    return s


# ══════════════════════════════════════════════════════ 실행
def main():
    sys.stdout.reconfigure(encoding="utf-8")
    B._mpl_init()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    builders = [s01_problem, s01a_why_lmm, s02_goal,
                s03_structure, s04_data, s05_pipeline, s05a_web,
                s06_valid_system, s07_selection, s08_perf, s09_finding,
                s10_why_D, s11_vector, s12_vs2020, s13_conclusion, s14_next]
    for i, fn in enumerate(builders, 1):
        fn(prs, i)
        print(f"  [{i:02d}] {fn.__name__}")

    B.audit(prs)
    B.OUT_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = B.OUT_DIR / f"{ts}_중간보고회_LMM_본문_{len(builders)}장.pptx"
    prs.save(out)
    print(f"\n저장: {out}")
    return out


if __name__ == "__main__":
    main()
