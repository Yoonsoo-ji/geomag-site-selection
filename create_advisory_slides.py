# -*- coding: utf-8 -*-
"""
자문회의 발표자료 — 연구 진행현황 종합 (입지 선정 + LMM).  Bauhaus Geometry 판.

    python create_advisory_slides.py
        -> docs/output/YYYYMMDD_HHMMSS_지자기_자문회의_발표자료.pptx

디자인 방침 (Style D — Bauhaus Geometry, 학술 절제 변형)
  · 배경 흰색 · 기본 차콜/네이비 · 강조 코발트 1 + 오렌지 1 (4색 제한)
  · 도형은 원 · 삼각형 · 사각형만. 그라디언트 · 그림자 · 라운드 · 유기곡선 금지
  · 면마다 도형 하나는 화면 밖으로 흘린다(bleed). 대각 구성.
  · 제목은 짧은 결론형 문장. 본문 산문 대신 도식 내부의 짧은 라벨
  · 연구 논리 자체를 도형으로 — 인과구조도 · 깔때기 · 레이어 · 파이프라인 · 매트릭스

글꼴 (설치본 기준 치환)
  · 디스플레이 숫자/라틴  Bahnschrift  (DIN 계열 — Bebas Neue 대체)
  · 한글 제목             Noto Sans KR Black
  · 한글 본문             Noto Sans KR
  · 모노 라벨             Consolas     (Space Mono 대체, 라틴·숫자 전용)

수치는 하드코딩하지 않고 아래에서 읽는다 (재실행 시 자동 갱신).
    docs/data/candidates_p1|p2|p3.geojson · lmm_model.json · lmm_diagnosis.json
    lmm_build.igrf_dif (축척별 도폭 내 편각 변화)
"""
import glob
import json
import math
import tempfile
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
DOCS_OUT = HERE / "docs" / "output"
DOCS_DATA = HERE / "docs" / "data"
MODEL_JSON = DOCS_DATA / "lmm_model.json"
DIAG_JSON = DOCS_DATA / "lmm_diagnosis.json"
BOUNDARY = HERE / "data" / "korea_boundary.geojson"
ASSETS = Path(tempfile.gettempdir()) / "advisory_bauhaus_figs"

# ── 팔레트 (백지 + 4색) ───────────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x1A, 0x1A, 0x1A)   # 차콜 — 글
NAVY    = RGBColor(0x16, 0x32, 0x4F)   # 구조 도형
COBALT  = RGBColor(0x1D, 0x4E, 0xD8)   # 강조 1 — 데이터
ORANGE  = RGBColor(0xE8, 0x53, 0x1F)   # 강조 2 — 미달·되먹임
GRAY    = RGBColor(0x6E, 0x76, 0x7E)   # 보조 글
FAINT   = RGBColor(0xEE, 0xF0, 0xF3)   # 연회색 면
LINE    = RGBColor(0xC9, 0xCE, 0xD4)   # 괘선

HX_INK, HX_NAVY = "#1A1A1A", "#16324F"
HX_COBALT, HX_ORANGE = "#1D4ED8", "#E8531F"
HX_GRAY, HX_FAINT = "#6E767E", "#E8EAEE"

# ── 글꼴 짝 (라틴, 한글) ──────────────────────────────────────────────────
F_DISP = ("Bahnschrift", "Noto Sans KR Black")
F_BODY = ("Bahnschrift", "Noto Sans KR")
F_MONO = ("Consolas", "Noto Sans KR")

SW, SH = 13.333, 7.5
M = 0.9
CW = SW - 2 * M            # 11.533
FOOT_Y = 6.94              # 꼬리말 괘선 — 본문은 이 위에서 끝낸다

DECK_TITLE = "GEOMAG REFERENCE PROJECT — ADVISORY 2026-07"
N_PAGES = 22


# ══════════════════════════════════════════════════════════════════════
# 자료 적재
# ══════════════════════════════════════════════════════════════════════
def load_candidates():
    out = {}
    for tier in (1, 2, 3):
        gj = json.loads((DOCS_DATA / f"candidates_p{tier}.geojson")
                        .read_text(encoding="utf-8"))
        out[tier] = [f["properties"] for f in gj["features"]]
    return out


def load_model():
    return json.loads(MODEL_JSON.read_text(encoding="utf-8"))


def load_diag():
    return json.loads(DIAG_JSON.read_text(encoding="utf-8"))


def scale_rows():
    """축척별 도폭 내 편각 변화 — IGRF-14 로 직접 계산."""
    import datetime as dt

    import numpy as np

    from lmm_build import igrf_dif

    d = dt.datetime(2027, 1, 1)
    rows = []
    for name, step in [("1:5,000", 0.025), ("1:10,000", 0.05),
                       ("1:25,000", 0.125), ("1:50,000", 0.25)]:
        var = []
        for la0 in np.arange(34.0, 38.0, 0.5):
            for lo0 in np.arange(126.5, 129.0, 0.5):
                la = np.array([la0, la0 + step, la0, la0 + step])
                lo = np.array([lo0, lo0, lo0 + step, lo0 + step])
                D, *_ = igrf_dif(la, lo, np.zeros(4), d)
                var.append(D.max() - D.min())
        rows.append({"name": name, "dD": float(np.median(var))})
    return rows


# ══════════════════════════════════════════════════════════════════════
# 도판 (matplotlib — 사스 세리프 · 코발트/오렌지 재채색)
# ══════════════════════════════════════════════════════════════════════
def _mpl():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager

    for pat in (r"C:\Windows\Fonts\NotoSansKR*",):
        for f in glob.glob(pat):
            try:
                font_manager.fontManager.addfont(f)
            except Exception:
                pass
    have = {f.name for f in font_manager.fontManager.ttflist}
    for cand in ("Noto Sans KR", "Malgun Gothic", "sans-serif"):
        if cand in have or cand == "sans-serif":
            plt.rcParams["font.family"] = cand
            break
    plt.rcParams.update({
        "axes.unicode_minus": False,
        "axes.edgecolor": HX_GRAY,
        "axes.linewidth": 0.8,
        "xtick.color": HX_GRAY,
        "ytick.color": HX_GRAY,
        "xtick.labelsize": 9,
        "ytick.labelsize": 9,
        "figure.facecolor": "white",
        "savefig.facecolor": "white",
    })
    return plt


def fig_buffers(path):
    """제외구역 이격 반경 — 코발트 막대 + 오렌지 법정 눈금."""
    plt = _mpl()

    items = [
        ("주거·취락", 0.30, None), ("핵심도심·산업", 0.50, None),
        ("파이프라인", 0.50, None), ("통신탑·기지국", 0.50, None),
        ("풍력발전기", 0.50, None), ("단층선", 0.50, None),
        ("송전탑", 1.00, 0.5), ("고압철탑", 1.00, 1.0),
        ("채석장·광산", 1.00, None), ("철도", 5.00, 2.0),
    ]
    fig, ax = plt.subplots(figsize=(7.0, 3.55))
    ys = range(len(items))
    ax.barh(list(ys), [it[1] for it in items], height=0.55,
            color=HX_COBALT, edgecolor="none")
    for y, (label, applied, legal) in zip(ys, items):
        ax.text(applied * 1.09, y, f"{applied:g}", va="center", ha="left",
                fontsize=9.5, color=HX_INK, family="Consolas")
        if legal is not None:
            ax.plot([legal, legal], [y - 0.36, y + 0.36],
                    color=HX_ORANGE, lw=2.4, solid_capstyle="butt", zorder=6)
    ax.set_yticks(list(ys))
    ax.set_yticklabels([it[0] for it in items], fontsize=10.5)
    ax.set_xscale("log")
    ax.set_xlim(0.2, 14)
    ax.set_xticks([0.2, 0.5, 1, 2, 5, 10])
    ax.set_xticklabels(["0.2", "0.5", "1", "2", "5", "10"])
    ax.set_xlabel("이격 반경 (km, 로그)", fontsize=9.5, labelpad=6)
    ax.plot([], [], color=HX_ORANGE, lw=2.4, label="법정 최소")
    ax.legend(fontsize=9, loc="lower right", frameon=False)
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    ax.tick_params(axis="y", length=0)
    ax.grid(axis="x", color=HX_FAINT, lw=0.7)
    ax.set_axisbelow(True)
    fig.subplots_adjust(left=0.185, right=0.985, top=0.985, bottom=0.13)
    fig.savefig(path, dpi=300)
    plt.close(fig)


def fig_scores(path, cands):
    """종합 점수 분포 — 코발트 계열, 오렌지 분위 경계."""
    plt = _mpl()
    import numpy as np

    allsc = np.array([p["score"] for t in (1, 2, 3) for p in cands[t]])
    p33, p66 = np.percentile(allsc, [33, 66])

    fig, ax = plt.subplots(figsize=(5.4, 2.0))
    bins = np.linspace(allsc.min(), allsc.max(), 34)
    for tier, alpha, label in ((3, 0.25, "3등급"), (2, 0.52, "2등급"),
                               (1, 1.00, "1등급")):
        ax.hist([p["score"] for p in cands[tier]], bins=bins, color=HX_COBALT,
                alpha=alpha, edgecolor="white", linewidth=0.4, label=label)
    top = ax.get_ylim()[1]
    for v, tag, dx, dy, ha in ((p33, "P33", -4, -6, "right"),
                               (p66, "P66", 4, -24, "left")):
        ax.axvline(v, color=HX_ORANGE, lw=1.2, ls=(0, (4, 3)))
        ax.annotate(f"{tag} {v:.1f}", xy=(v, top), xytext=(dx, dy),
                    textcoords="offset points", fontsize=8.5,
                    color=HX_ORANGE, va="top", ha=ha, family="Consolas")
    ax.legend(fontsize=8.5, frameon=False, loc="upper right")
    for side in ("top", "right", "left"):
        ax.spines[side].set_visible(False)
    ax.set_yticks([])
    ax.grid(axis="x", color=HX_FAINT, lw=0.7)
    ax.set_axisbelow(True)
    fig.subplots_adjust(left=0.02, right=0.99, top=0.95, bottom=0.13)
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return float(p33), float(p66)


def _boundary_rings():
    gj = json.loads(BOUNDARY.read_text(encoding="utf-8"))
    rings = []
    for feat in gj["features"]:
        geom = feat["geometry"]
        polys = (geom["coordinates"] if geom["type"] == "MultiPolygon"
                 else [geom["coordinates"]])
        for poly in polys:
            rings.append(poly[0])
    return rings


def fig_candidates_map(path, cands):
    plt = _mpl()

    fig, ax = plt.subplots(figsize=(4.55, 5.3))
    for ring in _boundary_rings():
        ax.fill([p[0] for p in ring], [p[1] for p in ring],
                facecolor=HX_FAINT, edgecolor="#B4BBc2", lw=0.6)
    styles = {3: (12, 0.22, "3등급"), 2: (18, 0.45, "2등급"),
              1: (34, 1.00, "1등급 P1")}
    for tier in (3, 2, 1):
        size, alpha, label = styles[tier]
        pts = cands[tier]
        ax.scatter([p["lon"] for p in pts], [p["lat"] for p in pts],
                   s=size, facecolor=HX_COBALT, alpha=alpha,
                   edgecolor="white", linewidth=0.4, zorder=5, label=label)
    ax.set_xlim(125.0, 131.2)
    ax.set_ylim(33.0, 38.8)
    ax.set_aspect(1.0 / math.cos(math.radians(36.0)))
    ax.axis("off")
    leg = ax.legend(fontsize=9.5, loc="upper left", frameon=False,
                    labelspacing=0.8, borderpad=0.1, handletextpad=0.5)
    leg._legend_box.align = "left"
    fig.subplots_adjust(left=0.01, right=0.99, top=0.99, bottom=0.01)
    fig.savefig(path, dpi=300)
    plt.close(fig)


def fig_subgrid(path):
    """서브격자 선정 개념도 — 실제 좌표가 아니다."""
    plt = _mpl()
    import numpy as np

    fig, ax = plt.subplots(figsize=(4.3, 4.3))
    R = 5.0
    th = np.linspace(0, 2 * math.pi, 400)
    ax.plot(R * np.cos(th), R * np.sin(th), color="#B4BBC2", lw=1.1)
    ax.text(0, R + 0.45, "R = 5 KM", ha="center", fontsize=9,
            color=HX_GRAY, family="Consolas")
    for cx, cy, r in ((-2.6, 2.2, 1.9), (2.9, -2.4, 2.3)):
        ax.add_patch(plt.Circle((cx, cy), r, facecolor=HX_ORANGE, alpha=0.12,
                                edgecolor=HX_ORANGE, lw=0.9, ls=(0, (3, 2))))
    xs = np.arange(-5, 5.1, 1.0)
    keep, drop = [], []
    for x in xs:
        for y in xs:
            if x * x + y * y > R * R:
                continue
            excluded = ((x + 2.6) ** 2 + (y - 2.2) ** 2 < 1.9 ** 2 or
                        (x - 2.9) ** 2 + (y + 2.4) ** 2 < 2.3 ** 2)
            (drop if excluded else keep).append((x, y))
    ax.scatter([p[0] for p in drop], [p[1] for p in drop], s=14,
               facecolor="white", edgecolor="#B4BBC2", linewidth=0.9, zorder=4)
    ax.scatter([p[0] for p in keep], [p[1] for p in keep], s=14,
               facecolor="#9FB6D8", edgecolor="none", zorder=4)
    keepset = {(round(x, 3), round(y, 3)) for x, y in keep}
    top5 = [p for p in ((-1, -1), (0, 1), (2, 0), (0, -3), (-3, -2))
            if (round(float(p[0]), 3), round(float(p[1]), 3)) in keepset]
    ax.scatter([p[0] for p in top5], [p[1] for p in top5], s=78,
               facecolor=HX_COBALT, edgecolor="white", linewidth=1.0, zorder=6)
    ax.scatter([0], [0], marker="+", s=150, color=HX_ORANGE, linewidth=2.0,
               zorder=7)
    ax.set_xlim(-6.3, 6.3)
    ax.set_ylim(-6.5, 6.6)
    ax.set_aspect("equal")
    ax.axis("off")
    handles = [
        plt.Line2D([], [], marker="+", ls="", color=HX_ORANGE, ms=10, mew=2.0,
                   label="사이트 중심"),
        plt.Line2D([], [], marker="o", ls="", mfc="#9FB6D8", mec="none", ms=5,
                   label="격자 통과"),
        plt.Line2D([], [], marker="o", ls="", mfc="white", mec="#B4BBC2",
                   ms=5, label="제외구역 저촉"),
        plt.Line2D([], [], marker="o", ls="", mfc=HX_COBALT, mec="white",
                   ms=8.5, label="상위 5개"),
    ]
    ax.legend(handles=handles, fontsize=9, frameon=False, loc="lower left",
              bbox_to_anchor=(-0.04, -0.04), labelspacing=0.55,
              handletextpad=0.4)
    fig.subplots_adjust(left=0.005, right=0.995, top=0.995, bottom=0.005)
    fig.savefig(path, dpi=300)
    plt.close(fig)


def fig_sites(path, model):
    plt = _mpl()
    sites = model["sites"]

    fig, ax = plt.subplots(figsize=(4.2, 4.9))
    for ring in _boundary_rings():
        ax.fill([p[0] for p in ring], [p[1] for p in ring],
                facecolor=HX_FAINT, edgecolor="#B4BBC2", lw=0.6)
    ax.scatter([s["lon"] for s in sites], [s["lat"] for s in sites],
               s=[24 + 16 * s["n_visit"] for s in sites], facecolor=HX_COBALT,
               edgecolor="white", linewidth=0.8, zorder=5)
    ax.set_xlim(125.0, 130.2)
    ax.set_ylim(33.0, 38.8)
    ax.set_aspect(1.0 / math.cos(math.radians(36.0)))
    ax.axis("off")
    for v in (2, 4):
        ax.scatter([], [], s=24 + 16 * v, facecolor=HX_COBALT,
                   edgecolor="white", linewidth=0.8, label=f"{v}회")
    leg = ax.legend(title="반복 관측", fontsize=9, title_fontsize=9,
                    loc="upper left", frameon=False, labelspacing=0.95,
                    borderpad=0.1)
    leg._legend_box.align = "left"
    fig.subplots_adjust(left=0.01, right=0.99, top=0.99, bottom=0.01)
    fig.savefig(path, dpi=300)
    plt.close(fig)


def fig_residuals(path, diag):
    plt = _mpl()
    from matplotlib.patches import Patch

    rows = sorted(diag["site_table"], key=lambda r: r["rD"])
    names = [r["name"] for r in rows]
    vals = [r["rD"] for r in rows]
    inl = [bool(r["inlier_D"]) for r in rows]

    fig, ax = plt.subplots(figsize=(10.9, 2.0))
    xs = range(len(rows))
    ax.bar(list(xs), vals, width=0.6,
           color=[HX_COBALT if k else "#C4CAD1" for k in inl],
           edgecolor="none")
    ax.axhspan(-6.0, 6.0, color=HX_ORANGE, alpha=0.14, zorder=0)
    ax.axhline(0, color=HX_GRAY, lw=0.9)
    ax.set_xticks(list(xs))
    ax.set_xticklabels(names, fontsize=9.5)
    ax.set_ylabel("편각 잔차 (′)", fontsize=9.5, labelpad=4)
    ax.set_xlim(-0.7, len(rows) - 0.3)
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    ax.tick_params(axis="x", length=0)
    ax.grid(axis="y", color=HX_FAINT, lw=0.7)
    ax.set_axisbelow(True)
    ax.legend(handles=[Patch(facecolor=HX_ORANGE, alpha=0.14,
                             label="KPI ±6′ (=0.1°)"),
                       Patch(facecolor="#C4CAD1", label="품질 미달 배제")],
              fontsize=8.5, frameon=False, loc="upper left", ncol=2,
              handlelength=1.3, columnspacing=1.3)
    fig.subplots_adjust(left=0.056, right=0.995, top=0.97, bottom=0.20)
    fig.savefig(path, dpi=300)
    plt.close(fig)


def fig_variogram(path, diag):
    plt = _mpl()

    v = diag["variogram"]
    bins = v["bins"]
    mids = [(b["lo"] + b["hi"]) / 2 for b in bins]
    gam = [b["semivariance"] for b in bins]

    fig, ax = plt.subplots(figsize=(4.9, 2.3))
    ax.scatter(mids, gam, s=[18 + 3.4 * b["n"] for b in bins],
               facecolor=HX_COBALT, edgecolor="white", linewidth=0.8, zorder=5)
    ax.axhline(v["sill"], color=HX_ORANGE, lw=1.3, ls=(0, (4, 3)))
    ax.annotate(f"SILL {v['sill']:,.0f}", xy=(max(mids), v["sill"]),
                xytext=(0, 6), textcoords="offset points", ha="right",
                fontsize=8.5, color=HX_ORANGE, family="Consolas")
    ax.set_xlabel("측점 간 거리 (km)", fontsize=9.5, labelpad=5)
    ax.set_xlim(0, max(mids) * 1.12)
    ax.set_ylim(0, max(gam) * 1.22)
    ax.set_yticks([])
    for side in ("top", "right", "left"):
        ax.spines[side].set_visible(False)
    ax.grid(axis="x", color=HX_FAINT, lw=0.7)
    ax.set_axisbelow(True)
    fig.subplots_adjust(left=0.02, right=0.99, top=0.94, bottom=0.22)
    fig.savefig(path, dpi=300)
    plt.close(fig)


def fig_scale(path, rows):
    plt = _mpl()

    names = [r["name"] for r in rows]
    vals = [r["dD"] for r in rows]

    fig, ax = plt.subplots(figsize=(5.4, 2.6))
    bars = ax.bar(names, vals, width=0.52,
                  color=[HX_COBALT if r["name"] == "1:25,000" else "#B7C3D6"
                         for r in rows], edgecolor="none")
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.004,
                f"{100 * v / 0.1:.0f}%", ha="center", fontsize=10.5,
                color=HX_INK, family="Consolas")
    ax.axhline(0.1, color=HX_ORANGE, lw=1.3, ls=(0, (4, 3)))
    ax.annotate("0.1° BUDGET", xy=(len(rows) - 0.5, 0.1), xytext=(0, 6),
                textcoords="offset points", ha="right", fontsize=8.5,
                color=HX_ORANGE, family="Consolas")
    ax.set_ylim(0, 0.128)
    ax.set_yticks([])
    ax.tick_params(axis="x", length=0, labelsize=10)
    for side in ("top", "right", "left"):
        ax.spines[side].set_visible(False)
    fig.subplots_adjust(left=0.02, right=0.99, top=0.95, bottom=0.12)
    fig.savefig(path, dpi=300)
    plt.close(fig)


def build_figures(cands, model, diag, scales):
    ASSETS.mkdir(parents=True, exist_ok=True)
    p = {k: ASSETS / f"bh_{k}.png" for k in
         ("buffers", "scores", "cmap", "subgrid", "sites", "resid",
          "vario", "scale")}
    fig_buffers(p["buffers"])
    pcts = fig_scores(p["scores"], cands)
    fig_candidates_map(p["cmap"], cands)
    fig_subgrid(p["subgrid"])
    fig_sites(p["sites"], model)
    fig_residuals(p["resid"], diag)
    fig_variogram(p["vario"], diag)
    fig_scale(p["scale"], scales)
    return p, pcts


# ══════════════════════════════════════════════════════════════════════
# 조판 헬퍼
# ══════════════════════════════════════════════════════════════════════
def _fonts(run, fonts):
    latin, ea = fonts
    run.font.name = latin
    rPr = run.font._rPr
    lat = rPr.find(qn("a:latin"))
    ea_el = rPr.find(qn("a:ea"))
    if ea_el is None:
        ea_el = rPr.makeelement(qn("a:ea"), {"typeface": ea})
        lat.addnext(ea_el)
    else:
        ea_el.set("typeface", ea)


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def run_in(p, text, size, color, *, fonts=F_BODY, bold=False, spacing=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    _fonts(r, fonts)
    if spacing is not None:                      # 자간 (1/100 pt)
        r.font._rPr.set("spc", str(spacing))
    return r


def para(tf, text, size, color, *, fonts=F_BODY, bold=False, line=1.35,
         space_before=0, align=None, first=False, spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.line_spacing = line
    p.space_before = Pt(space_before)
    run_in(p, text, size, color, fonts=fonts, bold=bold, spacing=spacing)
    return p


def geo(slide, kind, x, y, w, h, *, fill=None, line=None, weight=1.0, rot=0):
    """기하 도형 — 원 · 삼각형 · 사각형만. 그림자·라운드 없음."""
    shp = {"rect": MSO_SHAPE.RECTANGLE, "oval": MSO_SHAPE.OVAL,
           "tri": MSO_SHAPE.ISOSCELES_TRIANGLE}[kind]
    sh = slide.shapes.add_shape(shp, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(weight)
    sh.shadow.inherit = False
    style = sh._element.find(qn("p:style"))
    if style is not None:
        sh._element.remove(style)      # 테마 effectRef 까지 떼어야 그림자가 사라진다
    sh.rotation = rot
    sh.text_frame.word_wrap = True
    return sh


def seg(slide, x1, y1, x2, y2, *, color=INK, weight=1.0, dash=None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),
                                    Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    if dash is not None:                      # "dash" 등 prstDash 값
        el = ln.line._get_or_add_ln()
        el.append(el.makeelement(qn("a:prstDash"), {"val": dash}))
    return ln


def arrow(slide, x1, y1, x2, y2, *, color=INK, weight=1.2):
    cn = seg(slide, x1, y1, x2, y2, color=color, weight=weight)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"),
                             {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def zigzag(slide, x, y, w, n, amp, *, color=INK, weight=1.4):
    """지그재그 — 파장 표현. 유기곡선 대신 직선 분절."""
    step = w / n
    for i in range(n):
        y1 = y + (amp if i % 2 == 0 else -amp)
        y2 = y + (amp if i % 2 == 1 else -amp)
        seg(slide, x + i * step, y1, x + (i + 1) * step, y2,
            color=color, weight=weight)


def chrome(slide, page, eyebrow):
    """상단 모노 라벨 + 색 마커, 하단 러닝타이틀 + 쪽번호."""
    tf = textbox(slide, M, 0.30, 8.0, 0.24)
    para(tf, eyebrow, 9.5, GRAY, fonts=F_MONO, line=1.1, first=True,
         spacing=120)
    for i, c in enumerate((NAVY, COBALT, ORANGE)):
        geo(slide, "rect", SW - M - 0.54 + i * 0.20, 0.32, 0.13, 0.13, fill=c)
    seg(slide, M, FOOT_Y, SW - M, FOOT_Y, color=LINE, weight=0.75)
    tf = textbox(slide, M, FOOT_Y + 0.09, 8.0, 0.24)
    para(tf, DECK_TITLE, 8, GRAY, fonts=F_MONO, line=1.1, first=True,
         spacing=100)
    tf = textbox(slide, SW - M - 1.2, FOOT_Y + 0.09, 1.2, 0.24,
                 align=PP_ALIGN.RIGHT)
    para(tf, f"{page:02d} / {N_PAGES}", 8, GRAY, fonts=F_MONO, line=1.1,
         first=True)


def title_block(slide, num, runs, *, y=0.78):
    """결론형 제목 — 네이비 사각 번호 + 문장. runs = [(글, 색), ...]"""
    if num:
        sq = geo(slide, "rect", M, y, 0.46, 0.46, fill=NAVY)
        tfq = sq.text_frame
        tfq.margin_left = tfq.margin_right = 0
        tfq.margin_top = tfq.margin_bottom = 0
        tfq.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tfq.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run_in(p, num, 17, WHITE, fonts=F_MONO, bold=True)
        tx = M + 0.66
    else:
        geo(slide, "rect", M, y + 0.05, 0.16, 0.36, fill=ORANGE)
        tx = M + 0.36
    tf = textbox(slide, tx, y + 0.015, SW - M - tx, 0.5)
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    for text, color in runs:
        run_in(p, text, 24, color, fonts=F_DISP, bold=True)
    return y + 0.66


def cap(slide, x, y, w, fig_no, text):
    """캡션 — 모노 번호 + 본문."""
    tf = textbox(slide, x, y, w, 0.3)
    p = tf.paragraphs[0]
    p.line_spacing = 1.25
    run_in(p, fig_no + "  ", 8.5, COBALT, fonts=F_MONO, bold=True)
    run_in(p, text, 9, GRAY, fonts=F_BODY)
    return tf


def stat(slide, x, y, w, value, unit, label, sub=None, *, color=NAVY,
         size=37):
    """숫자 블록 — 색 사각 + Bahnschrift 큰 수."""
    geo(slide, "rect", x, y, 0.30, 0.055, fill=color)
    tf = textbox(slide, x, y + 0.13, w, size / 72 * 1.25)
    p = tf.paragraphs[0]
    p.line_spacing = 1.0
    run_in(p, value, size, color, fonts=F_DISP, bold=True)
    if unit:
        run_in(p, " " + unit, size * 0.38, GRAY, fonts=F_MONO)
    ly = y + 0.13 + size / 72 * 1.08
    tfl = textbox(slide, x, ly, w, 0.55)
    para(tfl, label, 11, INK, fonts=F_BODY, bold=True, line=1.25, first=True)
    if sub:
        para(tfl, sub, 8.5, GRAY, fonts=F_MONO, line=1.25, space_before=3)
    return ly


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ══════════════════════════════════════════════════════════════════════
# 표제 · 부 표지 · 대시보드
# ══════════════════════════════════════════════════════════════════════
def s_cover(prs):
    s = blank(prs)

    # 기하 구성 — 원 하나는 화면 밖으로
    geo(s, "oval", 9.4, -2.3, 5.8, 5.8, fill=NAVY)
    geo(s, "rect", 8.35, 4.05, 3.4, 0.85, fill=COBALT)
    geo(s, "tri", 10.55, 2.55, 1.35, 1.2, fill=ORANGE, rot=-15)
    seg(s, 7.3, 6.9, 12.9, 1.35, color=INK, weight=1.0)

    tf = textbox(s, M, 0.36, 8.0, 0.26)
    para(tf, "NGII · GEOMAGNETIC REFERENCE PROJECT", 10, GRAY, fonts=F_MONO,
         line=1.1, first=True, spacing=140)

    geo(s, "rect", M, 1.78, 0.95, 0.10, fill=ORANGE)
    tf = textbox(s, M, 2.06, 7.6, 0.34)
    para(tf, "지구자기장 측정 기준 정비", 14, GRAY, line=1.3, first=True)
    tf = textbox(s, M, 2.50, 8.2, 2.2)
    para(tf, "연구 진행현황과", 44, INK, fonts=F_DISP, bold=True, line=1.14,
         first=True)
    para(tf, "자문 요청사항", 44, INK, fonts=F_DISP, bold=True, line=1.14)

    parts = [("01", "측정 입지 선정", NAVY),
             ("02", "지역 자기장 모형 LMM", COBALT),
             ("03", "자문 질의사항", ORANGE)]
    for i, (n, name, c) in enumerate(parts):
        x = M + i * 2.75
        geo(s, "rect", x, 5.48, 0.15, 0.15, fill=c)
        tf = textbox(s, x + 0.27, 5.42, 2.4, 0.26)
        para(tf, n, 10, c, fonts=F_MONO, bold=True, line=1.15, first=True)
        tf = textbox(s, x + 0.27, 5.70, 2.4, 0.3)
        para(tf, name, 11.5, INK, bold=True, line=1.2, first=True)

    tf = textbox(s, M, 6.66, 8.0, 0.26)
    para(tf, "ADVISORY MEETING — " + datetime.now().strftime("%Y-%m"),
         9.5, GRAY, fonts=F_MONO, line=1.1, first=True, spacing=120)

    s.notes_slide.notes_text_frame.text = (
        "두 축은 독립된 사업이 아니다. 선점으로 확보되는 신규 측점이 "
        "LMM 의 Regional 층 입력이 된다.")
    return s


def s_dashboard(prs, ctx):
    s = blank(prs)
    chrome(s, 2, "OVERVIEW")

    # 흘림 도형 — 좌측 연회색 원
    geo(s, "oval", -1.7, 3.1, 3.4, 3.4, fill=FAINT)

    title_block(s, None, [("절차는 검증을 마쳤다 — ", INK),
                          ("남은 제약은 입력자료다", COBALT)])

    cands, model, diag = ctx["cands"], ctx["model"], ctx["diag"]
    cv = model["loo_cv"]
    n_all = sum(len(cands[t]) for t in (1, 2, 3))

    tiles = [
        (f"{n_all}", "개소", "제외구역 통과 후보지", f"P1 = {len(cands[1])}", NAVY),
        ("103", "개소", "도상선점 확정", "SURVEY CARDS ISSUED", NAVY),
        ("496", "점", "매칭 국가기준점", "733 → FILTERED", NAVY),
        (f"{len(model['sites'])}", "점", "지표 절대측정 측점",
         f"OBS {diag['dataset']['n_obs']} · 2019/2022-25", NAVY),
        (f"{cv['D']:.2f}", "°", "교차검증 편각 LOO", "KPI 0.1° — 미달", ORANGE),
        (f"{cv['F']:.0f}", "nT", "교차검증 총자력 LOO", "KPI 50 nT — 미달", ORANGE),
    ]
    colw = (CW - 2 * 0.5) / 3
    for i, (val, unit, label, sub, c) in enumerate(tiles):
        x = M + (i % 3) * (colw + 0.5)
        y = 1.78 + (i // 3) * 2.02
        stat(s, x, y, colw, val, unit, label, sub, color=c)

    geo(s, "rect", M, 6.06, CW, 0.56, fill=NAVY)
    tf = textbox(s, M + 0.26, 6.20, CW - 0.5, 0.3)
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    run_in(p, "RESULT  ", 9.5, ORANGE, fonts=F_MONO, bold=True, spacing=120)
    run_in(p, "구현 오류 0 건 — 성능을 묶는 것은 관측시각 · 원측선 해상도 · "
              "측점 수 세 가지다.", 11.5, WHITE, bold=True)
    return s


def s_causal(prs, ctx):
    """인과구조도 — 현황 → 문제 → 원인 → 연구의 두 축."""
    s = blank(prs)
    chrome(s, 3, "WHY THIS PROJECT")

    geo(s, "oval", 11.9, 5.3, 3.2, 3.2, fill=FAINT)   # 흘림 도형

    title_block(s, None, [("문제의 답이 ", INK), ("연구의 두 축", COBALT),
                          ("이다", INK)])

    boxes = [
        ("01 현황", "기존 측정점 33개 (2010–25)\n자침편차는 전역모형(IGRF) 의존",
         M + 0.0, 4.85),
        ("02 문제", "간섭원 급증 · 관측 공백\n편각 정확도 KPI 0.1° 미달",
         M + 3.05, 3.65),
        ("03 원인", "선점 기준 부재\n관측시각 미기록 · 지역 모형 부재",
         M + 6.10, 2.45),
    ]
    bw, bh = 2.65, 1.30
    for label, body, x, y in boxes:
        geo(s, "rect", x, y, bw, bh, fill=WHITE, line=INK, weight=1.2)
        geo(s, "rect", x, y, 0.5, 0.14, fill=INK)
        tf = textbox(s, x + 0.16, y + 0.22, bw - 0.32, 0.24)
        para(tf, label, 9.5, GRAY, fonts=F_MONO, bold=True, line=1.15,
             first=True, spacing=80)
        tf = textbox(s, x + 0.16, y + 0.50, bw - 0.32, 0.72)
        para(tf, body, 10, INK, line=1.32, first=True)

    arrow(s, M + bw + 0.06, 4.85 + bh / 2 - 0.4, M + 3.05 - 0.06,
          3.65 + bh / 2 + 0.28, color=INK, weight=1.4)
    arrow(s, M + 3.05 + bw + 0.06, 3.65 + bh / 2 - 0.4, M + 6.10 - 0.06,
          2.45 + bh / 2 + 0.28, color=INK, weight=1.4)

    # 두 축 — 결과 블록
    res = [("PART 01", "측정 입지 선정", "어디에서 잴 것인가", NAVY, 1.62),
           ("PART 02", "지역 자기장 모형", "어떻게 면(面)으로 만들 것인가",
            COBALT, 3.42)]
    rx, rw, rh = M + 9.35, 3.05, 1.42
    for tag, name, q, c, ry in res:
        geo(s, "rect", rx, ry, rw, rh, fill=c)
        tf = textbox(s, rx + 0.22, ry + 0.20, rw - 0.44, 0.24)
        para(tf, tag, 9, WHITE, fonts=F_MONO, bold=True, line=1.15,
             first=True, spacing=120)
        tf = textbox(s, rx + 0.22, ry + 0.48, rw - 0.44, 0.36)
        para(tf, name, 15, WHITE, fonts=F_DISP, bold=True, line=1.15,
             first=True)
        tf = textbox(s, rx + 0.22, ry + 0.92, rw - 0.44, 0.32)
        para(tf, q, 9.5, WHITE, line=1.25, first=True)

    ox = M + 6.10 + bw
    arrow(s, ox + 0.06, 2.45 + bh / 2 - 0.25, rx - 0.08, 1.62 + rh / 2 + 0.15,
          color=NAVY, weight=1.6)
    arrow(s, ox + 0.06, 2.45 + bh / 2 + 0.25, rx - 0.08, 3.42 + rh / 2 - 0.15,
          color=COBALT, weight=1.6)

    # 되먹임 라벨 — 신규 측점이 모형 입력이 된다
    seg(s, rx + rw / 2, 3.42 + rh + 0.10, rx + rw / 2, 5.35, color=ORANGE,
        weight=1.2)
    tf = textbox(s, rx - 0.6, 5.44, rw + 1.2, 0.5, align=PP_ALIGN.CENTER)
    para(tf, "선점된 신규 측점 = 모형의 Regional 층 입력", 9.5, ORANGE,
         bold=True, line=1.25, first=True)
    tf2 = textbox(s, rx - 0.6, 5.72, rw + 1.2, 0.3, align=PP_ALIGN.CENTER)
    para(tf2, "두 축은 함께 결정되어야 한다", 9, GRAY, line=1.25, first=True)
    return s


def s_part(prs, page, num, kind, color, title, question, items):
    s = blank(prs)

    if kind == "oval":
        geo(s, "oval", -2.9, 0.9, 6.3, 6.3, fill=color)
        nx, ny = 0.55, 2.75
    elif kind == "rect":
        geo(s, "rect", -1.6, -1.9, 5.6, 5.6, fill=color, rot=-15)
        nx, ny = 0.55, 0.95
    else:
        geo(s, "tri", -2.3, 2.3, 6.8, 5.9, fill=color)
        nx, ny = 0.85, 4.45

    tf = textbox(s, nx, ny, 3.4, 1.9)
    para(tf, num, 110, WHITE, fonts=F_DISP, bold=True, line=1.0, first=True)

    tx = 6.35
    tf = textbox(s, tx, 2.30, SW - tx - 0.7, 0.7)
    para(tf, title, 33, INK, fonts=F_DISP, bold=True, line=1.1, first=True)
    tf = textbox(s, tx, 3.06, SW - tx - 0.7, 0.34)
    para(tf, question, 13, GRAY, line=1.3, first=True)

    for i, it in enumerate(items):
        y = 3.78 + i * 0.60
        geo(s, "rect", tx, y + 0.05, 0.12, 0.12, fill=color)
        tf = textbox(s, tx + 0.26, y, SW - tx - 1.0, 0.3)
        para(tf, it, 12, INK, line=1.2, first=True)

    # 부 표지 — 흘림 도형이 왼쪽 아래를 덮으므로 쪽번호만 둔다
    tf = textbox(s, SW - M - 1.2, FOOT_Y + 0.09, 1.2, 0.24,
                 align=PP_ALIGN.RIGHT)
    para(tf, f"{page:02d} / {N_PAGES}", 8, GRAY, fonts=F_MONO, line=1.1,
         first=True)
    return s


# ══════════════════════════════════════════════════════════════════════
# 제1부 — 입지 선정
# ══════════════════════════════════════════════════════════════════════
def s_funnel(prs, ctx):
    """깔때기 — 좁히고(63) 다시 펼쳐(314) 확정한다(103)."""
    s = blank(prs)
    chrome(s, 5, "PART 01 · SELECTION PIPELINE")

    n_all = sum(len(ctx["cands"][t]) for t in (1, 2, 3))
    n_p1 = len(ctx["cands"][1])

    title_block(s, "1", [("약 1,800점이 ", INK), ("103개소", COBALT),
                         ("로 좁혀진다", INK)])

    stages = [
        ("전국 격자", "10 KM GRID", 1800, "약 1,800", "점", NAVY, "solid"),
        ("제외구역 필터", "11 ZONES", n_all, f"{n_all}", "개소", COBALT,
         "solid"),
        ("1등급 선별", "TOP TERCILE", n_p1, f"{n_p1}", "개소", COBALT,
         "solid"),
        ("서브격자 재확장", "1 KM × R5 KM", 314, "314", "개", COBALT,
         "outline"),
        ("도상선점 확정", "FIELD CARDS", 103, "103", "개소", NAVY, "solid"),
    ]
    x0 = M + 2.35
    wmax = 6.2
    bh, gap = 0.56, 0.31
    y = 1.72
    ys = []
    for name, mono, n, disp, unit, c, mode in stages:
        w = max(wmax * math.sqrt(n / 1800.0), 0.6)
        ys.append(y)
        if mode == "outline":
            geo(s, "rect", x0, y, w, bh, fill=WHITE, line=COBALT, weight=1.6)
        else:
            geo(s, "rect", x0, y, w, bh, fill=c)
        tf = textbox(s, M, y + 0.02, 2.2, 0.26, align=PP_ALIGN.RIGHT)
        para(tf, name, 11, INK, bold=True, line=1.15, first=True)
        tf = textbox(s, M, y + 0.29, 2.2, 0.22, align=PP_ALIGN.RIGHT)
        para(tf, mono, 8, GRAY, fonts=F_MONO, line=1.1, first=True,
             spacing=60)
        tf = textbox(s, x0 + w + 0.14, y - 0.02, 2.4, 0.5)
        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        run_in(p, disp, 25, c if mode != "outline" else COBALT,
               fonts=F_DISP, bold=True)
        run_in(p, " " + unit, 10.5, GRAY, fonts=F_MONO)
        y += bh + gap

    # 되먹임 — 확정 단계에서 서브격자로
    fx = x0 + wmax + 2.15
    y_sub = ys[3] + bh / 2
    y_fin = ys[4] + bh / 2
    seg(s, fx, y_fin, fx, y_sub, color=ORANGE, weight=1.4)
    seg(s, x0 + max(wmax * math.sqrt(103 / 1800.0), 0.6) + 2.0, y_fin, fx,
        y_fin, color=ORANGE, weight=1.4)
    arrow(s, fx, y_sub, x0 + wmax * math.sqrt(314 / 1800.0) + 1.9, y_sub,
          color=ORANGE, weight=1.4)
    tf = textbox(s, fx - 2.9, y_fin + 0.16, 4.2, 0.3)
    para(tf, "현장 부적합 → 차순위 재선정", 9.5, ORANGE, bold=True,
         line=1.2, first=True)

    tf = textbox(s, M, 6.28, CW, 0.3)
    p = tf.paragraphs[0]
    p.line_spacing = 1.25
    run_in(p, "NOTE  ", 8.5, COBALT, fonts=F_MONO, bold=True, spacing=100)
    run_in(p, "네 번째 단계는 축소가 아니라 재확장이다 — 1등급 63개소마다 "
              "반경 5 km 를 다시 훑어 설치 지점을 고른다.", 10, GRAY)
    return s


def s_exclusion(prs, ctx):
    s = blank(prs)
    chrome(s, 6, "PART 01 · EXCLUSION ZONES")

    title_block(s, "2", [("이격은 전부 ", INK), ("법정 최소 이상", COBALT),
                         ("이다", INK)])

    s.shapes.add_picture(str(ctx["figs"]["buffers"]), Inches(M),
                         Inches(1.72), width=Inches(7.0))
    cap(s, M, 5.36, 7.0, "FIG 01",
        "인공·지질 간섭원 11종의 적용 반경. 오렌지 눈금이 법정 최소값이다.")

    cx = M + 7.55
    cw = CW - 7.55
    stat(s, cx, 1.80, cw, "11", "종", "제외 항목",
         "인공 8 · 지질 2 · 자기이상 1", color=NAVY, size=33)
    stat(s, cx, 3.42, cw, "5.0", "km", "최대 이격 — 철도",
         "직류·교류 구분 불가 → 최대값", color=COBALT, size=33)

    geo(s, "rect", cx, 5.06, cw, 1.30, fill=FAINT)
    geo(s, "rect", cx, 5.06, 0.5, 0.12, fill=COBALT)
    tf = textbox(s, cx + 0.2, 5.30, cw - 0.4, 0.95)
    para(tf, "「지구물리측량 작업규정」 §13②", 9.5, INK, bold=True,
         line=1.3, first=True)
    para(tf, "모든 항목에서 법정 최소 이격을 충족하거나 초과", 9.5, GRAY,
         line=1.35, space_before=4)
    return s


def s_scoring(prs, ctx):
    s = blank(prs)
    chrome(s, 7, "PART 01 · SCORING")

    title_block(s, "3", [("배점의 4분의 1이 ", INK), ("공백 해소", COBALT),
                         ("다", INK)])

    # 배점 누적 막대 — 네이티브 사각형
    segs = [("① 희소성", 25, NAVY), ("② 지형", 15, COBALT),
            ("③ 전력철도", 15, COBALT), ("④ 인구", 15, COBALT),
            ("⑤ 모델기여", 10, COBALT), ("⑥ 암상", 5, COBALT)]
    alphas = {"② 지형": 0.85, "③ 전력철도": 0.7, "④ 인구": 0.55,
              "⑤ 모델기여": 0.4, "⑥ 암상": 0.3}

    bx, by, bw_total, bh = M, 1.92, CW, 0.62
    x = bx
    for name, val, c in segs:
        w = bw_total * val / 100.0
        # 단색 계열 — 코발트를 흰색과 섞어 단계화 (그라디언트 아님, 개별 단색)
        if name in alphas:
            a = alphas[name]
            mix = RGBColor(int(0x1D + (0xFF - 0x1D) * (1 - a)),
                           int(0x4E + (0xFF - 0x4E) * (1 - a)),
                           int(0xD8 + (0xFF - 0xD8) * (1 - a)))
            fill = mix
            tcol = WHITE if a >= 0.55 else INK
        else:
            fill = c
            tcol = WHITE
        sh = geo(s, "rect", x, by, w, bh, fill=fill, line=WHITE, weight=1.2)
        tfq = sh.text_frame
        tfq.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tfq.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run_in(p, str(val), 15, tcol, fonts=F_DISP, bold=True)
        tf = textbox(s, x, by - 0.28, max(w, 0.9), 0.24)
        para(tf, name, 8.5, INK, line=1.1, first=True)
        x += w
    for name, val in (("⑦ 지속성", 10), ("⑧ 접근성", 5)):
        w = bw_total * val / 100.0
        geo(s, "rect", x, by, w, bh, fill=WHITE, line=LINE, weight=1.0)
        for k in range(int(w / 0.14)):
            xx = x + 0.06 + k * 0.14
            seg(s, xx, by + bh - 0.06, min(xx + 0.10, x + w - 0.03),
                by + 0.06, color=LINE, weight=0.9)
        tf = textbox(s, x, by - 0.28, max(w, 0.9), 0.24)
        para(tf, name, 8.5, GRAY, line=1.1, first=True)
        x += w

    seg(s, bx, by + bh + 0.16, bx + bw_total * 0.85, by + bh + 0.16,
        color=NAVY, weight=1.4)
    tf = textbox(s, bx, by + bh + 0.24, bw_total * 0.85, 0.24)
    para(tf, "가용 85점 → 100점 정규화", 9, NAVY, bold=True, line=1.15,
         first=True)
    tf = textbox(s, bx + bw_total * 0.85, by + bh + 0.24, bw_total * 0.15,
                 0.24, align=PP_ALIGN.RIGHT)
    para(tf, "미산정 15점", 9, GRAY, line=1.15, first=True)

    # 분포 + 등급 규칙
    s.shapes.add_picture(str(ctx["figs"]["scores"]), Inches(M),
                         Inches(3.68), width=Inches(5.4))
    p33, p66 = ctx["pcts"]
    cap(s, M, 5.80, 5.4, "FIG 02",
        f"종합 점수 분포 — 3분위 경계 P33 {p33:.1f} · P66 {p66:.1f}. "
        "경계는 분포에 따라 자동 조정된다.")

    cx = M + 6.05
    cw = CW - 6.05
    notes = [("희소성 25점", "측정 공백 해소를 최우선으로 둔 설계", NAVY),
             ("모델기여 10점", "150–400 nT 구간 최고점 — 자기이상 반영",
              COBALT),
             ("미산정 15점", "⑦ 부지 지속성 · ⑧ 접근성 — 자료 미확보",
              ORANGE)]
    for i, (t, d, c) in enumerate(notes):
        y = 3.80 + i * 0.86
        geo(s, "rect", cx, y + 0.03, 0.13, 0.13, fill=c)
        tf = textbox(s, cx + 0.28, y, cw - 0.28, 0.26)
        para(tf, t, 11, INK, bold=True, line=1.2, first=True)
        tf = textbox(s, cx + 0.28, y + 0.28, cw - 0.28, 0.30)
        para(tf, d, 9.5, GRAY, line=1.3, first=True)
    return s


def s_results(prs, ctx):
    s = blank(prs)
    chrome(s, 8, "PART 01 · CANDIDATES")

    title_block(s, "4", [("후보는 ", INK), ("해안과 산간", COBALT),
                         ("에 몰린다", INK)])

    cands = ctx["cands"]
    fx = M + 7.0
    s.shapes.add_picture(str(ctx["figs"]["cmap"]), Inches(fx), Inches(1.45),
                         height=Inches(5.15))
    cap(s, fx - 0.2, 6.56, CW - 6.9, "FIG 03", "등급별 전국 분포.")

    for i, (tier, label) in enumerate(((1, "1등급 · P1"), (2, "2등급"),
                                       (3, "3등급"))):
        sc = [p["score"] for p in cands[tier]]
        stat(s, M, 1.72 + i * 1.30, 2.9, f"{len(cands[tier])}", "개소", label,
             f"{min(sc):.1f}–{max(sc):.1f} PTS",
             color=COBALT if tier == 1 else GRAY, size=30)

    geo(s, "rect", M + 3.35, 1.80, 3.15, 3.55, fill=FAINT)
    geo(s, "tri", M + 5.55, 1.55, 0.75, 0.65, fill=ORANGE)
    tf = textbox(s, M + 3.60, 2.10, 2.65, 3.1)
    para(tf, "판독", 9.5, COBALT, fonts=F_MONO, bold=True, line=1.2,
         first=True, spacing=100)
    para(tf, "제외구역이 조밀한 내륙 평야부에는 후보가 남지 않는다.",
         11.5, INK, bold=True, line=1.4, space_before=10)
    para(tf, "관측망 대표성에는 부합 — 접근성과 부지 확보 난이도는 상승.",
         10, GRAY, line=1.4, space_before=8)

    tf = textbox(s, M, 5.80, 6.3, 0.5)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    run_in(p, "Q  ", 9.5, ORANGE, fonts=F_MONO, bold=True)
    run_in(p, "이 쏠림을 받아들일 것인가, 배점을 손볼 것인가 — 질의 1-2",
           10, GRAY)
    return s


def s_subgrid(prs, ctx):
    s = blank(prs)
    chrome(s, 9, "PART 01 · SUB-GRID")

    title_block(s, "5", [("5 km 원 안에서 ", INK), ("다시 고른다", COBALT)])

    fx = M + 7.35
    s.shapes.add_picture(str(ctx["figs"]["subgrid"]), Inches(fx),
                         Inches(1.75), height=Inches(4.35))
    cap(s, fx - 0.25, 6.24, CW - 7.0, "FIG 04", "선정 개념도 — 실제 좌표가 아니다.")

    rows = [("서브격자 후보", "사이트당 상위 5개", "314", "개"),
            ("도로 접근", "고속도로 제외", "≤2", "km"),
            ("국가기준점", "733점 → 제외구역 필터", "496", "점"),
            ("임도망도", "3,322 노선 · 산림청", "9,018", "km")]
    for i, (label, sub, val, unit) in enumerate(rows):
        y = 1.86 + i * 1.12
        seg(s, M, y, M + 6.95, y, color=LINE, weight=0.75)
        geo(s, "rect", M, y + 0.15, 0.12, 0.12, fill=COBALT)
        tf = textbox(s, M + 0.26, y + 0.10, 3.2, 0.28)
        para(tf, label, 12, INK, bold=True, line=1.2, first=True)
        tf = textbox(s, M + 0.26, y + 0.44, 3.6, 0.26)
        para(tf, sub, 9, GRAY, line=1.2, first=True)
        tf = textbox(s, M + 3.9, y + 0.06, 3.05, 0.52, align=PP_ALIGN.RIGHT)
        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        run_in(p, val, 27, NAVY, fonts=F_DISP, bold=True)
        run_in(p, " " + unit, 11, GRAY, fonts=F_MONO)
    seg(s, M, 1.86 + 4 * 1.12, M + 6.95, 1.86 + 4 * 1.12, color=LINE,
        weight=0.75)
    return s


def s_fieldcard(prs, ctx):
    s = blank(prs)
    chrome(s, 10, "PART 01 · FIELD SURVEY")

    title_block(s, "6", [("판정은 조사자가 아니라 ", INK),
                         ("규칙이 한다", COBALT)])

    items = [("차량", "80 m", NAVY), ("전력", "60 m", NAVY),
             ("통신", "50 m", NAVY), ("금속", "30 m", NAVY),
             ("건축", "70 m", NAVY), ("매설물", "존재=부적합", ORANGE),
             ("상공 장애", "존재=부적합", ORANGE),
             ("방위표지", "불가=부적합", ORANGE)]
    cw0, gap = 1.55, 0.15
    for i, (name, crit, c) in enumerate(items):
        x = M + (i % 4) * (cw0 + gap)
        y = 1.80 + (i // 4) * 1.14
        geo(s, "rect", x, y, cw0, 0.98, fill=WHITE, line=LINE, weight=1.0)
        geo(s, "rect", x, y, cw0, 0.10, fill=c)
        tf = textbox(s, x + 0.14, y + 0.22, cw0 - 0.28, 0.28)
        para(tf, name, 11.5, INK, bold=True, line=1.15, first=True)
        tf = textbox(s, x + 0.14, y + 0.56, cw0 - 0.28, 0.28)
        para(tf, crit, 9, GRAY, fonts=F_MONO, line=1.15, first=True)

    tf = textbox(s, M, 4.18, 6.7, 0.26)
    para(tf, "8개 항목 — 존재 여부만 드롭다운으로 고르면 판정 자동", 9.5,
         GRAY, line=1.2, first=True)

    # 판정 산식 — 도형 수식
    jy = 4.72
    rules = [("0 건", "적합", "자기구배 조사 진행", NAVY),
             ("1–2 건", "조건부", "추가 조사 후 재판정", COBALT),
             ("3 건+", "부적합", "차순위 후보 선정", ORANGE)]
    for i, (cond, verdict, action, c) in enumerate(rules):
        x = M + i * 2.32
        geo(s, "oval", x, jy, 0.62, 0.62, fill=c)
        tf = textbox(s, x, jy + 0.14, 0.62, 0.3, align=PP_ALIGN.CENTER)
        para(tf, cond.split(" ")[0], 11, WHITE, fonts=F_DISP, bold=True,
             line=1.0, first=True)
        tf = textbox(s, x + 0.76, jy - 0.02, 1.5, 0.3)
        para(tf, verdict, 12.5, INK, bold=True, line=1.15, first=True)
        tf = textbox(s, x + 0.76, jy + 0.28, 1.55, 0.42)
        para(tf, action, 8.5, GRAY, line=1.2, first=True)

    tf = textbox(s, M, 5.80, 6.7, 0.7)
    p = tf.paragraphs[0]
    p.line_spacing = 1.35
    run_in(p, "CARD  ", 8.5, COBALT, fonts=F_MONO, bold=True, spacing=100)
    run_in(p, "기본정보(관할본부 자동 배정) → 자기교란 → 방위표지(방위각 "
              "자동 계산) → 종합 판정 → 사진·약도. 103건 본부별 분철.",
           9.5, GRAY)

    # 우측 — 기준 이격의 출처 질문
    qx = M + 7.35
    qw = CW - 7.35
    geo(s, "rect", qx, 1.80, qw, 4.55, fill=FAINT)
    geo(s, "tri", qx + qw - 1.0, 1.52, 0.78, 0.68, fill=ORANGE)
    tf = textbox(s, qx + 0.28, 2.14, qw - 0.56, 3.9)
    para(tf, "OPEN QUESTION", 9, ORANGE, fonts=F_MONO, bold=True, line=1.2,
         first=True, spacing=120)
    para(tf, "기준 이격 80·60·50·30·70 m 는 실측 근거가 아직 없다.",
         12.5, INK, bold=True, line=1.42, space_before=12)
    para(tf, "실측 검증이 필요한가, 판정 구간(1–2건 조건부)은 적절한가 — "
             "질의 1-5", 10, GRAY, line=1.42, space_before=10)
    return s


# ══════════════════════════════════════════════════════════════════════
# 제2부 — LMM
# ══════════════════════════════════════════════════════════════════════
def s_layers(prs, ctx):
    """4층 레이어 아키텍처 — 전부 네이티브 도형."""
    s = blank(prs)
    chrome(s, 12, "PART 02 · MODEL ARCHITECTURE")

    title_block(s, "7", [("자기장은 ", INK), ("네 층의 합", COBALT),
                         ("이다", INK)])

    # 식 — 모노
    tf = textbox(s, M, 1.62, CW, 0.3)
    p = tf.paragraphs[0]
    p.line_spacing = 1.15
    run_in(p, "B(r,t) = B", 14, INK, fonts=F_MONO, bold=True)
    run_in(p, "IGRF", 9, NAVY, fonts=F_MONO, bold=True)
    run_in(p, " + B", 14, INK, fonts=F_MONO, bold=True)
    run_in(p, "REGIONAL", 9, COBALT, fonts=F_MONO, bold=True)
    run_in(p, " + B", 14, INK, fonts=F_MONO, bold=True)
    run_in(p, "CRUSTAL", 9, COBALT, fonts=F_MONO, bold=True)
    run_in(p, " + B", 14, INK, fonts=F_MONO, bold=True)
    run_in(p, "EXTERNAL", 9, ORANGE, fonts=F_MONO, bold=True)

    bands = [
        ("①", "CORE", "IGRF-14 · 13차 구면조화", "≥ 3,000 KM", "확보",
         NAVY, WHITE, 1, 0.05),
        ("②", "REGIONAL", "지표 절대측정 17점 · 1차 다항식",
         "50 – 3,000 KM", "확보", COBALT, WHITE, 4, 0.11),
        ("③", "CRUSTAL", "KIGAM 항공자력 1.5′ 격자", "0.05 – 50 KM",
         "해상도 부족", FAINT, INK, 18, 0.07),
        ("④", "EXTERNAL", "청양 CYG 1분 자료", "시간 변동 (1분–수시간)",
         "부분 적용", WHITE, INK, 0, 0.0),
    ]
    bx, bw = M, 7.6
    by, bh, gap = 2.08, 0.94, 0.16
    for i, (n, name, src, band, state, fill, tcol, nzig, amp) in \
            enumerate(bands):
        y = by + i * (bh + gap)
        outline = (name == "EXTERNAL")
        geo(s, "rect", bx, y, bw, bh, fill=fill,
            line=LINE if outline else None, weight=1.2)
        tf = textbox(s, bx + 0.2, y + 0.14, 0.5, 0.3)
        para(tf, n, 13, tcol, fonts=F_MONO, bold=True, line=1.0, first=True)
        tf = textbox(s, bx + 0.62, y + 0.12, 2.4, 0.3)
        para(tf, name, 14.5, tcol, fonts=F_DISP, bold=True, line=1.1,
             first=True, spacing=60)
        tf = textbox(s, bx + 0.62, y + 0.50, 3.4, 0.3)
        para(tf, src, 9.5, tcol if fill in (NAVY, COBALT) else GRAY,
             line=1.2, first=True)
        # 파장 지그재그
        zx, zw = bx + 3.55, 1.55
        zc = WHITE if fill in (NAVY, COBALT) else NAVY
        if nzig == 1:
            seg(s, zx, y + bh * 0.62, zx + zw, y + bh * 0.38, color=zc,
                weight=1.6)
        elif nzig > 1:
            zigzag(s, zx, y + bh / 2, zw, nzig, amp, color=zc, weight=1.4)
        else:
            seg(s, zx, y + bh / 2, zx + zw, y + bh / 2, color=ORANGE,
                weight=1.6, dash="dash")
        # 파장·상태 라벨 — 띠 오른쪽 끝 안쪽에서 끝나도록 (밖으로 나가면 잘려 보인다)
        tf = textbox(s, bx + bw - 2.45, y + 0.14, 2.25, 0.26,
                     align=PP_ALIGN.RIGHT)
        para(tf, band, 8, tcol if fill in (NAVY, COBALT) else GRAY,
             fonts=F_MONO, line=1.1, first=True)
        ok = state == "확보"
        tf = textbox(s, bx + bw - 2.45, y + 0.50, 2.25, 0.28,
                     align=PP_ALIGN.RIGHT)
        para(tf, state, 10,
             (WHITE if fill in (NAVY, COBALT) else (INK if ok else ORANGE)),
             bold=True, line=1.1, first=True)

    # 우측 — 설계상 한계
    cx = M + 8.15
    cw2 = CW - 8.15
    geo(s, "rect", cx, 2.08, cw2, 2.45, fill=WHITE, line=INK, weight=1.2)
    geo(s, "rect", cx, 2.08, 0.5, 0.14, fill=ORANGE)
    tf = textbox(s, cx + 0.24, 2.36, cw2 - 0.48, 2.0)
    para(tf, "설계상 한계", 9.5, ORANGE, fonts=F_MONO, bold=True, line=1.2,
         first=True, spacing=80)
    para(tf, "지각층은 총자력 F 에만 더해진다.", 11.5, INK, bold=True,
         line=1.4, space_before=10)
    para(tf, "항공자력이 스칼라 이상만 제공 — 편각 D·복각 I 에는 미반영. "
             "이 선택이 §9 의 첫 검정 대상.", 9.5, GRAY, line=1.42,
         space_before=8)

    tf = textbox(s, cx, 4.85, cw2, 0.9)
    para(tf, "층마다 담당 파장이 달라", 10.5, INK, line=1.35, first=True)
    para(tf, "단순 내삽으로는 분리되지 않는다", 10.5, INK, line=1.35)
    return s


def s_pipeline(prs, ctx):
    """데이터 파이프라인 — 입력 → 적합 → 검증 → 산출."""
    s = blank(prs)
    chrome(s, 13, "PART 02 · DATA PIPELINE")

    title_block(s, "8", [("입력에서 계산기까지 ", INK), ("한 줄", COBALT),
                         ("로 흐른다", INK)])

    # 입력 카드 4
    inputs = [("성과표 2022–25", "15측점 · 연도만 기록", ORANGE),
              ("야장 2019", "검증 통과 6건 · 분 단위 시각", NAVY),
              ("KIGAM 항공자력", "1.5′ 격자 · 1982–2018", NAVY),
              ("CYG 1분 자료", "2019 구간 78일 확보", NAVY)]
    ix, iw, ih, igap = M, 2.55, 0.80, 0.12
    sy = 2.42
    sw2, sh2 = 2.05, 1.55
    for i, (name, sub, c) in enumerate(inputs):
        y = 1.70 + i * (ih + igap)
        geo(s, "rect", ix, y, iw, ih, fill=WHITE, line=LINE, weight=1.0)
        geo(s, "rect", ix, y, 0.14, ih, fill=c)
        tf = textbox(s, ix + 0.28, y + 0.10, iw - 0.42, 0.28)
        para(tf, name, 10.5, INK, bold=True, line=1.15, first=True)
        tf = textbox(s, ix + 0.28, y + 0.42, iw - 0.42, 0.3)
        para(tf, sub, 8.5, GRAY, line=1.2, first=True)
        arrow(s, ix + iw + 0.05, y + ih / 2, ix + iw + 0.68,
              sy + sh2 / 2, color=LINE, weight=1.1)

    # 적합 → 검증 → 산출
    stages = [("적합", "lmm_build.py", "층 결합 · 잔차면", NAVY),
              ("검증", "LOO 교차검증", "17측점 하나씩 제외", COBALT)]
    sx = ix + iw + 0.75
    for i, (name, mono, sub, c) in enumerate(stages):
        x = sx + i * (sw2 + 0.55)
        geo(s, "rect", x, sy, sw2, sh2, fill=c)
        tf = textbox(s, x + 0.2, sy + 0.22, sw2 - 0.4, 0.34)
        para(tf, name, 16, WHITE, fonts=F_DISP, bold=True, line=1.1,
             first=True)
        tf = textbox(s, x + 0.2, sy + 0.62, sw2 - 0.4, 0.26)
        para(tf, mono, 8.5, WHITE, fonts=F_MONO, line=1.15, first=True)
        tf = textbox(s, x + 0.2, sy + 0.94, sw2 - 0.4, 0.4)
        para(tf, sub, 9, WHITE, line=1.25, first=True)
        if i == 0:
            arrow(s, x + sw2 + 0.06, sy + sh2 / 2, x + sw2 + 0.49,
                  sy + sh2 / 2, color=INK, weight=1.5)

    # 산출물 3
    outs = [("웹 계산기", "lmm.html · 오프라인 단일 파일"),
            ("엑셀 계산기", "수식 공개 · 매크로 없음"),
            ("자침편차 (보류)", "KPI 미달 — 정식 산출 금지")]
    ox = sx + 2 * (sw2 + 0.55)
    ow = SW - M - ox
    for i, (name, sub) in enumerate(outs):
        y = 1.70 + i * 0.98
        hold = "보류" in name
        geo(s, "rect", ox, y, ow, 0.84, fill=FAINT if not hold else WHITE,
            line=ORANGE if hold else None, weight=1.4)
        tf = textbox(s, ox + 0.2, y + 0.11, ow - 0.4, 0.28)
        para(tf, name, 10.5, ORANGE if hold else INK, bold=True, line=1.15,
             first=True)
        tf = textbox(s, ox + 0.2, y + 0.43, ow - 0.4, 0.3)
        para(tf, sub, 8.5, GRAY, line=1.2, first=True)
    arrow(s, ox - 0.49, sy + sh2 / 2, ox - 0.06, sy + sh2 / 2, color=INK,
          weight=1.5)

    # 일치 검증 스트립
    gy = 5.45
    geo(s, "rect", M, gy, CW, 0.78, fill=NAVY)
    tf = textbox(s, M + 0.26, gy + 0.13, 3.4, 0.3)
    para(tf, "CROSS-CHECK", 9.5, ORANGE, fonts=F_MONO, bold=True, line=1.15,
         first=True, spacing=120)
    tf = textbox(s, M + 0.26, gy + 0.40, 3.4, 0.28)
    para(tf, "세 구현의 수치 일치", 10.5, WHITE, bold=True, line=1.15,
         first=True)
    checks = [("JS ↔ PY", "0.22″ · 0.27 nT"),
              ("XLSX ↔ PY", "0.33″ · 0.12 nT"),
              ("JS ↔ ppigrf", "0.17″ · 0.03 nT")]
    for i, (pair, tol) in enumerate(checks):
        x = M + 4.0 + i * 2.55
        tf = textbox(s, x, gy + 0.12, 2.4, 0.26)
        para(tf, pair, 9, ORANGE, fonts=F_MONO, bold=True, line=1.1,
             first=True)
        tf = textbox(s, x, gy + 0.40, 2.4, 0.26)
        para(tf, tol, 9.5, WHITE, fonts=F_MONO, line=1.1, first=True)
    return s


def s_validation(prs, ctx):
    s = blank(prs)
    chrome(s, 14, "PART 02 · VALIDATION")

    model = ctx["model"]
    cv = model["loo_cv"]
    cd = model["crustal_diagnostics"]

    title_block(s, "9", [("교차검증은 아직 ", INK), ("KPI 밖", ORANGE),
                         ("이다", INK)])

    stat(s, M, 1.74, 2.6, f"{cv['D']:.3f}", "°", "편각 D — LOO",
         "KPI 0.1° · 7.7×", color=ORANGE, size=34)
    stat(s, M + 2.9, 1.74, 2.6, f"{cv['F']:.1f}", "nT", "총자력 F — LOO",
         "KPI 50 nT · 1.2×", color=ORANGE, size=34)
    stat(s, M + 5.8, 1.74, 2.6, f"{cv['I']:.3f}", "°", "복각 I — 참고",
         "기준 없음", color=NAVY, size=34)

    fx = M + 8.75
    s.shapes.add_picture(str(ctx["figs"]["sites"]), Inches(fx), Inches(1.55),
                         height=Inches(2.35))
    cap(s, fx - 0.1, 3.96, CW - 8.6, "FIG 05", "측점 17 · 관측 33회.")

    geo(s, "rect", M, 3.60, 8.4, 0.72, fill=FAINT)
    tf = textbox(s, M + 0.24, 3.72, 8.0, 0.5)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    run_in(p, "CRUSTAL  ", 8.5, COBALT, fonts=F_MONO, bold=True, spacing=100)
    run_in(p, f"지각층이 총자력 잔차 {cd['rms_before_nT']:.0f}→"
              f"{cd['rms_after_nT']:.0f} nT ({cd['rms_reduction_pct']:.0f}% "
              f"감소). 단 표본 {cd['n']}개 — 측점 하나에 좌우, 단정 금지.",
           9.5, INK)

    s.shapes.add_picture(str(ctx["figs"]["resid"]), Inches(M), Inches(4.52),
                         width=Inches(CW))
    cap(s, M, 6.52, CW, "FIG 06",
        "측점별 편각 잔차 — KPI 대(帶) 안은 14개 중 3개뿐이다.")
    return s


def s_diagnosis(prs, ctx):
    s = blank(prs)
    chrome(s, 15, "PART 02 · RESIDUAL FORENSICS")

    diag = ctx["diag"]
    hyp = diag["hypotheses"]
    asym = diag["asymmetry"]
    vr = diag["vector_recovery"]

    # 서사를 자료에 맞춘다 — 감사·지각 벡터 반영 후 가설 A 는 유의해졌다.
    a_sig = hyp["A_vector_D"]["p"] < 0.05
    if a_sig:
        title_block(s, "10", [("가설 넷 중 ", INK), ("하나", ORANGE),
                              ("만 살아남았다", INK)])
    else:
        title_block(s, "10", [("남는 것은 측점에 고정된 ", INK), ("30′", ORANGE),
                              (" 다", INK)])

    # 가설 검정 매트릭스 (기각된 것에만 사선)
    tests = [("A", "지각 벡터 누락", hyp["A_vector_D"]),
             ("B", "측정 정밀도", hyp["B_precision"]),
             ("B′", "관측 횟수", hyp["B_nvisit"]),
             ("C", "국소 거칠기", hyp["C_roughness"])]
    bw0, gap = 1.78, 0.20
    for i, (tag, name, h) in enumerate(tests):
        x = M + i * (bw0 + gap)
        geo(s, "rect", x, 1.76, bw0, 1.30, fill=WHITE, line=LINE, weight=1.0)
        tf = textbox(s, x + 0.14, 1.90, bw0 - 0.28, 0.24)
        para(tf, f"H-{tag}", 8.5, GRAY, fonts=F_MONO, bold=True, line=1.1,
             first=True, spacing=80)
        tf = textbox(s, x + 0.14, 2.14, bw0 - 0.28, 0.30)
        para(tf, name, 10.5, INK, bold=True, line=1.18, first=True)
        tf = textbox(s, x + 0.14, 2.62, bw0 - 0.28, 0.28)
        p = tf.paragraphs[0]
        p.line_spacing = 1.1
        run_in(p, f"r {h['r']:+.2f} ", 9.5, COBALT, fonts=F_MONO, bold=True)
        run_in(p, f"p {h['p']:.2f}", 8.5, GRAY, fonts=F_MONO)
        if h["p"] >= 0.05:
            seg(s, x + 0.14, 2.30, x + bw0 - 0.14, 2.56, color=ORANGE,
                weight=1.8)   # 기각 사선
    tf = textbox(s, M + 4 * (bw0 + gap) + 0.1, 1.94, 3.0, 1.0)
    n_rej = sum(1 for _, _, h in tests if h["p"] >= 0.05)
    para(tf, f"{n_rej}건 기각" if a_sig else "전부 기각", 17, GRAY,
         fonts=F_DISP, bold=True, line=1.1, first=True)
    para(tf, ("A 는 유의 — 벡터를 복원해 모형에 넣었다"
              if a_sig else "n≈14 · 유의하려면 |r|>0.53"),
         8.5, GRAY, fonts=F_MONO, line=1.3, space_before=4)

    s.shapes.add_picture(str(ctx["figs"]["vario"]), Inches(M), Inches(3.40),
                         width=Inches(4.9))
    cap(s, M, 5.78, 4.9, "FIG 07",
        "잔차 반변량도 — 순수 nugget. 공간 상관이 없다.")

    # D/I 비대칭 — 도형 저울
    dx = M + 5.45
    geo(s, "rect", dx, 3.46, 0.95, 0.95, fill=ORANGE)
    tf = textbox(s, dx, 3.62, 0.95, 0.5, align=PP_ALIGN.CENTER)
    para(tf, "D", 21, WHITE, fonts=F_DISP, bold=True, line=1.0, first=True)
    geo(s, "rect", dx + 1.15, 3.72, 0.55, 0.55, fill=NAVY)
    tf = textbox(s, dx + 1.15, 3.80, 0.55, 0.36, align=PP_ALIGN.CENTER)
    para(tf, "I", 14, WHITE, fonts=F_DISP, bold=True, line=1.0, first=True)
    tf = textbox(s, dx, 4.62, 2.0, 0.6)
    p = tf.paragraphs[0]
    p.line_spacing = 1.25
    run_in(p, f"D/I = {asym['ratio']:.2f}", 11, INK, fonts=F_MONO, bold=True)
    tf = textbox(s, dx, 4.92, 2.2, 0.5)
    para(tf, f"자기 원인이면 {asym['expected_from_magnetic']:.2f} — "
             "편각에만 작용하는 오차", 8.5, GRAY, line=1.3, first=True)

    # 결론 박스
    cx2 = M + 8.0
    cw2 = CW - 8.0
    geo(s, "rect", cx2, 3.40, cw2, 2.65, fill=WHITE, line=INK, weight=1.2)
    geo(s, "rect", cx2, 3.40, 0.5, 0.14, fill=ORANGE)
    tf = textbox(s, cx2 + 0.22, 3.66, cw2 - 0.44, 2.3)
    para(tf, "남은 후보 3", 9.5, ORANGE, fonts=F_MONO, bold=True, line=1.2,
         first=True, spacing=80)
    for k, t in enumerate(("마크 진북 방위각 오차", "성과표 전사·환산 오류",
                           "편각 환산 절차 오류")):
        para(tf, f"{k + 1}  {t}", 10.5, INK, bold=True, line=1.5,
             space_before=6)
    para(tf, "자료만으로는 구분 불가 — 원본 감사 필요", 9, GRAY, line=1.35,
         space_before=8)

    tf = textbox(s, M, 6.30, 7.0, 0.4)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    run_in(p, "EVIDENCE  ", 8.5, COBALT, fonts=F_MONO, bold=True, spacing=80)
    run_in(p, f"벡터 복원 예측 {vr['predD_rms_arcmin']:.1f}′ ≪ 잔차 "
              f"{asym['rD_rms']:.1f}′ · 성산 2회 재현 0.6′ 인데 둘 다 34′ "
              "오차 — 계통 오차의 서명.", 9.5, GRAY)
    return s


def s_external(prs, ctx):
    s = blank(prs)
    chrome(s, 16, "PART 02 · EXTERNAL LAYER")

    title_block(s, "11", [("관측시각을 찾았고, 관측소를 ", INK),
                          ("넷", ORANGE), ("으로 늘렸다", INK)])

    # 사슬 도해 — 2026-08 에 전 구간이 이어졌다
    chain = [("야장 원본", "212세션 · 분 단위", False),
             ("관측 시각", "2019~2025 복원", False),
             ("관측소 4소", "청양·제주·강릉·이천", False),
             ("④ EXTERNAL", "편각에 적용", False)]
    bx, bw0, bh0, gap = M, 2.5, 1.05, 0.42
    for i, (name, sub, broken) in enumerate(chain):
        x = bx + i * (bw0 + gap)
        geo(s, "rect", x, 1.80, bw0, bh0, fill=WHITE,
            line=ORANGE if broken else INK, weight=1.4)
        tf = textbox(s, x + 0.18, 1.96, bw0 - 0.36, 0.3,
                     align=PP_ALIGN.CENTER)
        para(tf, name, 12, ORANGE if broken else INK, bold=True, line=1.15,
             first=True)
        tf = textbox(s, x + 0.18, 2.32, bw0 - 0.36, 0.28,
                     align=PP_ALIGN.CENTER)
        para(tf, sub, 9, GRAY, line=1.2, first=True)
        if i < len(chain) - 1:
            mid_x = x + bw0 + gap / 2
            arrow(s, x + bw0 + 0.04, 1.80 + bh0 / 2, x + bw0 + gap - 0.04,
                  1.80 + bh0 / 2, color=ORANGE, weight=1.4)
            if False:    # (2026-08) 사슬이 이어져 절단 표시를 쓰지 않는다
                seg(s, mid_x - 0.08, 1.74, mid_x + 0.02, 2.20,
                    color=ORANGE, weight=2.2)
                seg(s, mid_x + 0.02, 1.74, mid_x + 0.12, 2.20,
                    color=ORANGE, weight=2.2)

    stat(s, M, 3.30, 3.0, "0.112", "°", "편각 일변동 중앙값",
         "전형적 하루가 이미 KPI 0.1° 초과", color=ORANGE, size=33)
    stat(s, M + 3.35, 3.30, 3.0, "37.5", "nT", "총자력 일변동 중앙값",
         "최대 68.4 nT — KPI 50 nT", color=ORANGE, size=33)

    # 방침
    px = M + 6.85
    pw = CW - 6.85
    geo(s, "rect", px, 3.30, pw, 2.55, fill=FAINT)
    geo(s, "rect", px, 3.30, 0.5, 0.12, fill=NAVY)
    tf = textbox(s, px + 0.24, 3.56, pw - 0.48, 2.2)
    para(tf, "적용 방침", 9.5, NAVY, fonts=F_MONO, bold=True, line=1.2,
         first=True, spacing=80)
    para(tf, "관측소 4소 공간보간 · 정온 세션 · 편각에만 적용.", 11, INK,
         bold=True, line=1.4, space_before=8)
    para(tf, "청양 단독 균일장 근사는 편각을 악화시켰다. 4소 1차 평면으로 "
             "바꾸고 교란(Kp>2) 세션을 빼자 LOO 편각이 0.590°→0.323° 로 "
             "줄었다. 총자력에는 적용하지 않는다 — F 잔차는 지각장 불일치가 "
             "지배해 잡음만 는다.", 9.5, GRAY, line=1.4, space_before=7)

    geo(s, "rect", M, 6.10, 6.5, 0.52, fill=WHITE, line=COBALT, weight=1.2)
    tf = textbox(s, M + 0.2, 6.21, 6.1, 0.3)
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    run_in(p, "NEW METRIC  ", 8.5, COBALT, fonts=F_MONO, bold=True,
           spacing=80)
    run_in(p, "겹치는 3측점 시기간 차 → 0 수렴이면 외부장, 잔존이면 성과표 "
              "오류.", 9.5, INK)
    return s


def s_legal(prs, ctx):
    s = blank(prs)
    chrome(s, 17, "PART 02 · LEGAL BASIS & SCALE")

    title_block(s, "12", [("1:25,000 이 ", INK), ("단일값 표기", COBALT),
                          ("가 성립하는 한계다", INK)])

    rows = [("철도 이격 · 직류", "5.0", "5.0", "부합"),
            ("철도 이격 · 교류/일반", "2.0", "5.0", "보수적"),
            ("고압철탑 이격", "1.0", "1.0", "부합"),
            ("송전탑 이격", "0.5", "1.0", "보수적"),
            ("편각 정수차 §20 (′)", "30", "20", "엄격")]
    tw = 6.0
    tf = textbox(s, M, 1.76, tw, 0.24)
    p = tf.paragraphs[0]
    run_in(p, "LAW vs APPLIED (KM)", 8.5, GRAY, fonts=F_MONO, bold=True,
           spacing=100)
    hy = 2.08
    for i, (item, legal, ours, verdict) in enumerate(rows):
        y = hy + i * 0.62
        seg(s, M, y, M + tw, y, color=LINE, weight=0.75)
        tf = textbox(s, M, y + 0.13, tw * 0.44, 0.3)
        para(tf, item, 10.5, INK, line=1.2, first=True)
        tf = textbox(s, M + tw * 0.46, y + 0.11, tw * 0.16, 0.3,
                     align=PP_ALIGN.RIGHT)
        para(tf, legal, 11, GRAY, fonts=F_MONO, line=1.2, first=True)
        tf = textbox(s, M + tw * 0.64, y + 0.11, tw * 0.16, 0.3,
                     align=PP_ALIGN.RIGHT)
        para(tf, ours, 11, COBALT, fonts=F_MONO, bold=True, line=1.2,
             first=True)
        tf = textbox(s, M + tw * 0.82, y + 0.13, tw * 0.18, 0.3,
                     align=PP_ALIGN.RIGHT)
        para(tf, verdict, 9.5, ORANGE, line=1.2, first=True)
    seg(s, M, hy + 5 * 0.62, M + tw, hy + 5 * 0.62, color=LINE, weight=0.75)

    tf = textbox(s, M, 5.48, tw, 0.9)
    p = tf.paragraphs[0]
    p.line_spacing = 1.35
    run_in(p, "KPI  ", 8.5, ORANGE, fonts=F_MONO, bold=True, spacing=100)
    run_in(p, "D<0.1° 는 공학적 목표치 — 법정 정수차 30′ 와 성격이 다르다. "
              "KPI 미달은 법령 위반이 아니다.", 9.5, GRAY)

    fx = M + 6.55
    fw = CW - 6.55
    s.shapes.add_picture(str(ctx["figs"]["scale"]), Inches(fx), Inches(1.80),
                         width=Inches(fw))
    cap(s, fx, 4.28, fw, "FIG 08",
        "축척별 도폭 내 편각 변화 — IGRF-14 직접 계산.")

    sc = {r["name"]: r["dD"] for r in ctx["scales"]}
    geo(s, "rect", fx, 4.78, fw, 1.55, fill=FAINT)
    geo(s, "rect", fx, 4.78, 0.5, 0.12, fill=COBALT)
    tf = textbox(s, fx + 0.24, 5.02, fw - 0.48, 1.2)
    para(tf, f"1:50,000 은 도폭 내 변화만으로 예산의 "
             f"{100 * sc['1:50,000'] / 0.1:.0f}% 소진.", 10.5, INK,
         bold=True, line=1.4, first=True)
    para(tf, f"1:25,000 은 {100 * sc['1:25,000'] / 0.1:.0f}% — 도폭당 "
             "단일 편각값이 성립하는 가장 성긴 축척.", 10, GRAY, line=1.4,
         space_before=6)
    return s


def s_roadmap(prs, ctx):
    s = blank(prs)
    chrome(s, 18, "PART 02 · ROADMAP")

    title_block(s, "13", [("야장 확보가 ", INK), ("모든 것의 전제", COBALT),
                          ("다", INK)])

    tf = textbox(s, M, 1.60, CW, 0.3)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    run_in(p, "S/N  ", 8.5, ORANGE, fonts=F_MONO, bold=True, spacing=100)
    run_in(p, "제거 가능한 잡음 6.7′ 에 미설명 잔차 35′ — 순서를 뒤집으면 "
              "개선 여부를 판정할 수 없다.", 10.5, GRAY)

    steps = [("0", "2022–25 야장 확보", "관측 일시 복원 · 공통 전제", NAVY),
             ("1", "성과표 원본 감사", "편각 환산·전사 오류 확인", COBALT),
             ("2", "External 포함 재적합", "층 결합 — 순차 튜닝 금지",
              COBALT),
             ("3", "교차검증 · 독립 표본", "IGRF 대조는 Core 확인용",
              COBALT)]
    bw0, gap = 2.55, 0.42
    by = 2.20
    for i, (n, name, sub, c) in enumerate(steps):
        x = M + i * (bw0 + gap)
        first = (i == 0)
        geo(s, "rect", x, by, bw0, 1.55, fill=c if first else WHITE,
            line=None if first else LINE, weight=1.0)
        tf = textbox(s, x + 0.2, by + 0.14, bw0 - 0.4, 0.5)
        para(tf, n, 26, WHITE if first else COBALT, fonts=F_DISP, bold=True,
             line=1.0, first=True)
        tf = textbox(s, x + 0.2, by + 0.72, bw0 - 0.4, 0.3)
        para(tf, name, 11.5, WHITE if first else INK, bold=True, line=1.2,
             first=True)
        tf = textbox(s, x + 0.2, by + 1.06, bw0 - 0.4, 0.36)
        para(tf, sub, 8.5, WHITE if first else GRAY, line=1.25, first=True)
        if i < len(steps) - 1:
            arrow(s, x + bw0 + 0.05, by + 0.78, x + bw0 + gap - 0.05,
                  by + 0.78, color=INK, weight=1.4)

    for k, (tag, text) in enumerate((
            ("병렬 A", "NOC 공간투영 구현 — 2019 세션을 회귀 가드로"),
            ("병렬 B", "3측점 시기간 대조를 검증 지표로 고정 (즉시 가능)"))):
        y = 4.10 + k * 0.34
        geo(s, "rect", M, y + 0.04, 0.11, 0.11, fill=COBALT)
        tf = textbox(s, M + 0.26, y, 6.6, 0.3)
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        run_in(p, tag + "  ", 9.5, COBALT, bold=True)
        run_in(p, text, 9.5, INK)

    # 1단계 판정 게이트 — 삼거리
    gy = 5.02
    tf = textbox(s, M, gy, 4.0, 0.26)
    para(tf, "GATE — 감사 후 편각 잔차", 9, GRAY, fonts=F_MONO, bold=True,
         line=1.15, first=True, spacing=80)
    gates = [("≤ 10′", "외부장 보정 진행", NAVY),
             ("10–20′", "잔차 큰 측점 재측정", COBALT),
             ("> 20′", "중단 · 방위 기준 재결정", ORANGE)]
    for i, (cond, act, c) in enumerate(gates):
        x = M + i * 2.55
        geo(s, "rect", x, gy + 0.34, 0.72, 0.44, fill=c)
        tf = textbox(s, x, gy + 0.42, 0.72, 0.3, align=PP_ALIGN.CENTER)
        para(tf, cond, 9, WHITE, fonts=F_MONO, bold=True, line=1.05,
             first=True)
        tf = textbox(s, x + 0.84, gy + 0.42, 1.75, 0.42)
        para(tf, act, 9, INK, line=1.2, first=True)

    # 현재 위치
    cx = M + 8.3
    cw2 = CW - 8.3
    geo(s, "rect", cx, 4.10, cw2, 2.15, fill=WHITE, line=ORANGE, weight=1.6)
    geo(s, "tri", cx + cw2 - 0.95, 3.82, 0.72, 0.62, fill=ORANGE)
    tf = textbox(s, cx + 0.24, 4.32, cw2 - 0.48, 0.26)
    para(tf, "NOW", 9, ORANGE, fonts=F_MONO, bold=True, line=1.1,
         first=True, spacing=140)
    tf = textbox(s, cx + 0.24, 4.58, cw2 - 0.48, 0.75)
    p = tf.paragraphs[0]
    p.line_spacing = 1.0
    run_in(p, "35", 40, ORANGE, fonts=F_DISP, bold=True)
    run_in(p, " ′", 16, ORANGE, fonts=F_MONO)
    tf = textbox(s, cx + 0.24, 5.44, cw2 - 0.48, 0.7)
    para(tf, "세 번째 구간에서 시작 — 판정기준은 착수 전 확정 완료.", 9.5,
         GRAY, line=1.35, first=True)
    return s


# ══════════════════════════════════════════════════════════════════════
# 제3부 — 자문 질의 (의사결정 매트릭스)
# ══════════════════════════════════════════════════════════════════════
def s_qmatrix(prs, page, num, title_runs, rows):
    s = blank(prs)
    chrome(s, page, "PART 03 · ADVISORY MATRIX")

    title_block(s, num, title_runs)

    c1, c2, c3 = 2.35, 4.15, 4.55
    x1 = M + 0.42
    x2 = x1 + c1 + 0.20
    x3 = x2 + c2 + 0.20

    tf = textbox(s, x1, 1.66, c1, 0.24)
    para(tf, "쟁점", 9, GRAY, fonts=F_MONO, bold=True, line=1.1, first=True,
         spacing=100)
    tf = textbox(s, x2, 1.66, c2, 0.24)
    para(tf, "현재 상태", 9, GRAY, fonts=F_MONO, bold=True, line=1.1,
         first=True, spacing=100)
    tf = textbox(s, x3, 1.66, c3, 0.24)
    para(tf, "자문 질문", 9, GRAY, fonts=F_MONO, bold=True, line=1.1,
         first=True, spacing=100)
    seg(s, M, 1.96, SW - M, 1.96, color=INK, weight=1.2)

    rh = 0.80
    for i, (issue, state, q) in enumerate(rows):
        y = 1.96 + i * rh
        if i:
            seg(s, M, y, SW - M, y, color=LINE, weight=0.75)
        geo(s, "rect", M, y + 0.16, 0.26, 0.26,
            fill=COBALT if num == "14" else NAVY)
        tf = textbox(s, M, y + 0.20, 0.26, 0.2, align=PP_ALIGN.CENTER)
        para(tf, str(i + 1), 9, WHITE, fonts=F_MONO, bold=True, line=1.0,
             first=True)
        tf = textbox(s, x1, y + 0.13, c1, rh - 0.16)
        para(tf, issue, 10.5, INK, bold=True, line=1.22, first=True)
        tf = textbox(s, x2, y + 0.13, c2, rh - 0.16)
        para(tf, state, 9.5, GRAY, line=1.25, first=True)
        tf = textbox(s, x3, y + 0.13, c3, rh - 0.16)
        para(tf, q, 9.5, INK, line=1.25, first=True)
    seg(s, M, 1.96 + len(rows) * rh, SW - M, 1.96 + len(rows) * rh,
        color=INK, weight=1.2)
    return s


def s_closing(prs, ctx):
    s = blank(prs)

    geo(s, "oval", 10.7, -1.9, 4.4, 4.4, fill=FAINT)
    geo(s, "rect", -0.7, 5.95, 6.2, 0.9, fill=NAVY)
    geo(s, "tri", 11.9, 5.1, 1.15, 1.0, fill=ORANGE, rot=12)

    tf = textbox(s, M, 0.36, 8.0, 0.26)
    para(tf, "CLOSING", 9.5, GRAY, fonts=F_MONO, line=1.1, first=True,
         spacing=140)

    geo(s, "rect", M, 1.30, 0.95, 0.10, fill=ORANGE)
    tf = textbox(s, M, 1.58, 10.0, 0.7)
    para(tf, "다음 단계는 정해져 있다", 34, INK, fonts=F_DISP, bold=True,
         line=1.1, first=True)

    points = [
        ("01", "두 축은 하나의 문제", "선점 기준이 요구 정밀도를 정하고, "
         "확보된 측점이 다시 모형의 입력이 된다.", NAVY),
        ("02", "막힌 곳은 자료", "구현 오류 0건 — 관측시각·원측선·측점 수 "
         "세 가지가 성능을 묶고 있다.", COBALT),
        ("03", "판정기준을 먼저", "결과를 보고 기준을 정하면 자기확증이 "
         "된다. 1단계 기준은 착수 전에 확정하였다.", ORANGE),
    ]
    colw = (CW - 2 * 0.55) / 3
    for i, (n, t, d, c) in enumerate(points):
        x = M + i * (colw + 0.55)
        geo(s, "rect", x, 2.78, 0.15, 0.15, fill=c)
        tf = textbox(s, x + 0.27, 2.72, colw - 0.27, 0.26)
        para(tf, n, 10, c, fonts=F_MONO, bold=True, line=1.15, first=True)
        tf = textbox(s, x, 3.10, colw, 0.34)
        para(tf, t, 15, INK, fonts=F_DISP, bold=True, line=1.2, first=True)
        tf = textbox(s, x, 3.54, colw, 1.2)
        para(tf, d, 10.5, GRAY, line=1.5, first=True)

    tf = textbox(s, M + 0.0, 6.14, 5.2, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "NEXT — FIELDBOOKS 2022-25 → AUDIT → REFIT", 10.5, WHITE,
         fonts=F_MONO, bold=True, line=1.15, first=True, spacing=60)

    tf = textbox(s, SW - M - 4.4, 6.30, 4.4, 0.3, align=PP_ALIGN.RIGHT)
    para(tf, DECK_TITLE, 8.5, GRAY, fonts=F_MONO, line=1.15, first=True,
         spacing=80)
    return s


# ══════════════════════════════════════════════════════════════════════
def main():
    cands = load_candidates()
    model = load_model()
    diag = load_diag()
    scales = scale_rows()
    figs, pcts = build_figures(cands, model, diag, scales)

    ctx = {"cands": cands, "model": model, "diag": diag,
           "scales": scales, "figs": figs, "pcts": pcts}

    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    s_cover(prs)                                                  # 01
    s_dashboard(prs, ctx)                                         # 02
    s_causal(prs, ctx)                                            # 03
    s_part(prs, 4, "01", "oval", NAVY, "측정 입지 선정",
           "어디에서 측정할 것인가",
           ["전국 격자 → 제외구역 11종 필터",
            "다기준 점수 · 3등급 분류",
            "서브격자 · 국가기준점 매칭",
            "현장 자기교란 조사카드"])                             # 04
    s_funnel(prs, ctx)                                            # 05
    s_exclusion(prs, ctx)                                         # 06
    s_scoring(prs, ctx)                                           # 07
    s_results(prs, ctx)                                           # 08
    s_subgrid(prs, ctx)                                           # 09
    s_fieldcard(prs, ctx)                                         # 10
    s_part(prs, 11, "02", "rect", COBALT, "지역 자기장 모형",
           "측정값을 어떻게 면(面)으로 만들 것인가",
           ["4층 결합 — IGRF · Regional · Crustal · External",
            "Leave-One-Out 교차검증 · KPI 판정",
            "편각 잔차 원인 진단",
            "웹 · 엑셀 계산기 · 도엽별 자침편차"])                 # 11
    s_layers(prs, ctx)                                            # 12
    s_pipeline(prs, ctx)                                          # 13
    s_validation(prs, ctx)                                        # 14
    s_diagnosis(prs, ctx)                                         # 15
    s_external(prs, ctx)                                          # 16
    s_legal(prs, ctx)                                             # 17
    s_roadmap(prs, ctx)                                           # 18
    s_part(prs, 19, "03", "tri", ORANGE, "자문 질의사항",
           "판단을 구하고자 하는 12건",
           ["입지 선정 6건 — 기준의 타당성",
            "지역 자기장 모형 6건 — 다음 단계의 선택"])            # 19

    s_qmatrix(prs, 20, "14",
              [("질의 ① ", INK), ("입지 선정", COBALT), (" 6건", INK)], [
        ("이격거리 보수적 적용", "철도 5.0 km 전면 적용 — 내륙 평야부 후보 소멸",
         "직류·교류 구분 자료를 확보해 차등 적용해야 하는가"),
        ("배점 가중", "희소성 25점 vs 모델기여 10점",
         "신규 측점이 곧 모형 입력인데 모델기여 가중을 올릴 것인가"),
        ("자기이상 구간 점수", "150–400 nT 최고점 · >800 nT 조건부 제외",
         "구간 설정의 물리적 근거가 타당한가"),
        ("미산정 15점", "⑦ 부지 지속성 · ⑧ 접근성 — 자료 미확보",
         "국공유지 DB·도로망으로 대체 산정인가, 현장 판단인가"),
        ("현장 조사 기준", "기준 이격 80·60·50·30·70 m — 실측 근거 없음",
         "실측 검증이 필요한가, 판정 구간(1–2건 조건부)은 적절한가"),
        ("국가기준점 활용", "삼각점·통합기준점 496점을 위치 기준으로 제공",
         "표석의 철근·자성 재료가 총자력 측정에 영향을 주는가"),
    ])                                                            # 20

    s_qmatrix(prs, 21, "15",
              [("질의 ② ", INK), ("지역 자기장 모형", COBALT), (" 6건", INK)],
              [
        ("계통 오차 감사 순서", "잔차 30′ 후보 3 — 마크 방위각·전사·환산",
         "무엇부터 감사하는가, 자료만으로 가리는 진단법이 있는가"),
        ("지각층 각도 반영", "현행은 F 에만 반영 · 벡터 역산 규모는 8′",
         "주원인이 아닌데 벡터 역산을 구현할 가치가 있는가"),
        ("외부장 공간투영", "상시관측소 청양 1개소뿐",
         "NOC 공간투영으로 전역 복원에 몇 개소가 필요한가"),
        ("일변화 이중보정", "§21 은 법정 절차 — 성과표 적용 여부 미확인",
         "확인 방법과 이중보정 방지 안전장치를 어떻게 두는가"),
        ("Regional 함수 형태", "17측점 · 1차 다항식",
         "30점 도달 시 2차인가, 구면캡조화(SCHA)·크리깅인가"),
        ("KPI 와 2027 캠페인", "KPI 0.1° 는 일변동 중앙값 0.112° 보다 작음",
         "시각 보정 후 달성 가능한가, 재측정과 신규점 비율은"),
    ])                                                            # 21

    s_closing(prs, ctx)                                           # 22

    DOCS_OUT.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = DOCS_OUT / f"{stamp}_지자기_자문회의_발표자료.pptx"
    prs.save(out)
    print("saved:", out)
    print("slides:", len(prs.slides._sldIdLst))
    return out


if __name__ == "__main__":
    main()
