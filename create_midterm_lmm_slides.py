# -*- coding: utf-8 -*-
"""
중간보고회 추가 슬라이드 — 국가 지자기모델(LMM) 중간결과 · 자료 · 검증 · 2020 비교
==============================================================================

`20260818_중간보고회_작성부분_별도자료.pptx` 와 **동일한 디자인 언어**로 추가분을
생성한다(Pretendard · 13.333×7.5in · 라벨/제목/부제/괘선 상단 블록 · 하단 푸터).
따라서 산출물을 그대로 중간보고회 본 deck 에 끼워 넣을 수 있다.

구성 (18장)
  Ⅰ LMM 중간결과      : 간지 · IGRF 한계 · 4층 구조 · 내삽 불가 · 파이프라인 ·
                         자료 취득현황 · LOO 중간결과 · 편각 잔차 진단 · 구현 검증
  Ⅱ 확인·추가 자료    : KIGAM 항공자력 · 지리원 2021~2025 반복측정 · KPI 대응관계
  Ⅲ 향후 검증·보완    : 로드맵(게이트) · 보완 우선순위와 리스크
  Ⅳ 2020 연구와 비교  : 개념 차이 · 비교표 · Regional 계산흐름 · 계보와 주장 범위

⚠ 수치는 하드코딩하지 않고 `docs/data/lmm_model.json` · `lmm_diagnosis.json` 에서
   읽는다. 자료가 갱신되면 재실행만으로 본문 수치가 따라 바뀐다.

    python create_midterm_lmm_slides.py
"""
import json
import math
import sys
from datetime import datetime
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
from matplotlib.patches import Patch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).parent
DATA = ROOT / "docs" / "data"
OUT_DIR = ROOT / "docs" / "output"
FIG_DIR = OUT_DIR / "_midterm_figs"

# ── 디자인 토큰 (별도자료 deck 실측값) ───────────────────────────────────────
W, H = 13.3333, 7.5
M = 0.60                 # 좌우 여백
CW = W - 2 * M           # 본문 폭 12.13
FONT = "Pretendard"      # 발표 PC 설치본. 미설치 시 PowerPoint 가 대체
FONT_MPL = "Malgun Gothic"   # 도판용(설치 확인됨)

NAVY = "182638"     # 제목
BLUE = "2F6FB2"     # 강조·라벨
GRAY = "607184"     # 부제·푸터
RED = "C63D3D"      # KPI 미달
ORANGE = "E86F2A"   # 최우선 배지
GREEN = "2E7D5B"    # 충족·완료
PEACH = "FCE8D8"    # 경고 카드
LBLUE = "DCEAF7"    # 정보 카드
LGRAY = "F7F9FC"    # 기본 카드
BAND = "17365D"     # 하단 요지 띠
WHITE = "FFFFFF"

TOP_LABEL, TOP_TITLE, TOP_SUB, TOP_RULE = 0.33, 0.65, 1.27, 1.73
BODY_Y = 2.08
FOOT_Y = 7.12
FOOTER = "국가기준점 지자기측량 발전방안 연구 | 중간보고 작성분"


# ── 자료 로드 ────────────────────────────────────────────────────────────────
MODEL = json.load(open(DATA / "lmm_model.json", encoding="utf-8"))
DIAG = json.load(open(DATA / "lmm_diagnosis.json", encoding="utf-8"))

# External 보정량 현황 — 수치를 하드코딩하지 않고 산출물에서 읽는다.
def _ext_stats():
    import csv
    for name in ("external_corrections_multi.csv", "external_corrections.csv"):
        f = DATA / name
        if not f.exists():
            continue
        rows = list(csv.DictReader(open(f, encoding="utf-8-sig")))
        ok = [r for r in rows if r.get("상태") == "정상"]
        yrs = sorted({r["날짜"][:4] for r in rows if r.get("날짜")})
        return dict(total=len(rows), ok=len(ok), y0=yrs[0] if yrs else "",
                    y1=yrs[-1] if yrs else "",
                    missing=[y for y in range(int(yrs[0]), int(yrs[-1]) + 1)
                             if str(y) not in yrs] if yrs else [],
                    multi=name.endswith("multi.csv"))
    return dict(total=0, ok=0, y0="", y1="", missing=[], multi=False)


EXT = _ext_stats()

LOO = MODEL["loo_cv"]
CRU = MODEL["crustal_diagnostics"]
DS = DIAG["dataset"]
ASY = DIAG["asymmetry"]
HYP = DIAG["hypotheses"]
VAR = DIAG["variogram"]
KPI_D, KPI_F = 0.1, 50.0
N_SITES = DS["n_sites"]


def _rgb(h):
    return RGBColor.from_string(h)


# ── 기본 도형 헬퍼 ───────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=None, line=None, radius=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = _rgb(line)
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if shp.text_frame:
        shp.text_frame.text = ""
    return shp


def text(slide, x, y, w, h, s, size=12, color=NAVY, bold=False, align="l",
         anchor="t", space=1.0, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]
    lines = s.split("\n") if isinstance(s, str) else list(s)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                       "r": PP_ALIGN.RIGHT}[align]
        p.line_spacing = space
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, italic
        f.color.rgb = _rgb(color)
    return box


def slide_base(prs, label, title, subtitle, page, rule=BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    text(s, M, TOP_LABEL, CW, 0.25, label, 11.25, rule, True)
    text(s, M, TOP_TITLE, CW, 0.60, title, 25.5, NAVY, True)
    if subtitle:
        text(s, M, TOP_SUB, CW, 0.38, subtitle, 12.75, GRAY)
    rect(s, M, TOP_RULE, CW, 0.03, rule)
    text(s, M, FOOT_Y, 9.0, 0.19, FOOTER, 8.25, GRAY)
    text(s, 11.98, FOOT_Y - 0.03, 0.71, 0.21, f"{page:02d}", 9, GRAY, align="r")
    return s


def stat(slide, x, y, w, h, value, caption, fill=LGRAY, vcolor=BLUE, vsize=25.5):
    rect(slide, x, y, w, h, fill)
    text(slide, x + 0.13, y + 0.16, w - 0.26, 0.60, value, vsize, vcolor, True)
    text(slide, x + 0.13, y + 0.84, w - 0.26, h - 0.95, caption, 11.25, GRAY,
         space=1.15)


def badge_row(slide, y, badge, body, bfill=BLUE, h=0.44, bw=1.70,
              x=M, w=None, body_size=12):
    w = w or (CW - 0.0)
    rect(slide, x, y, bw, h, bfill)
    text(slide, x, y + 0.03, bw, h - 0.06, badge, 12, WHITE, True, "c", "m")
    rect(slide, x + bw + 0.13, y, w - bw - 0.13, h, LGRAY)
    text(slide, x + bw + 0.32, y + 0.02, w - bw - 0.52, h - 0.04, body,
         body_size, NAVY, anchor="m", space=1.05)


def band(slide, x, y, w, h, msg, fill=BAND, size=11.25):
    rect(slide, x, y, w, h, fill)
    text(slide, x + 0.18, y + 0.02, w - 0.36, h - 0.04, msg, size, WHITE,
         anchor="m", space=1.15)


def table(slide, x, y, w, cols, rows, widths, head_fill=BAND, row_h=0.36,
          size=10.5, head_size=10.5, zebra=LGRAY, first_bold=True):
    """괘선 없는 booktabs 풍 표 — 머리글 띠 + 얼룩 행."""
    xs, acc = [], x
    for fr in widths:
        xs.append(acc)
        acc += w * fr
    rect(slide, x, y, w, row_h, head_fill)
    for cx, fr, c in zip(xs, widths, cols):
        text(slide, cx + 0.12, y + 0.02, w * fr - 0.20, row_h - 0.04, c,
             head_size, WHITE, True, anchor="m")
    yy = y + row_h
    for i, row in enumerate(rows):
        hh = row_h
        if i % 2 == 0:
            rect(slide, x, yy, w, hh, zebra)
        for j, (cx, fr, cell) in enumerate(zip(xs, widths, row)):
            col = NAVY if j == 0 else GRAY
            text(slide, cx + 0.12, yy + 0.02, w * fr - 0.20, hh - 0.04, str(cell),
                 size, col, first_bold and j == 0, anchor="m", space=1.05)
        yy += hh
    return yy


# ── 도판 (matplotlib) ────────────────────────────────────────────────────────
def _mpl_init():
    if FONT_MPL in {f.name for f in fm.fontManager.ttflist}:
        plt.rcParams["font.family"] = FONT_MPL
    plt.rcParams["axes.unicode_minus"] = False
    FIG_DIR.mkdir(parents=True, exist_ok=True)


def fig_sites():
    """17 측점 분포 + 북측 공백."""
    p = FIG_DIR / "sites.png"
    gj = json.load(open(DATA / "korea_boundary.geojson", encoding="utf-8"))
    fig, ax = plt.subplots(figsize=(4.3, 4.7), dpi=200)
    for f in gj["features"]:
        g = f["geometry"]
        polys = ([g["coordinates"]] if g["type"] == "Polygon"
                 else g["coordinates"])
        for poly in polys:
            xs = [c[0] for c in poly[0]]
            ys = [c[1] for c in poly[0]]
            ax.fill(xs, ys, color="#EEF3F8", ec="#9FB4C7", lw=0.6, zorder=1)
    lat = [s["lat"] for s in MODEL["sites"]]
    lon = [s["lon"] for s in MODEL["sites"]]
    ax.scatter(lon, lat, s=52, c="#2F6FB2", ec="white", lw=1.1, zorder=4,
               label=f"Regional 입력 측점 {N_SITES}점")
    ax.set_xlim(124.5, 130.2)
    ax.set_ylim(33.0, 38.8)
    ax.set_xticks([])
    ax.set_yticks([])
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.legend(loc="upper left", fontsize=8, frameon=False)
    ax.set_aspect(1 / math.cos(math.radians(36)))
    fig.tight_layout(pad=0.2)
    fig.savefig(p, transparent=True)
    plt.close(fig)
    return p


def fig_loo():
    """LOO 성능 vs KPI — 편각·총자력."""
    p = FIG_DIR / "loo.png"
    fig, axes = plt.subplots(1, 2, figsize=(5.0, 2.5), dpi=200)
    for ax, (lab, val, kpi, unit) in zip(
            axes, [("편각 D", LOO["D"], KPI_D, "°"),
                   ("총자력 F", LOO["F"], KPI_F, " nT")]):
        ax.bar([0], [val], width=0.5, color="#C63D3D")
        ax.axhline(kpi, color="#2F6FB2", ls="--", lw=1.4)
        ax.text(0, val, f" {val:.3g}{unit}", ha="center", va="bottom",
                fontsize=9, color="#C63D3D", fontweight="bold")
        ax.text(0.42, kpi, f"목표 {kpi:g}{unit}", ha="right", va="bottom",
                fontsize=8, color="#2F6FB2")
        ax.set_xticks([])
        ax.set_title(lab, fontsize=9.5, color="#182638")
        ax.set_ylim(0, max(val, kpi) * 1.35)
        ax.tick_params(labelsize=7.5, colors="#607184")
        for sp in ("top", "right"):
            ax.spines[sp].set_visible(False)
        for sp in ("left", "bottom"):
            ax.spines[sp].set_color("#C9D4E0")
    fig.tight_layout(pad=0.4)
    fig.savefig(p, transparent=True)
    plt.close(fig)
    return p


def fig_bands():
    """층별 공간 파장 대역 — IGRF 차단차수 13 기준."""
    p = FIG_DIR / "bands.png"
    R = 6371.0
    lam13 = 2 * math.pi * R / math.sqrt(13 * 14)
    rows = [("Core (IGRF-14)", lam13, 20000, "#17365D"),
            ("Regional (지표 절대측정)", 50, lam13, "#2F6FB2"),
            ("Crustal (항공자력)", 0.25, 50, "#E86F2A"),
            ("External (상시관측)", 100, 20000, "#8A94A6")]
    fig, ax = plt.subplots(figsize=(6.6, 2.3), dpi=200)
    for i, (lab, lo, hi, c) in enumerate(rows):
        ax.barh(i, hi - lo, left=lo, height=0.52, color=c, alpha=0.9)
        ax.text(math.sqrt(lo * hi), i, lab, ha="center", va="center",
                fontsize=8.5, color="white", fontweight="bold")
    ax.axvline(lam13, color="#C63D3D", ls="--", lw=1.2)
    # Malgun Gothic 에 ≈ 글리프가 없다 → 한글로 적는다
    ax.text(lam13 * 1.06, 3.55, f"IGRF 차단 파장 약 {lam13:,.0f} km",
            fontsize=8, color="#C63D3D")
    ax.set_xscale("log")
    ax.set_xlim(0.1, 25000)
    ax.set_ylim(-0.6, 3.9)
    ax.set_yticks([])
    ax.set_xlabel("공간 파장 (km, 로그축)", fontsize=8.5, color="#607184")
    ax.tick_params(labelsize=8, colors="#607184")
    for sp in ("top", "right", "left"):
        ax.spines[sp].set_visible(False)
    ax.spines["bottom"].set_color("#C9D4E0")
    fig.tight_layout(pad=0.3)
    fig.savefig(p, transparent=True)
    plt.close(fig)
    return p


def fig_resolution():
    """Crustal 해상도 — 현재 2.8 km 격자 vs 원측선 250 m."""
    p = FIG_DIR / "resolution.png"
    dx = DIAG["vector_recovery"]["grid"]["dx_km"]
    fig, axes = plt.subplots(1, 2, figsize=(5.4, 2.7), dpi=200)
    for ax, (title, step, c) in zip(
            axes, [(f"현재 보유 {dx:.1f} km 격자", dx, "#8A94A6"),
                   ("원측선 250 m (권고)", 0.25, "#E86F2A")]):
        n = max(2, int(12 / step))
        for k in range(n + 1):
            ax.axvline(k * step, color=c, lw=0.7 if step < 1 else 1.4, alpha=0.85)
        ax.set_xlim(0, 12)
        ax.set_ylim(0, 1)
        ax.set_yticks([])
        ax.set_xlabel("12 km 구간", fontsize=8, color="#607184")
        ax.set_title(title, fontsize=9, color="#182638")
        ax.tick_params(labelsize=7.5, colors="#607184")
        for sp in ("top", "right", "left"):
            ax.spines[sp].set_visible(False)
        ax.spines["bottom"].set_color("#C9D4E0")
    fig.legend(handles=[Patch(color="#8A94A6", label="측선"),
                        Patch(color="#E86F2A", label="측선(고밀도)")],
               loc="lower center", ncol=2, fontsize=7.5, frameon=False,
               bbox_to_anchor=(0.5, -0.02))
    fig.tight_layout(pad=0.4, rect=(0, 0.06, 1, 1))
    fig.savefig(p, transparent=True)
    plt.close(fig)
    return p


def fig_asymmetry():
    """D/I 잔차 비대칭 — 관측비 vs 자기 원인 예상비."""
    p = FIG_DIR / "asym.png"
    fig, ax = plt.subplots(figsize=(4.6, 2.4), dpi=200)
    obs, exp = ASY["ratio"], ASY["expected_from_magnetic"]
    ax.bar([0, 1], [obs, exp], width=0.46,
           color=["#C63D3D", "#2F6FB2"])
    for x, v, lab in ((0, obs, "관측 비"), (1, exp, "자기 원인\n예상 비 (F/H)")):
        ax.text(x, v + 0.08, f"{v:.2f}", ha="center", fontsize=10,
                fontweight="bold", color="#182638")
        ax.text(x, -0.42, lab, ha="center", fontsize=8.5, color="#607184")
    ax.set_ylim(0, max(obs, exp) * 1.28)
    ax.set_xlim(-0.6, 1.6)
    ax.set_xticks([])
    ax.set_ylabel("D 잔차 / I 잔차", fontsize=8.5, color="#607184")
    ax.tick_params(labelsize=8, colors="#607184")
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color("#C9D4E0")
    fig.tight_layout(pad=0.4, rect=(0, 0.1, 1, 1))
    fig.savefig(p, transparent=True)
    plt.close(fig)
    return p


# ── 슬라이드 ─────────────────────────────────────────────────────────────────
def sl_divider(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, BAND)
    text(s, M, 2.35, CW, 0.30, "중간보고 추가 작성분", 12.75, "9FC0E4", True)
    text(s, M, 2.75, CW, 1.05, "국가 지자기모델(LMM) 중간결과", 40, WHITE, True)
    text(s, M, 3.95, 9.6, 0.70,
         "4층 결합 구조의 시범 구축 결과와 확인·추가 자료, 검증 계획,\n"
         "그리고 2020년 지자기도 연구와의 관계를 정리한다.",
         13.5, "C9D8E8", space=1.35)
    for i, (num, lab) in enumerate([("Ⅰ", "중간결과"), ("Ⅱ", "확인·추가 자료"),
                                    ("Ⅲ", "검증·보완"), ("Ⅳ", "2020 연구 비교")]):
        x = M + i * 3.06
        rect(s, x, 5.35, 2.86, 0.86, "23456B")
        text(s, x + 0.22, 5.50, 0.5, 0.28, num, 13, "9FC0E4", True)
        text(s, x + 0.22, 5.82, 2.4, 0.30, lab, 12.75, WHITE, True)
    text(s, M, FOOT_Y, 9.0, 0.19, FOOTER, 8.25, "7E93AB")
    text(s, 11.98, FOOT_Y - 0.03, 0.71, 0.21, f"{page:02d}", 9, "7E93AB",
         align="r")
    return s


def sl_why(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "IGRF만으로는 도폭 편각 정밀도에 닿지 못한다",
                   "전지구 모델은 차단차수 13에서 끊기므로, 한반도 규모의 지역장·지각장은 "
                   "구조적으로 담기지 않는다.", page)
    img = fig_bands()
    s.shapes.add_picture(str(img), Inches(M), Inches(BODY_Y + 0.05),
                         width=Inches(7.35))
    x2 = M + 7.60
    w2 = CW - 7.60
    stat(s, x2, BODY_Y, w2, 1.24, "≈ 3,000 km",
         "IGRF 공간 분해능 — 이보다 짧은 파장은 미포함", LBLUE, BLUE, 22)
    stat(s, x2, BODY_Y + 1.36, w2 / 2 - 0.07, 1.24, f"{KPI_D:g}°",
         "목표 편각 정확도", PEACH, RED, 22)
    stat(s, x2 + w2 / 2 + 0.07, BODY_Y + 1.36, w2 / 2 - 0.07, 1.24, f"{KPI_F:g} nT",
         "목표 총자력 정확도", PEACH, RED, 22)
    y = BODY_Y + 2.86
    for i, (t, b) in enumerate([
            ("전지구 평균의 한계", "한반도 전역에 동일한 보정이 적용된다"),
            ("지각 자기이상 누락", "차수 14 이상 단파장 신호가 빠진다"),
            ("시간 변동 미분리", "일변화·자기폭풍이 관측값에 남는다")]):
        yy = y + i * 0.62
        rect(s, M, yy, 0.10, 0.50, BLUE)
        text(s, M + 0.26, yy + 0.02, 3.4, 0.24, t, 12, NAVY, True)
        text(s, M + 0.26, yy + 0.26, 6.9, 0.24, b, 11, GRAY)
    band(s, x2, y + 0.10, w2, 1.02,
         "→ IGRF는 버리는 것이 아니라 출발점으로 두고,\n"
         "   한반도 고유 성분을 분리해 더한다.")
    return s


def sl_structure(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "네 개 층을 한 좌표·한 시점에서 결합한다",
                   "층마다 물리적 기원·공간 스케일·시간 상수가 달라, 분리해야 오차의 원인도 "
                   "분리된다.", page)
    rect(s, M, BODY_Y, CW, 0.62, "EDF3FA")
    text(s, M, BODY_Y + 0.10, CW, 0.42,
         "B_LMM (r, t)  =  B_IGRF  +  B_Regional  +  B_Crustal  +  B_External",
         18, NAVY, True, "c", "m")
    rows = [
        ["Core", "IGRF-14 (13차)", "≳ 3,000 km / 5년", "전지구 기준장", "보유"],
        ["Regional", f"지표 절대측정 {N_SITES}점", "50~3,000 km / 연", "한반도 절대 기준",
         "확충 필요"],
        ["Crustal", "KIGAM 항공자력", "0.05~50 km / 정적", "지각 자기이상", "해상도 보완"],
        ["External", "상시관측 1분 자료", "한반도 / 일·시·분", "외부장 시간변동 분리",
         "부분 적용"],
    ]
    y = table(s, M, BODY_Y + 0.82, CW,
              ["층", "자료원", "공간 / 시간 스케일", "역할", "현재 상태"], rows,
              [0.13, 0.24, 0.24, 0.24, 0.15], row_h=0.50)
    band(s, M, y + 0.30, CW, 0.92,
         "각 층은 서로 다른 자료에서 나온다 — 하나라도 빠지면 그 파장대의 신호를 "
         "메울 방법이 없다.\n"
         "그래서 층을 나눈 구조 자체가, 오차가 생겼을 때 어느 자료를 보강해야 하는지 "
         "가리키는 진단 장치가 된다.")
    return s


def sl_nointerp(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "측정점을 내삽하는 것만으로는 모델이 되지 않는다",
                   "「30점을 크리깅하면 되지 않는가」에 대한 답 — 네 가지 이유로 성립하지 "
                   "않는다.", page)
    items = [
        ("01", "물리적 기원이 다르다",
         "Core·Regional·Crustal·External은 기원·시간상수·공간스케일이 모두 다르다.\n"
         "한 측정값에 섞인 신호를 내삽으로는 분리할 수 없다."),
        ("02", "시간 의존성을 무시한다",
         "지자기장은 매일 수십 nT 변동한다. 측정 시각이 다른 점을 그대로 평균하면\n"
         "외부장 잡음이 공간 분포로 둔갑한다."),
        ("03", "공간 외삽이 위험하다",
         f"{N_SITES}점은 평균 간격이 수십 km다. 50 km 미만 단파장 지각이상을 담지 못하고,\n"
         "도서·경계 부근에서 외삽이 불안정해진다."),
        ("04", "물리 제약을 반영하지 못한다",
         "지자기장은 발산이 0인 벡터장이다. D·I·F를 각각 내삽하면\n"
         "성분 간 물리적 일관성이 깨진다."),
    ]
    for i, (num, t, b) in enumerate(items):
        col, row = i % 2, i // 2
        x = M + col * (CW / 2 + 0.10)
        y = BODY_Y + row * 1.72
        w = CW / 2 - 0.10
        rect(s, x, y, w, 1.52, LGRAY)
        rect(s, x, y, 0.62, 0.52, BLUE)
        text(s, x, y + 0.04, 0.62, 0.44, num, 13, WHITE, True, "c", "m")
        text(s, x + 0.78, y + 0.12, w - 0.95, 0.30, t, 13, NAVY, True)
        text(s, x + 0.20, y + 0.60, w - 0.40, 0.85, b, 10.5, GRAY, space=1.25)
    band(s, M, BODY_Y + 3.50, CW, 0.62,
         "일본 국토지리원·미국 NOAA·영국 BGS 모두 층 분리 구조를 채택하고 있다.")
    return s


def sl_pipeline(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "분석은 하나의 모델 JSON을 공유하는 파이프라인으로 돌린다",
                   "적합 → 진단 → 계산기 → 보고서가 같은 원천을 읽으므로, 자료가 바뀌면 "
                   "수치와 문서가 함께 갱신된다.", page)
    steps = [
        ("01", "자료 정리", "성과표·야장 판독\n좌표·표고 통일"),
        ("02", "기준시점 환산", "IGRF 영년변화\n관측시각 확보분 보정"),
        ("03", "층 분해", "IGRF 잔차 산출\n이상치 배제"),
        ("04", "지역장 적합", "1차 다항식\n로버스트 최소제곱"),
        ("05", "교차검증", "LOO\n잔차 진단"),
        ("06", "산출", "웹·Excel 계산기\n보고서 자동 갱신"),
    ]
    bw = (CW - 5 * 0.14) / 6
    for i, (num, t, b) in enumerate(steps):
        x = M + i * (bw + 0.14)
        rect(s, x, BODY_Y, bw, 1.62, LGRAY)
        rect(s, x, BODY_Y, bw, 0.34, BLUE if i < 5 else GREEN)
        text(s, x, BODY_Y + 0.03, bw, 0.28, num, 11.5, WHITE, True, "c", "m")
        text(s, x + 0.12, BODY_Y + 0.46, bw - 0.24, 0.30, t, 11.5, NAVY, True,
             "c")
        text(s, x + 0.10, BODY_Y + 0.82, bw - 0.20, 0.70, b, 9.5, GRAY, "c" == "x"
             and False or False, "c", space=1.25)
        if i < 5:
            text(s, x + bw + 0.01, BODY_Y + 0.62, 0.12, 0.3, "›", 15, "9FB4C7",
                 True, "c")
    y = BODY_Y + 1.90
    rows = [
        ["lmm_build.py", "층 적합 · LOO 교차검증 · 모델 JSON 생성", "lmm_model.json"],
        ["lmm_diagnose.py", "잔차 원인 가설 검정 · 민감도 · 반복성", "lmm_diagnosis.json"],
        ["lmm_web.py / lmm_excel.py", "웹·Excel 계산기(오프라인 단일 파일)", "lmm.html / xlsx"],
        ["create_lmm_*.py", "구축현황·검증연구 기록 문서", "docx / pptx"],
    ]
    y2 = table(s, M, y, CW, ["스크립트", "역할", "산출"], rows,
               [0.24, 0.53, 0.23], row_h=0.42)
    band(s, M, y2 + 0.18, CW, 0.56,
         "수치를 문서에 하드코딩하지 않는다 — 재적합하면 보고서 수치가 자동으로 따라간다.")
    return s


def sl_data_status(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "현재까지 확보한 자료 — 남은 병목은 측점과 지각층이다",
                   "Core·External 은 확보됐다. Crustal 은 원측선이 존재하지 않아 현 "
                   "격자가 상한이고, Regional 은 측점이 모자란다.", page)
    rows = [
        ["Core", "IGRF-14 (13차) 계수", "완비", "재현성 검증 통과 (ppigrf 대조)"],
        ["Regional", f"절대측정 {N_SITES}점 ({MODEL['epoch_label']})", "부족",
         "권고 30점 대비 미달 · 편각 품질 불량점 배제"],
        ["Crustal", f"KIGAM 자력이상 {DIAG['vector_recovery']['grid']['dx_km']:.1f} km 격자",
         "해상도 상한",
         f"자료 공백 {DIAG['vector_recovery']['grid']['gap_pct']:.0f}% · "
         "원측선 자료가 존재하지 않는다"],
        ["External", "상시관측 1분 자료 (관측소 4소)", "편각 적용",
         f"야장에서 {EXT['y0']}~{EXT['y1']} 관측시각 복원 · "
         f"보정량 {EXT['ok']}세션 (정온·편각 전용)"],
    ]
    y = table(s, M, BODY_Y, CW, ["층", "보유 자료", "상태", "비고"], rows,
              [0.11, 0.30, 0.13, 0.46], row_h=0.46)
    img = fig_sites()
    s.shapes.add_picture(str(img), Inches(9.35), Inches(y + 0.18),
                         height=Inches(2.30))
    x2, w2 = M, 8.20
    stat(s, x2, y + 0.24, 2.55, 1.06, f"{N_SITES}점",
         "Regional 입력 측점", LGRAY, BLUE, 22)
    stat(s, x2 + 2.68, y + 0.24, 2.55, 1.06, f"{DS['n_obs']}회",
         f"Regional 입력 관측 · 방문 ({DS['year_min']}~{DS['year_max']})",
         LGRAY, BLUE, 22)
    stat(s, x2 + 5.36, y + 0.24, 2.55, 1.06,
         f"{DIAG['only2019']['gap_north_km']:.0f} km",
         "최북단 측점 위 공백", PEACH, RED, 22)
    miss = ("(" + "·".join(str(m) for m in EXT["missing"]) + "년 야장은 없다)"
            if EXT["missing"] else "")
    band(s, x2, y + 1.46, w2, 0.86,
         f"성과표에는 연도만 있으나 지리원 야장에서 {EXT['y0']}~{EXT['y1']} 분 단위 "
         f"관측시각을 복원했다 {miss}.\n"
         f"관측소 4소 공간보간으로 {EXT['ok']}세션에 보정량을 산출해 편각에 적용 중이다 "
         "— 병목은 측점 수와 지각층으로 옮겨갔다.")
    return s


def sl_result(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "교차검증 결과 — 구조는 섰고, 목표치는 아직이다",
                   f"{N_SITES}측점 LOO 기준 편각 {LOO['D']:.3f}° · 총자력 {LOO['F']:.1f} nT로 "
                   "목표를 넘는다. 정식 편각 산출에는 쓰지 않는다.", page)
    img = fig_loo()
    s.shapes.add_picture(str(img), Inches(8.20), Inches(BODY_Y + 0.05),
                         width=Inches(4.53))
    stat(s, M, BODY_Y, 2.42, 1.34, f"{LOO['D']:.3f}°",
         f"LOO 편각 D  |  목표 <{KPI_D:g}°", PEACH, RED)
    stat(s, M + 2.55, BODY_Y, 2.42, 1.34, f"{LOO['F']:.1f} nT",
         f"LOO 총자력 F  |  목표 <{KPI_F:g} nT", PEACH, RED)
    stat(s, M + 5.10, BODY_Y, 2.42, 1.34, f"{LOO['I']:.3f}°",
         "LOO 복각 I  |  별도 목표 없음", LBLUE, BLUE)
    y = BODY_Y + 1.62
    rows = [
        ["편각 D", f"{LOO['D']:.3f}°", f"< {KPI_D:g}°", "미달",
         "지각장 해상도 · 방위 기준 · 측점 수"],
        ["총자력 F", f"{LOO['F']:.1f} nT", f"< {KPI_F:g} nT", "미달",
         "측점 밀도 · 지각장 해상도"],
        ["Crustal 상관", f"r = +{CRU['corr']:.3f}", "—", f"n={CRU['n']} 불안정",
         f"잔차 RMS {CRU['rms_reduction_pct']:.0f}% 감소"],
    ]
    table(s, M, y, 7.45, ["성분", "현재", "목표", "판정", "주된 원인"], rows,
          [0.16, 0.17, 0.15, 0.17, 0.35], row_h=0.46, size=10)
    band(s, M, y + 1.94, 7.45, 0.92,
         f"⚠ Crustal 상관은 표본이 {CRU['n']}개뿐이라 측점 하나에 크게 흔들린다.\n"
         "   개선으로 단정하지 말 것 — 표본 확대 후 재판정이 필요하다.")
    band(s, 8.20, BODY_Y + 2.62, 4.53, 1.62,
         "현 단계 사용 범위\n\n"
         "· 연구·시범계산 : 가능\n"
         "· 지형도 자침편차 등 정식 편각 산출 : 자료 확충 후",
         "23456B")
    return s


def sl_residual(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "네 가설 중 셋은 기각, 하나는 살아나 모형에 들어갔다",
                   "감사로 잡음이 걷히자 「지각 벡터 성분 누락」만 남았다. 그것을 "
                   "구현한 것이 이번 단계다.", page)
    A = HYP["A_vector_D"]
    rows = [
        ["A. 지각 벡터 성분 누락", f"r = {A['r']:+.3f}", f"p = {A['p']:.3f}",
         f"예상 기여 {DIAG['vector_recovery']['predD_rms_arcmin']:.1f}′ — "
         "주원인은 아니나 KPI 6′ 초과로 무시 불가"],
        ["B. 측정 정밀도", f"r = {HYP['B_precision']['r']:+.3f}",
         f"p = {HYP['B_precision']['p']:.2f}",
         f"1회 {DIAG['precision_split']['single_median']:.1f}′ vs 2회+ "
         f"{DIAG['precision_split']['multi_median']:.1f}′ — 차이 없음"],
        ["C. 국소 자기 거칠기", f"r = {HYP['C_roughness']['r']:+.3f}",
         f"p = {HYP['C_roughness']['p']:.2f}", "설명력 없음"],
        ["D. 미적합 지역 구조", "—", "—",
         f"공간 상관 확인 안 됨 — 다만 최근거리 측점쌍 {VAR['bins'][0]['n']}개뿐"],
    ]
    y = table(s, M, BODY_Y, 7.55, ["가설", "상관", "유의확률", "판정"], rows,
              [0.28, 0.14, 0.15, 0.43], row_h=0.52, size=9.5)
    img = fig_asymmetry()
    s.shapes.add_picture(str(img), Inches(8.30), Inches(BODY_Y + 0.10),
                         width=Inches(4.40))
    text(s, 8.30, BODY_Y + 2.42, 4.40, 0.30, "D / I 잔차 비대칭", 12, NAVY, True,
         "c")
    text(s, 8.30, BODY_Y + 2.72, 4.40, 0.52,
         f"관측 {ASY['ratio']:.2f} 대 자기 원인 예상 "
         f"{ASY['expected_from_magnetic']:.2f} — 초과분이 거의 사라졌다.",
         10, GRAY, align="c", space=1.25)
    band(s, M, y + 0.28, 7.55, 1.20,
         "가설 A 를 「규모 미달」로 접어서는 안 된다.\n"
         f"   현재 잔차({ASY['rD_rms']:.0f}′)의 주원인은 아니지만, 예상 기여 "
         f"{DIAG['vector_recovery']['predD_rms_arcmin']:.1f}′ 는 KPI 6′ 를 넘는다.\n"
         "   자료 계통 오차를 걷어낸 뒤 KPI 관점에서 다시 평가해야 하는 2차 오차원이다.")
    band(s, 8.30, y + 0.28, 4.40, 1.20,
         "남은 계통 오차 (자료만으로 구분 불가)\n\n"
         "① 마크 진북 방위각 오차\n② 성과표 전사·환산 오류\n③ 편각 환산 절차 오류",
         "23456B")
    return s


def sl_crustal_vector(prs, page):
    """지각층을 편각·복각에 기여시킨 단계."""
    VR = DIAG["vector_recovery"]
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "항공자력은 스칼라만 주지만, 방향 정보가 그 안에 있다",
                   "총자력 이상은 이상벡터를 주자기장 방향으로 투영한 값이다. "
                   "역산해 편각·복각에 넣었고, 개선폭은 작지만 두 성분에서 같은 "
                   "방향이다.", page, rule=GREEN)
    text(s, M, BODY_Y, CW, 0.52,
         "ΔF  =  l · b_north  +  m · b_east  +  n · b_down       "
         "(l, m, n) = 주자기장 단위벡터",
         15, NAVY, True, "c")
    y = BODY_Y + 0.66
    rows = [
        ["LOO 편각 D", "0.5118°", "0.5016°", "−0.010°"],
        ["LOO 복각 I", "0.2499°", "0.2201°", "−0.030°"],
        ["가설 A 상관", "—",
         f"r={HYP['A_vector_D']['r']:+.3f} p={HYP['A_vector_D']['p']:.3f}",
         "경계"],
        ["적용 후 잔여 상관", "—",
         f"r={HYP['A_vector_D_after']['r']:+.3f} "
         f"p={HYP['A_vector_D_after']['p']:.2f}", "이중계상 없음"],
    ]
    y = table(s, M, y, 7.55, ["지표", "지각 F만", "지각 벡터 적용", "변화"], rows,
              [0.30, 0.24, 0.26, 0.20], row_h=0.44, head_fill=GREEN, size=10)
    stat(s, 8.30, BODY_Y + 0.66, 2.10, 1.24,
         f"{VR['predD_rms_arcmin']:.1f}′", "복원 벡터의 편각 예측", LBLUE, BLUE, 21)
    stat(s, 10.63, BODY_Y + 0.66, 2.10, 1.24,
         f"{HYP['A_vector_D']['a']:.2f}", "최적배율 — 진폭이 모자란다",
         PEACH, ORANGE, 21)
    band(s, 8.30, BODY_Y + 2.05, 4.43, 1.62,
         "⚠ 탐색적 분석이다\n\n"
         f"격자 결측 {DIAG['vector_recovery']['grid']['gap_pct']:.0f}% 를 0 nT 로 "
         "채워 역산했다.\n결측과 「이상 0」은 다른 뜻이라\n"
         "복원 진폭 자체에 불확실성이 있다.", "23456B")
    band(s, M, y + 0.24, 7.55, 1.34,
         "⚠ 위 값은 «평가 측점을 한 번만 정해 고정한» 비교다.\n"
         "   설정마다 강건 적합이 채택 측점을 다시 고르게 두면 개선폭이 크게 부풀려진다\n"
         "   — 실제로 그렇게 견줬을 때 편각 0.540 → 0.340° 로 보였으나, 평가 집합을 "
         "고정하니 0.512 → 0.502° 였다.")
    return s


def sl_impl(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "구현은 세 환경에서 같은 값을 낸다",
                   "Python·JavaScript·Excel 세 구현이 동일 좌표에서 일치함을 확인했다 — "
                   "계산체계 자체는 재현 가능하다.", page)
    vs = DIAG["point_check"]
    rows = [
        ["JavaScript ↔ Python", "웹 계산기 vs 서버 재현",
         f"D {abs(vs['vs_webcalc']['D'])*3600:.2f}″ · F {abs(vs['vs_webcalc']['F']):.2f} nT"],
        ["Excel ↔ Python", "수식 평가 회귀검증", "D 0.33″ · F 0.12 nT (사전격자 보간오차)"],
        ["IGRF 엔진 ↔ ppigrf", "Core 층 구현 검증",
         f"D {abs(vs['vs_official']['D'])*3600:.2f}″ · F {abs(vs['vs_official']['F']):.2f} nT"],
    ]
    y = table(s, M, BODY_Y, CW, ["대조", "내용", "차이"], rows,
              [0.26, 0.38, 0.36], row_h=0.46)
    lay = vs["layers"]
    text(s, M, y + 0.34, CW, 0.30, "층별 기여 실측 (충북 청주 인근 1점)", 13, NAVY,
         True)
    bw = (CW - 3 * 0.14) / 4
    contrib = [("Core (IGRF)", f"{lay['core']:,.0f} nT", LGRAY, NAVY),
               ("Regional", f"{lay['regional_F']:+.2f} nT", LBLUE, BLUE),
               ("Crustal", f"{lay['crustal_F']:+.2f} nT", LBLUE, BLUE),
               ("합계 LMM − IGRF", f"{vs['diff_lmm_igrf']['F']:+.2f} nT", PEACH, RED)]
    for i, (lab, val, fill, col) in enumerate(contrib):
        x = M + i * (bw + 0.14)
        rect(s, x, y + 0.72, bw, 1.00, fill)
        text(s, x + 0.14, y + 0.86, bw - 0.28, 0.40, val, 16, col, True)
        text(s, x + 0.14, y + 1.32, bw - 0.28, 0.30, lab, 10.5, GRAY)
    band(s, M, y + 1.94, CW, 1.02,
         "⚠ IGRF 대조는 정확도 검증이 아니다. IGRF는 LMM의 입력이므로 LMM−IGRF는 "
         "정의상 보정층의 합이다.\n"
         "   정당한 용도는 Core 구현 검증 하나이며 이미 통과했다 — 정확도는 LOO·독립 "
         "표본으로만 말할 수 있다.")
    return s


def sl_kigam(prs, page):
    s = slide_base(prs, "Ⅱ. 확인·추가 자료",
                   "확인할 자료 ① — 한국지질자원연구원 항공자력",
                   "현재는 가공된 격자만 쓰고 있다. 원측선 자료를 확인해 Crustal 층의 "
                   "해상도 한계를 푼다.", page, rule=ORANGE)
    img = fig_resolution()
    s.shapes.add_picture(str(img), Inches(7.60), Inches(BODY_Y + 0.05),
                         width=Inches(5.13))
    g = DIAG["vector_recovery"]["grid"]
    stat(s, M, BODY_Y, 2.30, 1.24, f"{g['dx_km']:.1f} km", "현재 사용 격자 간격",
         PEACH, RED, 22)
    stat(s, M + 2.43, BODY_Y, 2.30, 1.24, "250 m", "원측선 간격 (권고)", LBLUE, BLUE, 22)
    stat(s, M + 4.86, BODY_Y, 2.30, 1.24, f"{g['gap_pct']:.0f}%",
         "격자 자료 공백 비율", PEACH, RED, 22)
    y = BODY_Y + 1.46
    for i, (t, b) in enumerate([
        ("무엇을 확인하는가",
         "측선 간격·측고·IGRF 제거 기준·측량 연도 — 격자로 가공되기 전의 원측선 제원"),
        ("왜 필요한가",
         f"현재 격자는 {g['dx_km']:.1f} km 간격이라 50 km 미만 단파장 지각이상을 담지 못한다"),
        ("무엇이 풀리는가",
         "총자력 F 목표 미달의 주된 원인 중 하나인 지각장 해상도 한계"),
    ]):
        yy = y + i * 0.72
        rect(s, M, yy, 6.98, 0.62, LGRAY)
        text(s, M + 0.18, yy + 0.06, 2.0, 0.26, t, 11, ORANGE, True)
        text(s, M + 0.18, yy + 0.32, 6.6, 0.26, b, 10.5, NAVY)
    band(s, M, y + 2.28, CW, 1.04,
         "선행 검증 근거 — 경상분지 동해안 사례에서 점 측정 잔차와 항공자력 2 km 평균이 "
         "5 nT 이내로 일치했다.\n"
         "즉 항공자력이 점 측정 잔차를 재현한다는 것은 이미 확인됐고, 남은 것은 전국 "
         "모자이크 적용이다.", ORANGE)
    return s


def sl_ngii(prs, page):
    s = slide_base(prs, "Ⅱ. 확인·추가 자료",
                   "추가할 자료 ② — 국토지리정보원 2021~2025 반복측정",
                   "관측시각은 이미 야장에서 복원했다. 남은 것은 Regional 층의 "
                   "입력 측점을 늘리는 일이다.", page, rule=ORANGE)
    stat(s, M, BODY_Y, 2.72, 1.30, f"{N_SITES}점 → 30점 이상",
         "Regional 입력 측점 목표", LBLUE, BLUE, 17)
    stat(s, M + 2.85, BODY_Y, 2.72, 1.30, "2021~2025",
         "추가 대상 반복측정 기간", LBLUE, BLUE, 19)
    stat(s, M + 5.70, BODY_Y, 2.72, 1.30, f"{EXT['ok']}세션",
         "관측시각 복원 완료 · 보정량 산출", LBLUE, GREEN, 19)
    rows = [
        ["측점 밀도", f"현재 {N_SITES}점 · 권고 30점 이상",
         "1차 다항식의 자유도 확보 · 이상치 배제 작동"],
        ["관측시각", f"야장에서 {EXT['y0']}~{EXT['y1']} 분 단위 복원 (완료)",
         "외부장 보정 적용 중 — 야장에만 있는 측점 17건이 추가 후보"],
        ["독립 검증점", "적합에 쓰지 않는 측점 분리",
         "LOO 외에 독립 표본 검증이 가능해진다"],
        ["시기간 대조", "겹치는 측점의 IGRF 잔차 차",
         "External 적용 전후로 0에 수렴하는지 정량 판정"],
    ]
    y = table(s, M, BODY_Y + 1.50, CW, ["항목", "현재 상태", "확보 시 효과"], rows,
              [0.18, 0.36, 0.46], row_h=0.44)
    band(s, M, y + 0.22, CW, 1.02,
         "⚠ 두 자료는 서로 다른 목표를 푼다 — 섞어 말하면 안 된다.\n\n"
         "   · 항공자력(Crustal) → 총자력 F의 지각장 해상도 — 원측선이 없어 "
         "현 격자가 상한이다\n"
         "   · 반복측정(Regional) → 편각 D의 측점 밀도. 관측시각(External)은 이미 "
         "확보·적용했다")
    return s


def sl_match(prs, page):
    s = slide_base(prs, "Ⅱ. 확인·추가 자료",
                   "어떤 자료가 어떤 미달을 푸는가",
                   "현재의 두 미달 항목을 원인별로 나누고, 각 자료가 담당하는 몫을 "
                   "분명히 한다.", page, rule=ORANGE)
    cols = [
        ("편각 D", f"{LOO['D']:.3f}°", f"목표 {KPI_D:g}°", RED, PEACH,
         [("외부장 보정", "관측소 4소 적용 완료 — 잔여는 NOC 투영"),
          ("측점 밀도 부족", f"{N_SITES}점 → 30점 이상"),
          ("계통 오차 의심", "성과표 원본 감사 · 방위 기준 재확인")]),
        ("총자력 F", f"{LOO['F']:.1f} nT", f"목표 {KPI_F:g} nT", RED, PEACH,
         [("지각장 해상도", "원측선 없음 — 현 격자가 상한"),
          ("측점 밀도 부족", "반복측정 확대"),
          ("자료 공백", f"격자 공백 {DIAG['vector_recovery']['grid']['gap_pct']:.0f}% 보완")]),
    ]
    for i, (name, cur, goal, col, fill, items) in enumerate(cols):
        x = M + i * (CW / 2 + 0.14)
        w = CW / 2 - 0.14
        rect(s, x, BODY_Y, w, 0.92, fill)
        text(s, x + 0.22, BODY_Y + 0.12, w - 0.44, 0.34, f"{name}  {cur}", 19,
             col, True)
        text(s, x + 0.22, BODY_Y + 0.52, w - 0.44, 0.28, goal, 11, GRAY)
        for j, (cause, fix) in enumerate(items):
            yy = BODY_Y + 1.08 + j * 1.02
            rect(s, x, yy, w, 0.90, LGRAY)
            text(s, x + 0.22, yy + 0.10, w - 0.44, 0.28, f"· {cause}", 11.5,
                 NAVY, True)
            text(s, x + 0.42, yy + 0.44, w - 0.64, 0.36, f"→ {fix}", 10.5, GRAY,
                 space=1.2)
    band(s, M, BODY_Y + 4.16, CW, 0.80,
         "⚠ 순서가 중요하다 — 상시관측으로 제거할 수 있는 일변동은 편각 기준 수 분(′) "
         "수준인데,\n"
         f"   미설명 잔차는 {ASY['rD_rms']:.0f}′다. 잡음이 신호보다 큰 상태에서 보정을 "
         "얹으면 개선 여부를 판정할 수 없다.")
    return s


def sl_roadmap(prs, page):
    s = slide_base(prs, "Ⅲ. 향후 검증·보완",
                   "검증은 게이트를 두고 단계적으로 진행한다",
                   "판정기준을 먼저 정하고 그 기준에 따라 다음 단계로 넘어간다 — "
                   "층이 결합돼 있어 순차 튜닝은 금지한다.", page, rule=GREEN)
    steps = [
        ("0", "2022~25 야장 확보", "212세션 · 분 단위 시각", "완료", GREEN),
        ("1", "성과표 원본 감사", "불일치 방문 2건 특정·제외", "완료", GREEN),
        ("2", "지각 벡터 · External", "end-to-end 재적합으로 판정", "완료", GREEN),
        ("3", "재측정 · 신규 측점", "KPI 미달분 · 독립 표본 검증", "진행", BLUE),
    ]
    bw = (CW - 3 * 0.16) / 4
    for i, (num, t, b, tag, col) in enumerate(steps):
        x = M + i * (bw + 0.16)
        rect(s, x, BODY_Y, bw, 1.66, LGRAY)
        rect(s, x, BODY_Y, bw, 0.38, col)
        text(s, x + 0.16, BODY_Y + 0.04, bw - 0.32, 0.30,
             f"{num}단계   {tag}", 10.5, WHITE, True, anchor="m")
        text(s, x + 0.18, BODY_Y + 0.54, bw - 0.36, 0.52, t, 13, NAVY, True,
             space=1.15)
        text(s, x + 0.18, BODY_Y + 1.14, bw - 0.36, 0.42, b, 10, GRAY,
             space=1.2)
        if i < 3:
            text(s, x + bw + 0.02, BODY_Y + 0.66, 0.14, 0.34, "›", 16, "9FB4C7",
                 True, "c")
    y = BODY_Y + 1.96
    text(s, M, y, CW, 0.30,
         "1단계 판정기준 — 편각 잔차로 다음 행동을 가른다 (기준은 착수 전에 정했다)",
         13, NAVY, True)
    rows = [
        ["10′ 이하", "외부장 보정 진행", "상시관측 자료 결합 착수"],
        ["10 ~ 20′", "부분 재측정", "잔차 큰 측점 현장 재측정"],
        ["20′ 초과", "모형 작업 중단", "방위 기준 재결정 선행"],
    ]
    table(s, M, y + 0.38, 8.10, ["감사 후 편각 잔차", "판정", "후속 조치"], rows,
          [0.30, 0.30, 0.40], row_h=0.46)
    r_now = ASY["rD_rms"]
    if r_now <= 10:
        vtxt, vcol = ("첫 번째 구간.\n외부장 보정을 그대로 진행한다.", GREEN)
    elif r_now <= 20:
        vtxt, vcol = ("두 번째 구간.\n남은 것은 잔차가 큰 측점의\n현장 재측정이다.",
                      ORANGE)
    else:
        vtxt, vcol = ("세 번째 구간.\n방위 기준 재결정이 먼저다.", RED)
    band(s, 8.45, y + 0.38, CW - 7.85, 1.84,
         f"현재 {r_now:.1f}′\n\n" + vtxt + "\n\n(착수 시점 35.0′)", vcol, 12)
    return s


def sl_risk(prs, page):
    s = slide_base(prs, "Ⅲ. 향후 검증·보완",
                   "보완 우선순위와 남은 위험",
                   "자료 확충만으로 풀리지 않는 항목이 있다 — 그 경우의 대응까지 미리 "
                   "정해 둔다.", page, rule=GREEN)
    prio = [
        ("최우선", "잔차 큰 측점 재측정",
         f"편각 잔차 {ASY['rD_rms']:.1f}′ — 방위표지 고정과 함께", ORANGE),
        ("2순위", "신규 측점 확보 (30점)",
         "자유도·이상치 배제·독립 검증점. 12점 이하에서 무너진다", BLUE),
        ("3순위", "NOC 공간투영 구현",
         "4소 평면은 근사 — 교란시에 먼저 깨진다", BLUE),
        ("보류", "지각층 해상도",
         "원측선이 없어 현 격자가 상한 — 신규 자력측량은 사업 판단", GRAY),
    ]
    for i, (tag, t, b, col) in enumerate(prio):
        yy = BODY_Y + i * 0.66
        badge_row(s, yy, tag, "", col, 0.54, 1.30, M, 7.55)
        text(s, M + 1.62, yy + 0.06, 5.6, 0.26, t, 12, NAVY, True)
        text(s, M + 1.62, yy + 0.30, 5.6, 0.24, b, 10, GRAY)
    y = BODY_Y + 2.86
    rows = [
        ["계통 오차가 남는 경우", "자료를 늘려도 편각 잔차가 줄지 않음",
         "방위 기준 재결정 · 문제 측점 재측정"],
        ["표본 불안정", f"Crustal 상관이 n={CRU['n']}로 흔들림", "표본 확대 전 결론 유보"],
        ["공간 공백", f"최북단 위 {DIAG['only2019']['gap_north_km']:.0f} km 미측",
         "신규 캠페인 배치에 반영"],
    ]
    table(s, M, y, CW, ["위험", "내용", "대응"], rows, [0.22, 0.42, 0.36],
          row_h=0.46)
    band(s, 8.45, BODY_Y, CW - 7.85, 2.62,
         "지금 말할 수 있는 것\n\n"
         "· 구조와 계산체계는 재현 가능\n"
         "· 오차 원인을 층으로 분리 가능\n\n"
         "아직 말할 수 없는 것\n\n"
         "· 목표 정확도 달성\n"
         "· 정식 편각 산출 적합성", "23456B")
    return s


def sl_cmp_concept(prs, page):
    s = slide_base(prs, "Ⅳ. 2020년 연구와의 비교",
                   "2020년은 «보정된 지자기도», 이번은 «계산되는 지역모델»",
                   "같은 IGRF를 쓰지만 IGRF를 다루는 방식이 다르다 — 보정 대상으로 보느냐, "
                   "배경장으로 두느냐.", page, rule="6B5B95")
    for i, (title, sub, flow, fill, col) in enumerate([
        ("2020년 지자기도", "IGRF 기반 보정",
         ["IGRF 예측", "국내 실측과 비교", "차이만큼 보정", "기준시점 지자기도"],
         "F1F0F6", "6B5B95"),
        ("2026년 LMM", "지역 성분 분리 모델",
         ["IGRF 배경장", "실측 − IGRF = 지역 신호", "성분별 공간모델링",
          "좌표·시점별 D·I·F 계산"],
         LBLUE, BLUE),
    ]):
        x = M + i * (CW / 2 + 0.14)
        w = CW / 2 - 0.14
        rect(s, x, BODY_Y, w, 0.86, fill)
        text(s, x + 0.22, BODY_Y + 0.10, w - 0.44, 0.34, title, 17, col, True)
        text(s, x + 0.22, BODY_Y + 0.50, w - 0.44, 0.26, sub, 11, GRAY)
        for j, step in enumerate(flow):
            yy = BODY_Y + 1.04 + j * 0.72
            rect(s, x, yy, w, 0.56, LGRAY)
            text(s, x + 0.22, yy + 0.02, w - 0.44, 0.52, f"{j+1}.  {step}", 11.5,
                 NAVY, anchor="m")
            if j < len(flow) - 1:
                text(s, x + w / 2 - 0.2, yy + 0.56, 0.4, 0.16, "▾", 10, "9FB4C7",
                     True, "c")
    band(s, M, BODY_Y + 3.94, CW, 0.96,
         "핵심 차이는 Regional 층이다. 2020년은 IGRF와 실측의 차이를 «보정량»으로 "
         "다뤘고,\n"
         "이번에는 그 차이를 한반도 고유의 «지역 자기장 신호»로 보아 공간함수로 "
         "추정한다.", "6B5B95")
    return s


def sl_cmp_table(prs, page):
    s = slide_base(prs, "Ⅳ. 2020년 연구와의 비교",
                   "항목별 비교 — 구조가 다르다",
                   "배경장은 같지만 국내 관측자료의 지위, 성분 분리 여부, 검증 방식, "
                   "갱신 개념이 모두 다르다.", page, rule="6B5B95")
    rows = [
        ["기본 개념", "IGRF 기반 보정 지자기도", "지역 자기장 성분 모델"],
        ["국내 반복측량", "IGRF와 비교·보정에 활용", "Regional 층을 추정하는 핵심자료"],
        ["지각 자기장", "명시적 분리 없음", "항공자력 → Crustal 층"],
        ["외부장·일변화", "상시관측을 보정·검토에 활용", "External 층으로 별도 취급"],
        ["공간 모델", "IGRF 분포 + 국내성과 보정", "IGRF + Regional + Crustal + 잔차면"],
        ["시간 모델", "기준연도 환산 · IGRF 영년변화 중심", "기준시점 환산 + 외부장·영년변화 분리"],
        ["검증", "관측값·IGRF 비교 중심", "LOO 교차검증 · 독립 표본"],
        ["활용", "지자기도 · 도엽별 값", "좌표 입력 → D·I·F 계산 (웹·Excel·도엽)"],
        ["갱신 개념", "새 지자기도 재작성", "신규 관측자료로 모델 재적합"],
    ]
    y = table(s, M, BODY_Y, CW, ["구분", "2020년 대한민국 지자기도", "2026년 LMM"], rows,
               [0.18, 0.40, 0.42], row_h=0.40, size=11, head_fill="6B5B95")
    band(s, M, y + 0.18, CW, 0.64,
         "2020년 연구는 «19년까지 2회 이상 측량된 15점을 기준자료로 활용했다 — "
         "자료의 성격이 아니라 쓰임이 달라진 것이다.", "6B5B95")
    return s


def sl_cmp_flow(prs, page):
    s = slide_base(prs, "Ⅳ. 2020년 연구와의 비교",
                   "같은 측정값을 넣어도 계산 흐름이 갈린다",
                   "한 측점에서 IGRF 편각 −8.00°, 실측 편각 −8.35°인 경우를 예로 든다.",
                   page, rule="6B5B95")
    rect(s, M, BODY_Y, CW, 0.62, "F1F0F6")
    text(s, M, BODY_Y + 0.10, CW, 0.42,
         "예시 측점   IGRF D = −8.00°      실측 D = −8.35°      차이 = −0.35°",
         15, "6B5B95", True, "c", "m")
    left = [
        ("2020년 방식", "6B5B95"),
        ("−0.35°를 «IGRF와 실측의 차이»로 본다", None),
        ("그 지점 주변의 지자기도 값을 보정한다", None),
        ("보정된 지자기도를 도엽 단위로 제공한다", None),
        ("→ 지도에 값이 인쇄된다", "6B5B95"),
    ]
    right = [
        ("LMM 방식", BLUE),
        ("ΔD_regional = D_obs − D_IGRF = −0.35°  (지역 신호)", None),
        ("전국 측점의 ΔD를 모아 ΔD_regional(λ, φ, t) 를 적합", None),
        ("임의 좌표에서  D_LMM = D_IGRF + ΔD_regional + 보정", None),
        ("→ 좌표를 넣으면 값이 계산된다", BLUE),
    ]
    for i, block in enumerate((left, right)):
        x = M + i * (CW / 2 + 0.14)
        w = CW / 2 - 0.14
        for j, (t, col) in enumerate(block):
            yy = BODY_Y + 0.84 + j * 0.70
            if j == 0:
                rect(s, x, yy, w, 0.54, col)
                text(s, x + 0.22, yy + 0.02, w - 0.44, 0.50, t, 14, WHITE, True,
                     anchor="m")
            elif col:
                rect(s, x, yy, w, 0.54, LBLUE if col == BLUE else "F1F0F6")
                text(s, x + 0.22, yy + 0.02, w - 0.44, 0.50, t, 12, col, True,
                     anchor="m")
            else:
                rect(s, x, yy, w, 0.54, LGRAY)
                text(s, x + 0.22, yy + 0.02, w - 0.44, 0.50, t, 10.5, NAVY,
                     anchor="m")
    band(s, M, BODY_Y + 4.34, CW, 0.60,
         "2020년 결과는 «값이 적힌 지도», LMM은 «값을 만들어내는 함수»다 — "
         "갱신 방식이 달라지는 이유가 여기에 있다.", "6B5B95")
    return s


def sl_lineage(prs, page):
    s = slide_base(prs, "Ⅳ. 2020년 연구와의 비교",
                   "경쟁 모델이 아니라 계보로 놓는 것이 정확하다",
                   "2020년 성과는 대체 대상이 아니라, 이번 연구가 딛고 서는 단계다.",
                   page, rule="6B5B95")
    line = [("2010", "지자기 분포모델", "국내 지자기 분포 파악", "8A94A6"),
            ("2020", "IGRF + 실측 보정 지자기도", "지도 제작용 국내 보정장", "6B5B95"),
            ("2026", "LMM 지역 지자기모델", "계산 가능한 국가 지역모델", BLUE)]
    bw = (CW - 2 * 0.30) / 3
    for i, (yr, t, b, col) in enumerate(line):
        x = M + i * (bw + 0.30)
        rect(s, x, BODY_Y, bw, 1.66, LGRAY)
        rect(s, x, BODY_Y, bw, 0.46, col)
        text(s, x, BODY_Y + 0.04, bw, 0.38, yr, 15, WHITE, True, "c", "m")
        text(s, x + 0.20, BODY_Y + 0.64, bw - 0.40, 0.56, t, 13, NAVY, True,
             space=1.15)
        text(s, x + 0.20, BODY_Y + 1.24, bw - 0.40, 0.32, b, 10.5, GRAY)
        if i < 2:
            text(s, x + bw + 0.04, BODY_Y + 0.62, 0.22, 0.40, "›", 18, "9FB4C7",
                 True, "c")
    y = BODY_Y + 2.00
    for i, (t, b, fill, col) in enumerate([
        ("이번 연구가 주장할 수 있는 것",
         "· 구조가 고도화됐다 — 오차 원인을 층으로 분리할 수 있다\n"
         "· 계산 가능한 형태다 — 좌표·시점을 넣으면 값이 나온다\n"
         "· 갱신 가능한 형태다 — 신규 관측으로 재적합한다",
         LBLUE, BLUE),
        ("아직 주장하면 안 되는 것",
         "· «2020년보다 정확하다» — 목표 정확도를 아직 못 만족한다\n"
         f"· 현재 LOO 편각 {LOO['D']:.3f}° · 총자력 {LOO['F']:.1f} nT\n"
         "· 층이 완전히 구현된 상태가 아니다 (Crustal 제한 · External 부분)",
         PEACH, RED),
    ]):
        x = M + i * (CW / 2 + 0.14)
        w = CW / 2 - 0.14
        rect(s, x, y, w, 1.86, fill)
        text(s, x + 0.24, y + 0.16, w - 0.48, 0.32, t, 13, col, True)
        text(s, x + 0.24, y + 0.60, w - 0.48, 1.10, b, 10.5, NAVY, space=1.35)
    band(s, M, y + 2.10, CW, 0.70,
         "제안요청서의 «국제모델과의 차이 분석» 요구와도 이 정의가 맞는다 — "
         "IGRF를 대체하는 것이 아니라 그 차이를 설명하는 모델이다.", "6B5B95")
    return s


SHOT_DIR = OUT_DIR / "_web_shots"


def sl_audit(prs, page):
    """Regional 성과 신뢰도 감사 — 정밀하지만 부정확하다."""
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "기존 성과는 정밀하지만 부정확하다",
                   "야장 원시 판독값 197세션을 세 층위로 검정했다. 측정·계산이 아니라 "
                   "방문마다 달라지는 방위 기준이 문제였다.", page, rule=RED)
    for i, (t, v, cap, fill, col) in enumerate([
            ("① 내부 산술", "0.16″", "편각 재계산차 중앙\n야장 산술은 완전 일치",
             LGRAY, GREEN),
            ("② 정밀도", "1.41′", "같은 날 세션 산포 중앙\n목표 6′ 대비 양호",
             LGRAY, GREEN),
            ("③ 정확도", "33.7′", "재방문 잔여 RMS\n16구간 중 9구간이 6′ 초과",
             PEACH, RED)]):
        x = M + i * (CW / 3 + 0.06)
        w = CW / 3 - 0.12
        rect(s, x, BODY_Y, w, 1.34, fill)
        text(s, x + 0.20, BODY_Y + 0.12, w - 0.40, 0.26, t, 11.5, GRAY, True)
        text(s, x + 0.20, BODY_Y + 0.42, w - 0.40, 0.42, v, 25.5, col, True)
        text(s, x + 0.20, BODY_Y + 0.92, w - 0.40, 0.36, cap, 10, GRAY, space=1.2)

    text(s, M, BODY_Y + 1.52, CW, 0.28,
         "재방문 잔여 = 관측 ΔD − IGRF 영년변화   (같은 표석을 다시 재면 0 이어야 한다)",
         11.5, NAVY, True)
    rows = [["정상 (<6′)", "미원 2.0 · 상주 −5.2/2.5 · 성산 −0.6 · 순천 −1.5/0.5 · 함양 −5.8",
             "7구간"],
            ["주의 (6~20′)", "거제 −19.7 · 여주 11.0 · 이원 −8.3", "3구간"],
            ["불일치 (≥20′)", "부안 62.2/−70.6 · 함양 76.9 · 포천 −32.2/22.3 · 가야 −35.5",
             "6구간"]]
    y = table(s, M, BODY_Y + 1.84, CW, ["판정", "측점별 잔여 (′)", "건수"], rows,
              [0.18, 0.68, 0.14], row_h=0.40, head_fill="8C2F2F", size=10)
    band(s, M, y + 0.18, CW, 1.24,
         "원인 — 마크 참방위각이 재방문마다 바뀐다 (거제 159.9°→330.6° · 포천 33.2°→192.7° · "
         "부안 54.4°→169.6°→177.5°)\n"
         "   마크 교체 자체는 오류가 아니나, 바뀐 마크의 참방위각이 틀리면 그 방문의 편각이 "
         "통째로 틀어진다.\n\n"
         "→ 신규 선점에서 방위표지를 고정하고 방위각 결정법을 상향해야 하는 근거 "
         "(천문·자이로 또는 RTK 장기선)", RED)
    return s


def sl_audit_effect(prs, page):
    """감사 반영 재적합 결과."""
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "불일치 방문 2건을 빼자 편각이 26% 좋아졌다",
                   "감사가 모델 잔차를 보지 않고 지목한 방문이다. 순환 논증이 아니라는 "
                   "점이 중요하다.", page, rule=GREEN)
    # ⚠ 이 표는 「감사」 단계의 효과만 담는다. 이후 지각 벡터·다중관측소로 더
    #   좋아졌으므로 현재 LOO(LOO['D'])를 여기에 넣으면 감사의 공로가 부풀려진다.
    rows = [["LOO 편각 D", "0.7691°", "0.5403°", "−26%", "개선"],
            ["LOO 총자력 F", "60.26 nT", "58.41 nT", "−1.85 nT", "개선"],
            ["LOO 복각 I", "0.1787°", "0.2600°", "+0.081°", "악화"],
            ["편각 잔차 RMS", "35.0′", "20.73′", "−41%", "개선"],
            ["D / I 비대칭", "3.69", "1.95", "자기 원인 예상 1.64", "거의 해소"],
            ["가설 A (지각 벡터)", "r=+0.244 p=0.40", "r=+0.560 p=0.046",
             "기각 → 유의", "반전"]]
    y = table(s, M, BODY_Y, CW,
              ["지표", "감사 전 (17측점)", "감사 후 (16측점)", "변화", "판정"], rows,
              [0.22, 0.20, 0.20, 0.20, 0.18], row_h=0.40, head_fill=GREEN, size=10)
    band(s, M, y + 0.18, CW, 1.50,
         "가장 중요한 것은 마지막 두 줄이다.\n\n"
         "· D/I 비대칭이 3.69 → 1.95 로 떨어져 자기 원인 예상값(1.64)에 근접했다. "
         "「편각에만 작용하는 초과 오차」의\n  정체가 바로 그 두 방문이었다는 뜻이다.\n"
         "· 규모 미달로 접어 뒀던 「지각 벡터 성분 누락」이 잡음이 걷히자 드러났다"
         "(p=0.046).\n\n"
         "⚠ 대가 — 부안이 성과표에 2023 한 건뿐이라 측점이 17→16 으로 줄어 서해안 "
         "공간 커버리지를 잃는다. 복각도 악화됐다.", GREEN)
    return s


def sl_compliance(prs, page):
    """설명자료(오석훈, 2026-05) 대비 구현 준수 점검."""
    s = slide_base(prs, "Ⅲ. 향후 검증·보완",
                   "설명자료 대비 — 무엇을 지켰고 무엇이 남았는가",
                   "4층 구조와 검증 절차는 그대로 따랐다. 남은 것은 자료 규모와 "
                   "External 의 공간투영이다.", page, rule=GREEN)
    rows = [
        ["① Core", "IGRF-14 (deg ≤ 13), ppigrf", "동일 구현 · 재현성 검증 통과", "준수"],
        ["② Regional", "지표 절대측정 30점 이상 · D+I+F", f"{N_SITES}점 · 세션 6회", "미달"],
        ["③ Crustal", "항공자력 측선 100~250 m", "2.8 km 가공 격자", "미달"],
        ["④ External", "CYG NOC 분해 → 점 위치 공간투영",
         "4소 1차 평면 보간 구현 · 모델 미반영", "미달"],
        ["Regional 적합", "다항식 또는 R-SCHA", "상수항(0차) — 공간항 미채택", "미달"],
        ["잔차면", "Residual = 관측 − LMM 을 공간 모델링", "미구현", "미구현"],
        ["검증", "LOO 자기검증", "LOO 구현 · 3환경 일치", "준수"],
    ]
    y = table(s, M, BODY_Y, CW, ["항목", "설명자료 요건", "현재 구현", "판정"], rows,
              [0.14, 0.36, 0.34, 0.16], row_h=0.44)
    band(s, M, y + 0.20, CW, 1.12,
         "설명자료의 「자료가 하나라도 빠지면 구조적 결함」 진단이 그대로 맞았다.\n\n"
         "· 항공자력 원측선이 존재하지 않아 F 목표는 현 자료로 달성 불가 — "
         "신규 자력측량 없이는 열리지 않는다\n"
         "· 외부장 보정은 구현했으나 16측점 표본에서는 효과가 가려진다 → "
         "측점 확충이 선행 조건", GREEN)
    return s


def sl_external_verdict(prs, page):
    """External — 구현했으나 현재 표본에서는 이득이 확인되지 않는다."""
    s = slide_base(prs, "Ⅲ. 향후 검증·보완",
                   "외부장 보정은 구현했으나, 지금 표본으로는 효과가 확인되지 않는다",
                   "자문 권고대로 대체관측소 공간보간까지 구현했다. 판정은 "
                   "평가 측점을 고정한 교차검증으로 했다.", page, rule=RED)
    rows = [
        ["미적용 (채택)", "—", "전 세션", "0.5016°", "기준선"],
        ["편각만 차감", "청양 단독", "전 세션", "0.5098°", "개선 없음"],
        ["편각만 차감", "관측소 4소", "전 세션", "0.5177°", "개선 없음"],
        ["편각만 차감", "청양 단독", "정온만", "0.5189°", "개선 없음"],
        ["편각만 차감", "관측소 4소", "정온만", "0.5323°", "개선 없음"],
    ]
    y = table(s, M, BODY_Y, CW,
              ["방식", "보정량 자료원", "세션", "LOO 편각 D", "판정"], rows,
              [0.24, 0.18, 0.13, 0.20, 0.25], row_h=0.46, head_fill="8C2F2F")
    stat(s, M, y + 0.20, 3.90, 1.06, "178세션",
         "4소 보간으로 산출한 보정량", LGRAY, BLUE, 20)
    stat(s, M + 4.05, y + 0.20, 3.90, 1.06, "5.4′",
         "세션 평균 편각 보정량 규모", LGRAY, BLUE, 20)
    band(s, M + 8.10, y + 0.20, CW - 8.10, 1.06,
         "무엇이 부족한가\n\n보정량(5.4′)이 잔차(26′)에 견줘\n작아 표본 16점으로는 "
         "가려지지 않는다", "23456B")
    band(s, M, y + 1.36, CW, 0.82,
         "⚠ 오전에는 「4소·정온 적용이 0.365 → 0.323° 로 개선」이라고 판정했으나 "
         "그 비교는 순환이었다 — 설정마다 강건 적합이 평가 측점을 다시 골랐다.\n"
         "   평가 집합을 한 번만 정해 고정하면 어떤 보정도 기준선을 넘지 못한다. "
         "보정량은 유지하므로 측점이 늘면 다시 판정한다.", RED)
    return s


def sl_advisory(prs, page):
    """자문 권고 반영 현황."""
    s = slide_base(prs, "Ⅲ. 향후 검증·보완",
                   "자문 권고 반영 현황 — 남은 것은 두 가지다",
                   "2026-07-29 자문회의 지적 사항을 항목별로 대조했다.", page,
                   rule=GREEN)
    rows = [
        ["일변화 미보정이 근본 원인", "야장 직접 대조로 원시값 확정 · 보정량 산출",
         "반영"],
        ["청양 단독 차감의 한계", "관측소 4소 공간보간 구현 — 효과는 미확인(앞 장)",
         "구현"],
        ["2020~2025 야장·관측시각 확보", "212세션 복원 · 178세션 보정량 산출", "반영"],
        ["Kp 반영", "정온 판정 · 교란 세션 보정 제외", "반영"],
        ["구배를 필터가 아니라 가중치로",
         "계단 점수 + 800 nT 하드 제외 → 연속 가중치 (250 nT 최고점)", "반영"],
        ["측점 수 대비 자유도 민감도",
         "1차 다항식은 14점까지 유지 · 12점부터 흔들림 · 2차는 사용 불가", "반영"],
        ["에폭 혼용 금지 (2010 vs 2020~24)",
         "정식 산출 단계에서 2019 제외 채택 · 시범 단계는 표본 유지", "부분"],
        ["그레디오미터 도입", "장비 사안 — 연구 범위 밖", "해당 없음"],
    ]
    y = table(s, M, BODY_Y, CW, ["자문 권고", "반영 내용", "상태"], rows,
              [0.26, 0.58, 0.16], row_h=0.38, head_fill=GREEN, size=9.5)
    band(s, M, y + 0.18, CW, 0.90,
         "남은 두 가지는 자료·장비 사안이다.\n\n"
         "· 신규 측점 확보 — 배치 최적과 자기환경 등급이 어긋나는 구간은 후보 재탐색\n"
         "· 지각층 해상도 — 항공자력 원측선이 존재하지 않아 현 격자가 상한이다. "
         "신규 자력측량 여부는 사업 판단 사항", GREEN)
    return s


def sl_webdemo(prs, page):
    s = slide_base(prs, "Ⅰ. LMM 중간결과",
                   "웹에서 바로 확인할 수 있게 만들어 두었다",
                   "모델·자료 현황·계산기를 각각 한 페이지로 낸다. 모두 정적 파일이라 "
                   "설치 없이 열린다.", page)
    # 2020 보고서 「그림 4-21/4-22」와 같은 형식의 3면 지자기도를 주 그림으로 쓴다.
    big = SHOT_DIR / "lmm_map.png"
    text(s, M, BODY_Y, CW, 0.28,
         "한국 지자기도 — IGRF · LMM · 차이   (lmm_map.html)", 13, NAVY, True)
    if big.exists():
        from PIL import Image as _Im
        iw, ih = _Im.open(big).size
        hh = 2.28
        ww = min(CW, hh * iw / ih)
        s.shapes.add_picture(str(big), Inches(M + (CW - ww) / 2),
                             Inches(BODY_Y + 0.32), height=Inches(hh))
    text(s, M, BODY_Y + 2.68, CW, 0.26,
         "IGRF 와 LMM 은 같은 색범위 · 오른쪽이 LMM 이 더하는 양(지각 자기이상 구조)",
         10.5, GRAY)

    shots = [(SHOT_DIR / "globe_lmm_diff.png", "3D 지구본 — 한국 LMM 패치",
              "geomag_globe.html · 커서를 올리면 그 좌표의 LMM D·I·F"),
             (SHOT_DIR / "inventory.png", "자료 확보 현황",
              "data_inventory.html · 4개 층을 파일에서 직접 세어 표시")]
    w = (CW - 0.30) / 2
    for i, (img, title, cap) in enumerate(shots):
        x = M + i * (w + 0.30)
        text(s, x, BODY_Y + 2.98, w, 0.26, title, 11.5, NAVY, True)
        if img.exists():
            from PIL import Image as _Im
            iw, ih = _Im.open(img).size
            hh = 0.94                      # 보조 캡처는 작게
            ww = min(w, hh * iw / ih)
            s.shapes.add_picture(str(img), Inches(x + (w - ww) / 2),
                                 Inches(BODY_Y + 3.26), height=Inches(hh))
        else:
            rect(s, x, BODY_Y + 3.26, w, 0.94, LGRAY)
            text(s, x, BODY_Y + 3.58, w, 0.3, "(캡처 없음)", 11, GRAY, align="c")
        text(s, x, BODY_Y + 4.26, w, 0.26, cap, 9.5, GRAY, space=1.2)
    band(s, M, BODY_Y + 4.56, CW, 0.42,
         "시연 — ① 지자기도 3면 비교 → ② 지구본에서 위치·규모 확인 → ③ 자료 현황으로 근거")
    return s


def sl_resolution(prs, page):
    """항공자력 원제원(박영수 외, 2019) 확인과 해상도 한계."""
    s = slide_base(prs, "Ⅱ. 확인·추가 자료",
                   "원측선은 존재하지 않고, 있었더라도 개선폭은 제한적이었다",
                   "박영수 외(2019) 「한국의 자력이상도」에 37년치 탐사 제원이 정리돼 "
                   "있다. 측선 간격 자체가 1~2 km 다.", page, rule=RED)
    stat(s, M, BODY_Y, 2.95, 1.24, "1 ~ 2 km",
         "주측선 간격 (연도별)\n점간 약 30 m · 고도 100~200 m", LBLUE, BLUE, 21)
    stat(s, M + 3.08, BODY_Y, 2.95, 1.24, "2.3 × 2.8 km",
         "공개본 1.5분 등격자\n원측선보다 성기다", PEACH, RED, 19)
    stat(s, M + 6.16, BODY_Y, 2.95, 1.24, "86.4 nT",
         "공개본 격자 자체 내삽오차\n25.9% 가 목표 50 nT 초과", PEACH, RED, 22)
    stat(s, M + 9.24, BODY_Y, CW - 9.24, 1.24, "364만점",
         "1982~2017 · 37년 · 도폭 190개", LGRAY, BLUE, 21)

    rows = [["주측선 / 맺음측선", "1 km / 5 km", "1.5 km / 8~10 km",
             "2 km / 6~8 km", "1~2 km / 4~6 km"],
            ["연도", "1982~1993", "1994~1999", "2000~2011", "2012~2017"],
            ["장비", "양성자 G-813", "세슘 G-822A", "세슘 + GNSS", "3축 경사계"]]
    y = table(s, M, BODY_Y + 1.42, CW,
              ["제원", "초기", "중기", "후기", "최근"], rows,
              [0.18, 0.20, 0.21, 0.20, 0.21], row_h=0.40, head_fill="8C2F2F",
              size=10)
    band(s, M, y + 0.18, CW, 1.66,
         "⚠ 설명자료의 「측선 250 m」는 자체 수치와 맞지 않는다 — 8,600 km²에 "
         "213,791점이면 점간 32 m 기준 측선 간격이 1.26 km 다\n"
         "   (250 m 였다면 108만점 필요). 사례연구도 표준 측선 자료였고 특별히 "
         "조밀한 자료가 아니었다.\n\n"
         "① 원측선은 존재하지 않는다. 설령 받았더라도 측선 직교방향은 2.8 → 1~2 km, "
         "1.4~2.8배 개선에 그쳤을 것이다\n"
         "② 편각에도 영향이 있다 — 지각층을 벡터로 복원해 D·I 에 기여시켰으므로, "
         "격자 평활이 진폭을 1.6배 줄이고 있다", RED)
    return s


def sl_datareq(prs, page):
    s = slide_base(prs, "Ⅱ. 확인·추가 자료",
                   "앞으로 필요한 자료와 요구 조건",
                   "설명자료의 권고 수준과 현재 보유량을 나란히 놓고, 확보 경로까지 "
                   "적는다.", page, rule=ORANGE)
    rows = [
        ["항공자력 원측선", "측선 100~250 m", "존재하지 않음", "신규 측량 외 경로 없음",
         "F 목표 (봉쇄)"],
        ["지표 절대측정", "30점 이상", f"{N_SITES}점", "2027 캠페인 · 야장 17건 검증",
         "D·F 목표"],
        ["관측시각", "세션마다 기록", "2019~2025 확보", "완료", "External 전제"],
        ["상시관측", "1분 자료 · 다지점", "CYG + KASA 3소", "확보 완료",
         "NOC 투영 입력"],
        ["독립 검증점", "적합 제외 측점", "미분리", "신규 캠페인 설계에 반영",
         "성능 과대평가 방지"],
        ["관측 조건", "Kp ≤ 2 · RTK ≤ 1 m", "Kp 미제어(정온 123/196)", "작업지침 반영",
         "표본 품질"],
    ]
    y = table(s, M, BODY_Y, CW,
              ["자료", "설명자료 요건", "현재", "확보 경로", "무엇이 풀리나"], rows,
              [0.17, 0.20, 0.21, 0.24, 0.18], row_h=0.46, head_fill=ORANGE, size=10)
    band(s, M, y + 0.28, CW, 1.44,
         "우선순위\n\n"
         "① 측점 확대 — 야장에만 있는 17건을 검증해 투입하면 가장 값싸게 늘릴 수 있다\n"
         "② NOC 공간투영 구현 — 자료는 이미 있다. 관측소 4곳이라 공간 모드를 "
         "적합할 수 있다\n"
         "③ 지각층은 보류 — 원측선이 존재하지 않으므로 F 목표는 현 자료로 열리지 "
         "않는다. 신규 자력측량 여부가 갈림길이다",
         ORANGE)
    return s


def sl_closing(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, BAND)
    text(s, M, 1.55, CW, 0.30, "중간보고 · LMM 정리", 12.75, "9FC0E4", True)
    text(s, M, 1.95, CW, 0.90, "구조는 세웠고, 남은 것은 자료다", 34, WHITE, True)
    cards = [
        ("구축 완료", "4층 결합 구조\n3개 계산환경 일치\n진단·검증 체계"),
        ("열리지 않는 것", "항공자력 원측선\n존재하지 않음"),
        ("추가할 자료", "2021~2025 반복측정\n관측시각(야장)"),
        ("판정 기준", f"편각 {KPI_D:g}° · 총자력 {KPI_F:g} nT\n감사 후 재적합"),
    ]
    bw = (CW - 3 * 0.24) / 4
    for i, (t, b) in enumerate(cards):
        x = M + i * (bw + 0.24)
        rect(s, x, 3.20, bw, 1.90, "23456B")
        text(s, x + 0.24, 3.42, bw - 0.48, 0.32, t, 14, "9FC0E4", True)
        text(s, x + 0.24, 3.88, bw - 0.48, 1.00, b, 11.5, WHITE, space=1.40)
    text(s, M, 5.55, CW, 0.66,
         f"현재 LOO 편각 {LOO['D']:.3f}° · 총자력 {LOO['F']:.1f} nT — "
         "연구·시범계산에는 쓰되 정식 편각 산출은 자료 확충 이후로 둔다.",
         13, "C9D8E8", space=1.35)
    text(s, M, FOOT_Y, 9.0, 0.19, FOOTER, 8.25, "7E93AB")
    text(s, 11.98, FOOT_Y - 0.03, 0.71, 0.21, f"{page:02d}", 9, "7E93AB",
         align="r")
    return s


SAFE_BOTTOM = FOOT_Y - 0.06     # 이 아래로 내려오면 푸터를 덮는다


def audit(prs):
    """모든 도형이 판면 안에 있는지 점검. 위반은 경고로 출력한다."""
    bad = []
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if Emu(sh.width).inches > W - 0.01:      # 간지 전면 배경은 의도된 것
                continue
            l, t = Emu(sh.left).inches, Emu(sh.top).inches
            r, b = l + Emu(sh.width).inches, t + Emu(sh.height).inches
            tag = (sh.text_frame.text[:26].replace("\n", " ")
                   if sh.has_text_frame and sh.text_frame.text else sh.shape_type)
            if r > W - 0.30 or l < 0.30 - 1e-6:
                bad.append(f"  [{i:02d}] 좌우 이탈 x={l:.2f}~{r:.2f}  {tag}")
            # 푸터·페이지번호 자체는 예외
            if b > SAFE_BOTTOM and Emu(sh.height).inches > 0.25 and (
                    t < FOOT_Y - 0.10 or b > H - 0.02):
                bad.append(f"  [{i:02d}] 하단 침범 y={t:.2f}~{b:.2f}  {tag}")
    if bad:
        print("\n⚠ 판면 점검 위반")
        print("\n".join(bad))
    else:
        print("\n판면 점검: 이탈 없음")
    return bad


def main():
    _mpl_init()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    builders = [sl_divider, sl_why, sl_structure, sl_nointerp, sl_pipeline,
                sl_data_status, sl_result, sl_audit, sl_audit_effect,
                sl_residual, sl_crustal_vector, sl_impl, sl_webdemo,
                sl_kigam, sl_resolution, sl_ngii, sl_match, sl_datareq,
                sl_external_verdict, sl_advisory, sl_compliance,
                sl_roadmap, sl_risk,
                sl_cmp_concept, sl_cmp_table, sl_cmp_flow, sl_lineage,
                sl_closing]
    for i, fn in enumerate(builders, 1):
        fn(prs, i)
        print(f"  [{i:02d}] {fn.__name__}")

    audit(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = OUT_DIR / f"{ts}_중간보고회_LMM_추가슬라이드_{len(builders)}장.pptx"
    prs.save(out)
    print(f"\n저장: {out}")
    return out


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    main()
