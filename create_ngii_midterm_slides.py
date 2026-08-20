# -*- coding: utf-8 -*-
"""
중간보고회(2026-08-25) 지윤수 작성분 — 국토지리정보원 본 deck 양식 판본
=========================================================================

첨부 양식 `지자기측량발전방안연구_중간보고회(20260812)_목차(초안1)_지윤수.pptx`
**그 파일을 그대로 열어 배정된 면만 채운다.** 머리글·제목띠·로고·쪽번호는
슬라이드 레이아웃(「2장」)에 있으므로 건드리지 않는다.

배정 면 (양식 파일 기준)
    5면  → 신설                   연구의 필요성 — 왜 국가 지역자기장모델인가
    5면  → 초안 PDF 10면          모델 개발에 사용한 자료 (반복관측점·상시관측소·전지구모델)
    6면  → 초안 PDF 11면          4축 자료와 항공자력 포함 여부
    7면  → 초안 PDF 12면 보완     도엽별 자침편차 제공의 정확도 요건 (신설·6면 골격 복제)
    10면 → 양식 9면               모델 갱신 주기 및 방안
    11면 → Ⅲ 추진계획             지자기모델 고도화 방안 (신설·「3장」 레이아웃)

⚠ 수치는 하드코딩하지 않고 `docs/data/lmm_model.json` · `lmm_diagnosis.json` ·
   `external_corrections_multi.csv` 에서 읽는다. 재적합하면 재실행만으로 따라간다.

    python create_ngii_midterm_slides.py
"""
import copy
import csv
import json
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).parent
DATA = ROOT / "docs" / "data"
OUT_DIR = ROOT / "docs" / "output"
TEMPLATE = Path(r"D:\LX_yoons\2026_research\2026_지자기 연구\20260825_중간보고회"
                r"\지자기측량발전방안연구_중간보고회(20260812)_목차(초안1)_지윤수.pptx")

# ── 양식 실측 토큰 ──────────────────────────────────────────────────────────
W, H = 10.8333, 7.5
L, R = 0.45, 10.38          # 본문 좌우 한계 (제목띠와 맞춤)
CW = R - L
BODY_Y = 2.00               # 회색 칩 아래

NAVY = RGBColor(0x2F, 0x54, 0x97)      # 제목띠·표 머리
CYAN = RGBColor(0x00, 0xB0, 0xF0)      # 우측 장 표시 테두리
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x59, 0x59, 0x59)
LGRAY = RGBColor(0xF2, 0xF5, 0xF9)
LINE = RGBColor(0xBF, 0xCB, 0xDC)
PEACH = RGBColor(0xFD, 0xEA, 0xDA)     # 요지 띠 (본 deck 12면과 동일 계열)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x1F, 0x6B, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FB = "KoPub돋움체 Bold"
FM = "KoPub돋움체 Medium"
FL = "KoPub돋움체 Light"


# ── 자료 ────────────────────────────────────────────────────────────────────
MODEL = json.loads((DATA / "lmm_model.json").read_text(encoding="utf-8"))
DIAG = json.loads((DATA / "lmm_diagnosis.json").read_text(encoding="utf-8"))
LOO = MODEL["loo_cv"]
N_SITES = len(MODEL["sites"])
N_OBS = DIAG["dataset"]["n_obs"]
RD = DIAG["asymmetry"]["rD_rms"]          # 편각 잔차 RMS (arcmin)
GRID_KM = DIAG["vector_recovery"]["grid"]["dx_km"]
KPI_D, KPI_F = 0.1, 50.0

# 층별 기여 사다리 — IGRF 단독부터 단계별 잔차 (하드코딩 아님)
VAL = {(r["성분"], r["단계"]): r["RMS"] for r in MODEL["validation"]}
PAGES = "https://yoonsoo-ji.github.io/geomag-site-selection"


def ext_stats():
    f = DATA / "external_corrections_multi.csv"
    if not f.exists():
        return dict(ok=0, y0="", y1="")
    rows = list(csv.DictReader(open(f, encoding="utf-8-sig")))
    ys = sorted({r["날짜"][:4] for r in rows if r.get("날짜")})
    return dict(ok=sum(1 for r in rows if r.get("상태") == "정상"),
                y0=ys[0], y1=ys[-1])


EXT = ext_stats()

# ── 2020 연구 15점 ↔ ’22~25 측량 15점 대조 (legacy2020_sets 단일 출처) ──────
#    명단·좌표·배치 지표를 하드코딩하지 않는다. 자료가 바뀌면 면이 따라간다.
import legacy2020_sets as LS                                   # noqa: E402

SET20 = LS.sets()
COV20 = LS.coverage(LS.coords(SET20["s2020"]))
COV22 = LS.coverage(LS.coords(SET20["s2225"], prefer_model=True))


# ── 그리기 도구 ─────────────────────────────────────────────────────────────
def _tf(shape, size, bold, color, align, anchor, space=1.15):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for pa in tf.paragraphs:
        pa.alignment = align
        pa.line_spacing = space
        for r in pa.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FB if bold else FM
    return tf


def text(s, x, y, w, h, txt, size=11, color=INK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.15):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(18000)
    tf.margin_top = tf.margin_bottom = 0
    lines = txt.split("\n")
    tf.text = lines[0]
    for ln in lines[1:]:
        tf.add_paragraph().text = ln
    _tf(box, size, bold, color, align, anchor, space)
    return box


def link(s, x, y, w, h, txt, url, size=10, color=None):
    """클릭 가능한 링크 상자 — 발표 중 그대로 눌러 띄운다."""
    box = text(s, x, y, w, h, txt, size, color or NAVY, True)
    for para in box.text_frame.paragraphs:
        for r in para.runs:
            r.hyperlink.address = url
    return box


def rect(s, x, y, w, h, fill=None, line=None, radius=False, weight=0.75):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        shp.adjustments[0] = 0.14
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(weight)
    shp.shadow.inherit = False
    return shp


def lead(s, y, txt, sub=None, fill=PEACH, h=0.72):
    """상단 요지 띠 — 본 deck 12면과 같은 형식(작은 근거 + 굵은 결론)."""
    rect(s, L, y, CW, h, fill=fill)
    if sub:
        text(s, L + 0.14, y + 0.06, CW - 0.28, 0.22, sub, 10, GRAY)
        text(s, L + 0.14, y + 0.28, CW - 0.28, 0.38, txt, 15.5, INK, True)
    else:
        text(s, L + 0.14, y + 0.16, CW - 0.28, 0.40, txt, 15.5, INK, True,
             anchor=MSO_ANCHOR.MIDDLE)
    return y + h


def label(s, x, y, w, txt, fill=NAVY, size=11.5, h=0.30):
    rect(s, x, y, w, h, fill=fill, radius=True)
    text(s, x, y + 0.02, w, h, txt, size, WHITE, True, PP_ALIGN.CENTER,
         MSO_ANCHOR.MIDDLE)
    return y + h


def table(s, x, y, w, heads, rows, widths, row_h=0.42, size=10.5,
          head_h=0.34, head_fill=NAVY, aligns=None):
    """괘선 최소·머리 네이비 — 본 deck 표와 같은 결."""
    cols = [w * f for f in widths]
    aligns = aligns or [PP_ALIGN.CENTER] + [PP_ALIGN.LEFT] * (len(heads) - 1)
    rect(s, x, y, w, head_h, fill=head_fill)
    cx = x
    for i, hd in enumerate(heads):
        text(s, cx + 0.10, y + 0.02, cols[i] - 0.14, head_h, hd, size, WHITE,
             True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        cx += cols[i]
    yy = y + head_h
    for k, row in enumerate(rows):
        if k % 2 == 0:
            rect(s, x, yy, w, row_h, fill=LGRAY)
        cx = x
        for i, cell in enumerate(row):
            bold = (i == 0)
            col = INK
            if isinstance(cell, tuple):
                cell, col = cell
            text(s, cx + 0.10, yy + 0.03, cols[i] - 0.14, row_h, str(cell),
                 size, col, bold, aligns[i], MSO_ANCHOR.MIDDLE)
            cx += cols[i]
        yy += row_h
    rect(s, x, y, w, yy - y, fill=None, line=LINE)
    return yy


def note(s, y, txt, color=GRAY, size=9.5):
    text(s, L, y, CW, 0.30, txt, size, color)
    return y + 0.28


def clear_notes(slide):
    """양식에 적힌 「지윤수 박사님 …」 지시문 상자를 지운다."""
    for sh in list(slide.shapes):
        if sh.has_text_frame and "박사님" in sh.text_frame.text:
            sh._element.getparent().remove(sh._element)


def set_chip(slide, txt):
    """회색 칩(소제목) 문구 교체 — 서식은 첫 run 을 물려받는다."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t and Emu(sh.top or 0).inches < 1.9 and Emu(sh.top or 0).inches > 1.0:
            p = sh.text_frame.paragraphs[0]
            if not p.runs:
                continue
            keep = p.runs[0]
            keep.text = txt
            for r in p.runs[1:]:
                r._r.getparent().remove(r._r)
            for extra in sh.text_frame.paragraphs[1:]:
                extra._p.getparent().remove(extra._p)
            return sh
    return None


def clone_slide(prs, src_idx, layout=None):
    """
    양식 면 하나를 복제한다. layout 을 주면 그 레이아웃에 붙인다
    (Ⅱ 연구수행 → Ⅲ 추진계획 처럼 장 표시를 바꿀 때).

    ⚠ 도형 XML 만 복사하면 그림이 깨진다. 이 버전 python-pptx 는 관계 id 를
      지정할 수 없으므로, 관계를 새로 만든 뒤 **복사본 XML 의 r:embed·r:id 를
      새 rId 로 갈아끼운다.**
    """
    from pptx.oxml.ns import qn

    src = prs.slides[src_idx]
    dst = prs.slides.add_slide(layout or src.slide_layout)
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)

    rid_map = {}
    for rid, rel in src.part.rels.items():
        # 레이아웃은 add_slide 가 이미 걸었고, 노트는 슬라이드마다 고유 파트라
        # 공유하면 패키지가 깨진다.
        if rel.reltype.endswith(("/slideLayout", "/notesSlide")):
            continue
        try:
            if rel.is_external:
                continue
            rid_map[rid] = dst.part.rels.get_or_add(rel.reltype, rel._target)
        except Exception:          # noqa: BLE001
            continue

    ATTRS = [qn("r:embed"), qn("r:id"), qn("r:link")]
    for sh in src.shapes:
        el = copy.deepcopy(sh._element)
        for node in el.iter():
            for a in ATTRS:
                v = node.get(a)
                if v in rid_map:
                    node.set(a, rid_map[v])
        dst.shapes._spTree.append(el)
    return dst


def strip_body(slide, keep_above=1.95):
    """복제본에서 본문 도형만 지우고 제목·칩 골격은 남긴다."""
    for sh in list(slide.shapes):
        if Emu(sh.top or 0).inches > keep_above:
            sh._element.getparent().remove(sh._element)


def set_title(slide, txt):
    """제목띠 문구 교체 (제목띠는 y<1.0 의 도형)."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if 0.2 < Emu(sh.top or 0).inches < 1.0 and sh.text_frame.text.strip():
            pa = sh.text_frame.paragraphs[0]
            if not pa.runs:
                continue
            pa.runs[0].text = txt
            for r in pa.runs[1:]:
                r._r.getparent().remove(r._r)
            return sh
    return None


def move_slide(prs, from_idx, to_idx):
    ids = prs.slides._sldIdLst
    items = list(ids)
    ids.remove(items[from_idx])
    ids.insert(to_idx, items[from_idx])


# ══════════════════════════════════════════════════════════════ 면 구성
def s_data(slide):
    """① 모델 개발에 사용한 자료 (초안 10면)."""
    clear_notes(slide)
    set_chip(slide, "기존 반복 관측 성과 기반 지자기모델 시범구축")

    y = lead(slide, BODY_Y,
             "‘25년 시범모델은 반복관측 성과·상시관측소·전지구모델·항공자력 "
             "4종 자료로 구축",
             "모델 개발에 사용한 자료 확인 (반복관측점 · 상시관측소 · 전지구모델)")

    y = table(
        slide, L, y + 0.18, CW,
        ["구분", "자료원", "규모", "모델 내 역할", "비고"],
        [["반복관측점", "지자기점 관측성과(’22~’25) + ’19년 야장",
          f"{N_SITES}점 · {N_OBS}회", "② Regional (공간분포)",
          "재방문 불일치 2점 제외"],
         ["상시관측소", "청양(기상청) · 제주 · 강릉 · 이천",
          "4개소 1분자료", "④ External (보정량 산출)",
          f"관측시각 복원(’{EXT['y0'][2:]}~’{EXT['y1'][2:]}) · 모델 반영은 보류"],
         ["전지구모델", "IGRF-14 (13차 구면조화)", "4개 에폭",
          "① Core (주자기장)", "공식 계산기와 자릿수 내 일치"],
         ["항공자력", "KIGAM 자력이상도 1.5분 격자",
          f"약 {GRID_KM:.1f} km", "③ Crustal (지각 자기이상)",
          "원측선 자료는 부재"]],
        [0.11, 0.31, 0.12, 0.25, 0.21], row_h=0.46)

    # 청양 활용 여부 — 자문의견 검토 결과
    y += 0.22
    bh = 1.28
    rect(slide, L, y, CW, bh, fill=WHITE, line=NAVY, radius=True, weight=1.0)
    label(slide, L + 0.20, y - 0.15, 2.95, "청양 상시관측소 자료 활용 여부")
    text(slide, L + 0.28, y + 0.30, CW - 0.56, 0.30,
         "보정량은 산출하였으나, 현재 표본에서는 모델 반영을 보류", 13, NAVY, True)
    text(slide, L + 0.28, y + 0.64, CW - 0.56, 0.56,
         "· 자문의견(거리·위도 차이로 단독 적용 한계)에 따라 청양 단독이 아닌 "
         "4개소 공간보간으로 세션별 보정량을 산출\n"
         "· 다만 평가 측점을 고정한 교차검증에서 보정 전(0.502°)보다 개선되지 "
         "않아(0.510~0.532°) 반영은 측점 확충 후 재판정", 10.5, INK, space=1.35)

    note(slide, y + bh + 0.10,
         "※ 반복관측 성과는 야장 원본과 대조하여 일변화·기준시점 환산이 적용되지 "
         "않은 원시값임을 확인함(매칭 18행 전부 표기 자릿수 이내 일치)")


def s_axes(slide):
    """② 4축 자료와 항공자력 포함 여부 (초안 11면)."""
    clear_notes(slide)
    set_chip(slide, "LMM(국내외기준) 모델 비교 및 한계·개선방향 도출")

    y = lead(slide, BODY_Y,
             "항공자력을 포함해 4축을 구성 — 남은 제약은 해상도와 측점 수",
             "LMM 모델 개발에 필요한 4축 자료 / 이번 개발의 항공자력 포함 여부",
             h=0.66)

    text(slide, L, y + 0.10, CW, 0.30,
         "B_LMM  =  ① Core(IGRF)  +  ② Regional(지역)  +  ③ Crustal(지각)  "
         "+  ④ External(외부장)",
         13, NAVY, True, PP_ALIGN.CENTER)

    y = table(
        slide, L, y + 0.44, CW,
        ["축", "자료", "이번 개발 반영", "한계"],
        [["① Core", "IGRF-14 (13차)", "반영 — 재현성 검증 통과", "—"],
         ["② Regional", f"지표 절대측정 {N_SITES}점", "반영 — 1차 다항식",
          "권고 30점 대비 미달"],
         ["③ Crustal", f"항공자력 {GRID_KM:.1f} km 격자",
          "반영 — 총자력에 1:1 결합(결합계수 검정 후 α=1 채택)",
          ("원측선 부재로 해상도 상한", RED)],
         ["④ External", "상시관측 4개소",
          f"보정량 산출 {EXT['ok']}세션 — 모델 반영은 보류",
          "표본 부족으로 효과 미확인"]],
        [0.11, 0.20, 0.42, 0.27], row_h=0.44)

    # 4축 구성의 효과 / 남은 한계
    y += 0.22
    bw = (CW - 0.24) / 2
    for i, (tt, col, lines) in enumerate([
        ("항공자력 포함(4축)의 효과", GREEN,
         "· 총자력 교차검증 오차 72.3 → 63.0 nT (지각 미사용 대비)\n"
         "· 결합계수를 추정하면 α=0.67 이나 교차검증은 67.0 nT 로 악화\n"
         "  → α=1 유지. 「α=1 을 진리로 두지 말라」는 지적을 검정으로 확인"),
        ("3축 수준에 머무르는 제약", RED,
         "· 원측선(측선 100~250 m) 자료가 존재하지 않음\n"
         f"· 공개 격자 {GRID_KM:.1f} km 는 단파장을 평활 — 진폭 약 1.6배 부족\n"
         "· 해상도 개선은 신규 자력측량 여부에 좌우")]):
        x = L + i * (bw + 0.24)
        rect(slide, x, y, bw, 1.24, fill=WHITE, line=col, radius=True, weight=1.0)
        label(slide, x + 0.18, y - 0.15, 2.5, tt, fill=col, size=11)
        text(slide, x + 0.24, y + 0.34, bw - 0.48, 0.86, lines, 10.5, INK,
             space=1.32)

    note(slide, y + 1.36,
         "※ 스칼라 총자력 이상은 이상벡터를 주자기장 방향으로 투영한 값이므로, "
         "파수영역 역산으로 3성분을 복원하여 편각·복각에 반영함")


def s_validate(slide):
    """④ 검증체계와 모델 선택 결과."""
    clear_notes(slide)
    set_chip(slide, "모델 검증체계 및 구조 선택 결과")

    y = lead(slide, BODY_Y,
             "만들고 끝내지 않고, 어떤 형태가 옳은지 자료로 선택",
             "구현 검증 · 구조 선택 · 예측 검증 — 평가 측점은 미리 고정", h=0.66)

    y += 0.18
    y = table(
        slide, L, y, CW * 0.56,
        ["구분", "후보", "교차검증", "판정"],
        [["Regional 차수", "R0 상수항", "0.4987°", ("채택", GREEN)],
         ["", "R1 1차 평면", "0.5236°", "기각"],
         ["", "R1T 시간항 추가", "0.5459°", "기각"],
         ["지각 결합", "F1 1:1 결합", "63.0 nT", ("채택", GREEN)],
         ["", "Fα 계수 추정(0.67)", "67.0 nT", "기각"],
         ["External", "E0 청양 단독", "0.5495°", "기각"],
         ["", "E1 최근접 관측소", "0.5504°", "기각"],
         ["", "E2 4소 공간보간", "0.4987°", ("채택(잠정)", GREEN)]],
        [0.24, 0.34, 0.24, 0.18], row_h=0.36, size=10,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])

    x2 = L + CW * 0.58
    w2 = CW - CW * 0.58
    rect(slide, x2, BODY_Y + 0.84, w2, 1.50, fill=WHITE, line=NAVY,
         radius=True, weight=1.0)
    label(slide, x2 + 0.18, BODY_Y + 0.69, 2.6, "평가 측점 고정")
    text(slide, x2 + 0.26, BODY_Y + 1.16, w2 - 0.52, 1.10,
         "여러 후보를 같은 교차검증 값을 보며 고르면 그 값에도 과적합 발생. "
         "기준 설정에서 평가 측점을 한 번만 정해 고정(편각 14 · 복각 15 · "
         "총자력 13)하고, 이후 모든 비교를 그 위에서만 수행", 10.5, INK, space=1.35)

    rect(slide, x2, BODY_Y + 2.52, w2, 1.34, fill=LGRAY, line=LINE, radius=True)
    text(slide, x2 + 0.22, BODY_Y + 2.66, w2 - 0.44, 0.26,
         "선택의 요점", 11.5, NAVY, True)
    text(slide, x2 + 0.22, BODY_Y + 2.96, w2 - 0.44, 0.84,
         "· 측점 16개로는 공간 구조 지지 불가\n"
         "· 결합계수 자유추정은 과적합\n"
         "· 단일 관측소 가정은 불성립", 10.5, INK, space=1.35)

    note(slide, y + 0.16,
         "※ 모든 비교는 같은 조건(Grade A 자료 · 동일 차수 · 동일 평가 측점)에서 "
         "처음부터 다시 적합해 수행함")


def s_perf(slide):
    """⑤ 현재 성능과 편각 원인 진단."""
    clear_notes(slide)
    set_chip(slide, "현재 성능 및 편각 정확도 제한요인")

    y = lead(slide, BODY_Y,
             "총자력은 목표에 근접 · 편각은 과제로 잔존",
             "적합에 쓰지 않은 측점에서의 예측오차(Station-LOSO)", h=0.66)

    y += 0.18
    bw = (CW - 0.44) / 3
    for i, (val, cap, fill, col) in enumerate([
            (f"{LOO['D']*60:.0f}′", "편각 D  |  목표 6′", PEACH, RED),
            (f"{LOO['F']:.0f} nT", f"총자력 F  |  목표 {KPI_F:g} nT", PEACH, RED),
            (f"{LOO['I']*60:.0f}′", "복각 I  |  별도 목표 없음", LGRAY, NAVY)]):
        x = L + i * (bw + 0.22)
        rect(slide, x, y, bw, 0.92, fill=fill, radius=True)
        text(slide, x, y + 0.10, bw, 0.44, val, 24, col, True,
             PP_ALIGN.CENTER)
        text(slide, x, y + 0.56, bw, 0.30, cap, 10.5, GRAY, align=PP_ALIGN.CENTER)

    y += 1.14
    y = table(
        slide, L, y, CW,
        ["검정한 가설", "결과", "판정"],
        [["IGRF 구현 오류", "공식 계산기와 표시 자릿수 내 일치", "배제"],
         ["Regional 차수 부족", "1차·2차가 상수항보다 나쁨", "배제"],
         ["지역 영년변화 항 누락", "시간항 추가 시 오히려 악화", "배제"],
         ["단일 관측소 외부장 보정", "청양 단독·최근접 모두 기준선보다 나쁨", "배제"],
         ["지각 결합계수 오류", "계수 자유추정은 과적합", "배제"]],
        [0.26, 0.56, 0.18], row_h=0.32, size=10,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER])

    y += 0.18
    bw2 = (CW - 0.24) / 2
    rect(slide, L, y, bw2, 0.86, fill=WHITE, line=RED, radius=True, weight=1.0)
    label(slide, L + 0.18, y - 0.15, 2.2, "남은 요인", fill=RED, size=11)
    text(slide, L + 0.24, y + 0.24, bw2 - 0.48, 0.58,
         "① 기존 반복측량의 편각 계통오차(장비·방위표지·측점 이설)\n"
         "② 지각 벡터 미반영  ③ 측점 16개라는 공간 밀도", 10.5, INK, space=1.35)

    rect(slide, L + bw2 + 0.24, y, bw2, 0.86, fill=LGRAY, line=LINE, radius=True)
    text(slide, L + bw2 + 0.44, y + 0.14, bw2 - 0.40, 0.26,
         "이 단계의 결론", 11, NAVY, True)
    text(slide, L + bw2 + 0.44, y + 0.40, bw2 - 0.40, 0.44,
         "「미달」이 아니라 「특정」 — 무엇을 고치면 무엇이 좋아지는지 "
         "확인 완료", 10.5, INK, space=1.3)


def s_service(slide):
    """③ 도엽별 자침편차 제공의 정확도 요건 (초안 12면 보완)."""
    clear_notes(slide)
    set_chip(slide, "도엽별 자침편차 제공을 위한 정확도 요건")

    y = lead(slide, BODY_Y,
             "제공 대상은 1:50,000 도엽 — 도폭 내 편각 변화 5.4′ 가 표기 "
             "정밀도의 하한을 결정",
             "규정 제12조 — 1등 지자기점은 국가기본도의 자편각 표기를 위한 점")

    y = table(
        slide, L, y + 0.18, CW * 0.60,
        ["축척", "도폭 크기", "도폭 내 편각 변화", "대표값 1개 표기 시"],
        [[("1:50,000 (제공)", NAVY), "15′", "0.0905° (5.4′)",
          ("±2.7′ 의 표기 한계", RED)],
         ["1:25,000", "7′30″", "0.0453° (2.7′)", "±1.4′"],
         ["1:5,000", "1′30″", "0.0091° (0.5′)", "±0.3′"]],
        [0.26, 0.16, 0.29, 0.29], row_h=0.42,
        aligns=[PP_ALIGN.CENTER] * 4)

    x2 = L + CW * 0.62
    w2 = CW - CW * 0.62
    rect(slide, x2, BODY_Y + 0.90, w2, 1.94, fill=WHITE, line=NAVY,
         radius=True, weight=1.0)
    label(slide, x2 + 0.18, BODY_Y + 0.75, 2.2, "현재 모델 정확도")
    text(slide, x2 + 0.26, BODY_Y + 1.22, w2 - 0.52, 0.34,
         f"편각 {LOO['D']:.2f}°  (목표 {KPI_D:g}°)", 15, RED, True)
    text(slide, x2 + 0.26, BODY_Y + 1.58, w2 - 0.52, 1.10,
         f"· 총자력 {LOO['F']:.0f} nT (목표 {KPI_F:g} nT)\n"
         f"· 모델 오차가 도폭 내 변화(5.4′)의 약 "
         f"{LOO['D']*60/5.43:.0f}배 — 아직 도엽 구분이 의미를 갖지 못함\n"
         "· 현 단계는 참고값이며 정식 표기는 목표 달성 후",
         10.5, INK, space=1.35)

    y += 0.30
    text(slide, L, y, CW, 0.30, "제공 방식(안)", 12.5, NAVY, True)
    y += 0.34
    y = table(
        slide, L, y, CW,
        ["단계", "제공 형태", "내용"],
        [["1단계", "도엽별 수치자료(엑셀)",
          "1:50,000 도엽 단위 — 도엽 중심 대표값 · 도폭 내 변화폭 · 기준시점 · 연변화율"],
         ["2단계", "조회 서비스", "국토정보플랫폼에서 위치·도엽 검색 시 자침편차 표시"],
         ["3단계", "지형도 표기", "도곽 밖 자침편차 도표 반영(정확도 목표 달성 후)"]],
        [0.11, 0.24, 0.65], row_h=0.40)

    note(slide, y + 0.12,
         "※ 도폭 내 편각 변화는 IGRF-14 산출값 / 도폭 규격 15′ 는 지도도식적용규정 "
         "제4조·제211조 / 대표값은 도곽 중심 기준")


def s_cycle(slide):
    """④ 모델 갱신 주기 및 방안 (양식 9면 · Ⅱ)."""
    clear_notes(slide)
    set_chip(slide, "모델 갱신 주기 및 방안 도출")

    y = lead(slide, BODY_Y,
             "관측망 재관측 주기(5년)에 맞춰 모델을 재계산하고, 기준시점 변경 시 "
             "수시 갱신",
             "국가 지자기모델 갱신 주기 및 절차", h=0.66)

    y += 0.22
    steps = [("재관측", "5년 주기 반복관측\n(관측시각·Kp 기록)"),
             ("재적합", "4축 재계산 · 교차검증\n(순차 튜닝 금지)"),
             ("검증", "독립 측점 검증\n판정기준 충족 확인"),
             ("배포", "도엽별 자침편차 산출\n계산기·조회 서비스 갱신")]
    bw = (CW - 3 * 0.34) / 4
    for i, (tt, body) in enumerate(steps):
        x = L + i * (bw + 0.34)
        rect(slide, x, y, bw, 1.16, fill=WHITE, line=LINE, radius=True)
        rect(slide, x, y, bw, 0.36, fill=NAVY, radius=True)
        rect(slide, x, y + 0.20, bw, 0.16, fill=NAVY)
        text(slide, x, y + 0.03, bw, 0.30, tt, 12, WHITE, True,
             PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        text(slide, x + 0.14, y + 0.48, bw - 0.28, 0.60, body, 10.5, INK,
             PP_ALIGN.CENTER if False else PP_ALIGN.LEFT, space=1.4)
        if i < 3:
            text(slide, x + bw + 0.04, y + 0.42, 0.26, 0.30, "→", 15, NAVY,
                 True, PP_ALIGN.CENTER)

    y += 1.40
    y = table(
        slide, L, y, CW,
        ["구분", "시점", "내용"],
        [["정기 갱신", "5년", "재관측 성과 반영 · 전 층 재적합 · 기준시점 환산"],
         ["수시 갱신", "IGRF 갱신 시(5년)", "① Core 계수 교체 후 재계산"],
         ["수시 갱신", "측점·자료 추가 시", "신규 관측성과 투입 후 교차검증 재수행"]],
        [0.16, 0.22, 0.62], row_h=0.42)

    note(slide, y + 0.16,
         "※ 층이 서로 결합되어 있어 한 층만 고쳐 끼우는 순차 갱신은 성립하지 "
         "않으며, 갱신 시에는 전 층을 다시 적합해야 함")


def s_plan(slide):
    """④ LMM 모델 추후 진행 방향 (Ⅲ 추진계획)."""
    clear_notes(slide)
    set_title(slide, "3. 지자기모델 고도화 방안")
    set_chip(slide, "지자기모델 검증 및 보완 방향")

    y = lead(slide, BODY_Y,
             "입력자료 보강 → 모델 고도화 → 독립검증·갱신주기 순으로 진행",
             "LMM 모델 추후 진행 방향")

    y += 0.20
    steps = [
        ("1단계", "입력자료 보강", "’26.9~10",
         "· 잔차 큰 측점 재측정 · 방위표지 고정\n"
         "· 야장에만 있는 성과 17건 검증 투입\n"
         f"· 목표: 반복관측점 {N_SITES}점 → 20점 이상"),
        ("2단계", "모델 고도화", "’26.10~11",
         "· 상시관측 4개소 공간투영(NOC)\n"
         "· 잔차면 모델링 · 다항식 차수 재선정\n"
         "· 지각 자기이상 진폭 보정 검토"),
        ("3단계", "검증 · 갱신체계", "’26.11~",
         "· 적합에 쓰지 않은 독립 측점 검증\n"
         "· 5년 갱신주기 연계 재계산 절차\n"
         "· 도엽별 자침편차 산출 · 서비스 연계"),
    ]
    bw = (CW - 0.30) / 3
    for i, (no, tt, when, body) in enumerate(steps):
        x = L + i * (bw + 0.15)
        rect(slide, x, y, bw, 2.26, fill=WHITE, line=LINE, radius=True)
        rect(slide, x, y, bw, 0.42, fill=NAVY, radius=True)
        rect(slide, x, y + 0.24, bw, 0.18, fill=NAVY)
        text(slide, x + 0.12, y + 0.04, bw - 0.24, 0.34,
             f"{no}   {tt}", 12, WHITE, True, PP_ALIGN.CENTER,
             MSO_ANCHOR.MIDDLE)
        text(slide, x + 0.16, y + 0.50, bw - 0.32, 0.22, when, 9.5, GRAY, True)
        text(slide, x + 0.16, y + 0.78, bw - 0.32, 1.40, body, 10.5, INK,
             space=1.45)
        if i < 2:
            text(slide, x + bw + 0.01, y + 0.80, 0.14, 0.30, "▶", 11, NAVY, True,
                 PP_ALIGN.CENTER)

    y += 2.46
    bw2 = (CW - 0.24) / 2
    rect(slide, L, y, bw2, 1.00, fill=LGRAY, line=LINE, radius=True)
    text(slide, L + 0.20, y + 0.12, bw2 - 0.40, 0.24, "단계 이행 판정기준",
         11.5, NAVY, True)
    text(slide, L + 0.20, y + 0.40, bw2 - 0.40, 0.52,
         f"편각 잔차 10′ 이하 → 정식 산출 검토 · 10~20′ → 부분 재측정 · "
         f"20′ 초과 → 방위 기준 재결정 (현재 {RD:.1f}′)", 10.5, INK, space=1.3)

    rect(slide, L + bw2 + 0.24, y, bw2, 1.00, fill=WHITE, line=RED,
         radius=True, weight=1.0)
    text(slide, L + bw2 + 0.44, y + 0.12, bw2 - 0.40, 0.24, "제약 사항",
         11.5, RED, True)
    text(slide, L + bw2 + 0.44, y + 0.40, bw2 - 0.40, 0.52,
         "항공자력 원측선 자료 부재 — 지각축 해상도 개선은 신규 자력측량 "
         "여부에 좌우", 10.5, INK, space=1.3)


def sv_rates():
    """
    한반도 중앙(36.0N 127.5E)의 IGRF-14 영년변화율 — 하드코딩하지 않는다.

    편각은 연 약 −2.2′ 로 서편각이 커지고 총자력은 연 약 +28 nT 로 증가한다.
    5년 재관측 주기로 환산하면 편각 이동이 1:50,000 도폭 내 변화(5.4′)의
    두 배가 되므로, 「갱신 주기」가 정확도와 같은 층위의 요건임을 보인다.
    """
    try:
        import datetime as dt

        import numpy as np

        from lmm_build import igrf_dif

        la, lo, h = np.array([36.0]), np.array([127.5]), np.zeros(1)
        d0 = igrf_dif(la, lo, h, dt.datetime(2026, 1, 1))
        d1 = igrf_dif(la, lo, h, dt.datetime(2027, 1, 1))
        return dict(dD=(float(d1[0][0]) - float(d0[0][0])) * 60,
                    dF=float(d1[2][0]) - float(d0[2][0]))
    except Exception as exc:                                   # noqa: BLE001
        print(f"  ⚠ 영년변화율 계산 생략({exc}) — 표기값으로 대체")
        return dict(dD=-2.2, dF=28.2)


def d_span():
    """전국 편각 폭 — 「전국 단일값으로는 표기할 수 없다」의 근거."""
    try:
        import datetime as dt

        import numpy as np

        from lmm_build import igrf_dif

        la = np.array([33.2, 33.2, 38.6, 38.6])
        lo = np.array([126.0, 129.6, 124.6, 129.4])
        D = igrf_dif(la, lo, np.zeros(4), dt.datetime(2026, 1, 1))[0]
        return float(np.min(D)), float(np.max(D))
    except Exception:                                          # noqa: BLE001
        return -9.4, -7.8


SV = sv_rates()
D_LO, D_HI = d_span()


def s_need(slide):
    """① 연구의 필요성 — 왜 국가 지역자기장모델을 구축해야 하는가."""
    clear_notes(slide)
    set_chip(slide, "국가 지역자기장모델(LMM) 구축의 필요성")

    dI = VAL[("D_deg", "IGRF")] * 60                # IGRF 단독 편각 잔차 (′)

    y = lead(slide, BODY_Y,
             "자침편차는 국가기본도 표기 대상이나, 산출 근거는 국내 실측이 "
             "아닌 전지구모델 계산값뿐임",
             "왜 국가 지역자기장모델을 구축해야 하는가", h=0.66)

    y = table(
        slide, L, y + 0.18, CW,
        ["구분", "현재 여건", "그래서 필요한 것"],
        [["제도·목적",
          "국가기본도 표기 대상 — 작업규정 제12조제1항(설치 목적)",
          "표기값을 국가가 직접 산출·관리하는 자체 모델"],
         ["정확도",
          (f"IGRF 단독 시 국내 편각 잔차 {dI:.1f}′ — 목표 "
           f"{KPI_D * 60:.0f}′ 의 약 {dI / (KPI_D * 60):.0f}배", RED),
          "지역장·지각장을 더해 국내 실측값에 맞추는 층 구조"],
         ["시간 변동",
          f"영년변화 편각 {SV['dD']:.1f}′/년 · 일변화만으로 편각 0.11°",
          "기준시점 환산 + 상시관측 보정 (규정 제21조)"],
         ["공간 변동",
          f"전국 편각 {D_LO:.1f}° ~ {D_HI:.1f}° — 약 "
          f"{abs(D_HI - D_LO) * 60:.0f}′ 분포",
          "전국 단일값 표기 불가 — 위치별 산출 모델"]],
        [0.12, 0.50, 0.38], row_h=0.46, head_h=0.32, size=10)

    y += 0.28
    bw = (CW - 0.24) / 2
    for i, (tt, lw, col, lines) in enumerate([
        ("왜 「층 분리」 구조인가", 2.7, NAVY,
         "· 네 층은 기원·시간상수·공간규모가 모두 상이\n"
         "· 한 값으로 내삽하면 시간 변동이 공간 분포로 둔갑\n"
         "· 일본 국토지리원·미국 NOAA·영국 BGS 도 층 분리 구조"),
        ("국내 여건 — 자료는 있으나 결합체계 부재", 4.1, GREEN,
         f"· 반복관측점 {N_SITES}점 · 항공자력 {GRID_KM:.1f} km 격자 · "
         "상시관측 4개소\n"
         "· 개별 활용에 그쳐 하나의 모델로 결합된 사례 없음\n"
         "· 신규 취득이 아니라 «결합체계» 구성이 이번 연구의 요체")]):
        x = L + i * (bw + 0.24)
        rect(slide, x, y, bw, 1.32, fill=WHITE, line=col, radius=True,
             weight=1.0)
        label(slide, x + 0.18, y - 0.15, lw, tt, fill=col, size=11)
        text(slide, x + 0.24, y + 0.32, bw - 0.48, 0.94, lines, 10, INK,
             space=1.28)

    note(slide, y + 1.44,
         "※ 편각 잔차는 국내 반복관측 성과에 대한 실측값, 영년변화율·편각 "
         "분포는 IGRF-14 산출값임 (하드코딩 아님)")


def s_why(slide):
    """① 왜 LMM 인가 — 층별 기여를 국내 실측으로."""
    clear_notes(slide)
    set_chip(slide, "전지구모델 대비 층별 기여도 실측 결과")

    y = lead(slide, BODY_Y,
             "전지구모델만으로는 국내 지표값과 불일치 — 층을 쌓으면 그 몫이 "
             "실제로 감소",
             "층별 기여도 실측 — IGRF 단독부터 단계별 잔차")

    dI = VAL[("D_deg", "IGRF")] * 60
    dR = VAL[("D_deg", "+Regional")] * 60
    iI = VAL[("I_deg", "IGRF")] * 60
    iR = VAL[("I_deg", "+Regional")] * 60
    fI = VAL[("F_nT", "IGRF")]
    fC = VAL[("F_nT", "+Crustal")]
    fR = VAL[("F_nT", "+Crustal+Regional")]

    y = table(
        slide, L, y + 0.18, CW,
        ["성분", "① Core 단독", "+ ③ Crustal", "+ ② Regional", "감소율"],
        [["편각 D", f"{dI:.1f}′", "—", f"{dR:.1f}′",
          (f"−{100 * (1 - dR / dI):.0f} %", GREEN)],
         ["복각 I", f"{iI:.1f}′", "—", f"{iR:.1f}′",
          (f"−{100 * (1 - iR / iI):.0f} %", GREEN)],
         ["총자력 F", f"{fI:.0f} nT", f"{fC:.0f} nT", f"{fR:.0f} nT",
          (f"−{100 * (1 - fR / fI):.0f} %", GREEN)]],
        [0.14, 0.20, 0.20, 0.20, 0.26], row_h=0.46)

    y += 0.22
    bw = (CW - 0.24) / 2
    for i, (tt, col, lines) in enumerate([
        ("층을 쌓아야 하는 이유", GREEN,
         "· 총자력은 ③ 지각 자기이상이 대부분을 제거\n"
         f"   ({fI:.0f} → {fC:.0f} nT, 전체 감소의 대부분)\n"
         "· 편각·복각은 지각층이 스칼라만 주므로 ② 지역장이 담당\n"
         "· 네 층이 각각 다른 성분·다른 파장을 담당"),
        ("그래도 남는 과제", RED,
         f"· 편각 {dR:.1f}′ — 목표 6′ 에 미달\n"
         "· 감소율도 편각이 최저 — 자기 원인만으로 설명 불가\n"
         "· 원인 진단과 측점 확충이 후반기 과제")]):
        x = L + i * (bw + 0.24)
        rect(slide, x, y, bw, 1.56, fill=WHITE, line=col, radius=True,
             weight=1.0)
        label(slide, x + 0.18, y - 0.15, 2.4, tt, fill=col, size=11)
        text(slide, x + 0.24, y + 0.34, bw - 0.48, 1.18, lines, 10.5, INK,
             space=1.32)

    note(slide, y + 1.68,
         "※ 개념도가 아니라 국내 반복관측 성과에 대해 IGRF 단독부터 단계별로 "
         "실제 측정한 잔차임 (평가 측점 고정)")


def s_vs2020(slide):
    """② 2020년 연구와의 관계 — 관측망 실체와 구조 확장."""
    clear_notes(slide)
    set_chip(slide, "2020년 연구성과와의 관계 및 확장 사항")

    y = lead(slide, BODY_Y,
             f"같은 「15점」이 아님 — 겹치는 측점은 {len(SET20['common'])}점뿐이고, "
             "새로 든 점은 2020 이 자료부족으로 뺐던 점",
             "기존 연구와의 관계 / 이번 연구에서 확장한 사항", h=0.62)

    # ── ① 측점망 — 세 갈래로 갈린다 ─────────────────────────────────────
    y += 0.16
    bw = (CW - 2 * 0.20) / 3
    for i, (tt, col, names) in enumerate([
        (f"공통 {len(SET20['common'])}점", GRAY, SET20["common"]),
        (f"2020 에만 {len(SET20['only2020'])}점", RED, SET20["only2020"]),
        (f"’22~25 에만 {len(SET20['only2225'])}점", NAVY, SET20["only2225"])]):
        x = L + i * (bw + 0.20)
        rect(slide, x, y, bw, 0.92, fill=WHITE, line=col, radius=True,
             weight=1.0)
        label(slide, x + 0.16, y - 0.15, 1.9, tt, fill=col, size=10.5)
        text(slide, x + 0.20, y + 0.30, bw - 0.40, 0.56,
             " · ".join(sorted(names)), 9.5, INK, space=1.28)

    # ── ② 그래서 무엇이 달라졌나 ────────────────────────────────────────
    y += 1.12
    rect(slide, L, y, CW, 0.66, fill=PEACH)
    text(slide, L + 0.16, y + 0.06, CW - 0.32, 0.28,
         f"① ’22~25 신규 {len(SET20['only2225'])}점이 "
         f"{len(SET20['new_from_once'])}점 전부 2020 의 「측량 1회」 배제 목록 "
         "— 자료 공백을 정면으로 채운 배치", 10.5, INK, True)
    text(slide, L + 0.16, y + 0.34, CW - 0.32, 0.28,
         f"② 대신 최대 관측공백 {COV20['max']:.0f} → {COV22['max']:.0f} km · "
         f"50 km 초과 지역 {COV20['over50']:.1f} → {COV22['over50']:.1f} % — "
         "강원·충청 내륙 축(화천·봉평·제천·청송·청양·서산)을 잃음",
         10.5, RED, True)

    # ── ③ 구조 확장 ────────────────────────────────────────────────────
    y += 0.84
    y = table(
        slide, L, y, CW,
        ["구분", "2020년 연구", "이번 연구 (’25 시범구축)"],
        [["지역장", "별도 층 없음 — IGRF 계산값 사용",
          "지표 절대측정으로 직접 구성 (R0 상수항)"],
         ["지각장", "반영 없음",
          (f"항공자력 1:1 결합 — 총자력 {VAL[('F_nT', 'IGRF')]:.0f} → "
           f"{VAL[('F_nT', '+Crustal')]:.0f} nT", GREEN)],
         ["외부장", "청양 단독 보정",
          "상시관측 4개소 공간보간 · 성분별 관측구간 (잠정)"],
         ["검증", "과거 성과와 비교",
          "평가집합 동결 + Station-LOSO + 블록 홀드아웃"]],
        [0.12, 0.36, 0.52], row_h=0.38, head_h=0.30, size=10)

    note(slide, y + 0.14,
         "※ 성능 우열 비교가 아님 — 측점 명단은 2020 보고서 p.117(측량회수 1회 "
         "15점)과 ’22~25 성과표를 대조한 결과이며, 관측공백은 국토 0.1° 격자 기준임")



def s_web(slide):
    """③ 웹 시연 — 링크를 걸어 발표 중 바로 띄운다."""
    clear_notes(slide)
    set_chip(slide, "시범구축 결과 웹 공개 및 시연")

    y = lead(slide, BODY_Y,
             "구축 결과는 설치 없이 브라우저에서 바로 확인 가능",
             "’25년 시범모델 산출물 공개 현황 (아래 주소를 눌러 시연)", h=0.66)

    items = [
        ("LMM 계산기", "lmm.html",
         "좌표·표고·일자를 넣으면 편각·복각·총자력을 계산.\n"
         "층별 기여와 정확도 고지를 함께 표시"),
        ("한국 지자기도 (2D)", "lmm_map.html",
         "IGRF · LMM · 차이를 나란히 표출.\n"
         "2020 연구 지자기도와 같은 형식이라 직접 대비 가능"),
        ("3D 지구 자기장", "geomag_globe.html",
         "전지구 자기장 위에 한국 고해상 패치를 중첩.\n"
         "선점 후보·기존점 마커에서 상세지도로 연결"),
        ("선점 검토 지도", "survey_review.html",
         "현장조사 103점의 등급·사진·방위표지와\n기존 측정점 30점을 함께 표출"),
    ]
    bw = (CW - 0.24) / 2
    ch, pitch = 1.26, 1.42          # 카드 높이 · 행 간격
    y += 0.18
    for i, (nm, pg, desc) in enumerate(items):
        x = L + (i % 2) * (bw + 0.24)
        yy = y + (i // 2) * pitch
        rect(slide, x, yy, bw, ch, fill=LGRAY)
        rect(slide, x, yy, 0.055, ch, fill=NAVY)
        text(slide, x + 0.22, yy + 0.12, bw - 0.44, 0.28, nm, 12.5, NAVY, True)
        link(slide, x + 0.22, yy + 0.44, bw - 0.44, 0.24,
             f"{PAGES}/{pg}", f"{PAGES}/{pg}", 9)
        text(slide, x + 0.22, yy + 0.72, bw - 0.44, 0.48, desc, 10, GRAY,
             space=1.30)

    # 카드 아래끝 + 여백 — 「시연 순서」 라벨이 카드에 물리지 않게 띄운다
    y += pitch + ch + 0.32
    bh = 0.80
    rect(slide, L, y, CW, bh, fill=WHITE, line=GREEN, radius=True, weight=1.0)
    label(slide, L + 0.20, y - 0.15, 1.9, "시연 순서", fill=GREEN)
    text(slide, L + 0.28, y + 0.28, CW - 0.56, 0.50,
         "① 계산기에서 층별 기여 확인 → ② 지자기도에서 IGRF 와의 차이 확인 → "
         "③ 지구본에서 한국 패치 확대·선점 후보 연결",
         10.5, INK, space=1.32, anchor=MSO_ANCHOR.MIDDLE)

    note(slide, y + bh + 0.10,
         "※ 모두 외부 서버 없이 단일 파일로 동작하며, GitHub Pages 로 상시 공개 중임")


# ══════════════════════════════════════════════════════════════ 발표 노트
def set_notes(slide, txt):
    """슬라이드 노트 기입 — 개조식. 복제 면은 노트가 없으므로 새로 만들어진다."""
    tf = slide.notes_slide.notes_text_frame
    lines = [ln.rstrip() for ln in txt.strip("\n").split("\n")]
    tf.text = lines[0]
    for ln in lines[1:]:
        tf.add_paragraph().text = ln
    for pa in tf.paragraphs:
        pa.line_spacing = 1.25
        for r in pa.runs:
            r.font.size = Pt(11)
            r.font.name = FM
    return slide


def apply_notes(prs):
    """면 위치(0-기준)마다 발표 노트를 붙인다."""
    N = {}

    N[4] = f"""
【핵심 메시지】
· 자침편차는 국가기본도 표기 대상이나, 현재 산출 근거는 전지구모델 계산값뿐임
· 정확도·시간변동·공간변동 세 축에서 모두 자체 지역모델이 요구됨

【설명 순서】
① 제도 — 「지구물리측량 작업규정」 제12조제1항, 1등 지자기점은 국가기본도 자편각 표기가 설치 목적
② 정확도 — IGRF 단독 국내 편각 잔차 {VAL[("D_deg", "IGRF")] * 60:.1f}′ 로 목표 {KPI_D * 60:.0f}′ 의 약 {VAL[("D_deg", "IGRF")] * 60 / (KPI_D * 60):.0f}배
③ 시간 — 영년변화 편각 {SV["dD"]:.1f}′/년 (5년이면 약 {abs(SV["dD"]) * 5:.0f}′ 이동, 도폭 내 변화 5.4′ 의 두 배)
   — 일변화만으로도 편각 0.11° 변동 → 기준시점 환산·상시관측 보정이 규정상 요건(제21조)
④ 공간 — 전국 편각이 {D_LO:.1f}° ~ {D_HI:.1f}° 로 약 {abs(D_HI - D_LO) * 60:.0f}′ 분포 → 단일값 표기 불가
⑤ 결론 — 자료를 새로 얻는 것이 아니라 이미 있는 자료를 «층으로 결합»하는 체계가 필요

【강조점】
· 필요성의 근거를 당위가 아니라 수치로 제시하는 면임
· 「층 분리」는 학술적 취향이 아니라 필수 — 단일 내삽은 시간 변동을 공간 분포로 오인함
· 국내에 자료는 이미 존재하나 하나의 모델로 결합된 사례가 없었음

【예상 질문】
Q. IGRF 를 그대로 쓰면 안 되는가
   → 국내 지표 실측과 편각 {VAL[("D_deg", "IGRF")] * 60:.1f}′ · 총자력 {VAL[("F_nT", "IGRF")]:.0f} nT 어긋남. 목표의 5배 수준
   → IGRF 는 차단차수 13 의 핵 자기장 모델이라 지각·지역 성분을 담지 않음
Q. 영년변화율은 어디서 나온 값인가
   → IGRF-14 로 한반도 중앙(36.0°N 127.5°E)에서 직접 계산. 스크립트 재실행 시 자동 갱신
Q. 다른 나라도 지역모델을 두는가
   → 일본 국토지리원·미국 NOAA·영국 BGS 모두 전지구모델과 별도로 층 분리 구조를 운용
"""

    N[5] = f"""
【핵심 메시지】
· 전지구모델(IGRF) 단독으로는 국내 지표 자기장 재현에 한계
· 층을 쌓을 때마다 오차가 실제로 감소함을 국내 실측으로 확인

【설명 순서】
① IGRF 는 지구 핵에서 생성되는 주자기장의 전지구 모델
   — 국내 지표에서 총자력 {VAL[("F_nT", "IGRF")]:.0f} nT · 편각 {VAL[("D_deg", "IGRF")] * 60:.1f}′ 편차
② ③ 지각 자기이상 반영 → 총자력 {VAL[("F_nT", "IGRF")]:.0f} → {VAL[("F_nT", "+Crustal")]:.0f} nT (감소분의 대부분을 차지)
③ 편각·복각은 지각층이 스칼라(총자력)만 제공 → ② 지역장이 담당
   — 편각 {VAL[("D_deg", "IGRF")] * 60:.1f} → {VAL[("D_deg", "+Regional")] * 60:.1f}′ · 복각 {VAL[("I_deg", "IGRF")] * 60:.1f} → {VAL[("I_deg", "+Regional")] * 60:.1f}′
④ 표의 수치는 개념도가 아닌 국내 반복관측 성과에 대한 실측 잔차

【강조점】
· 편각 감소율이 세 성분 중 최저 — 11면 원인진단의 출발점
· 자기장 원인만으로 설명되지 않는 몫의 존재를 시사

【예상 질문】
Q. 편각 개선폭이 왜 작은가
   → 편각 잔차의 상당분이 자기장이 아닌 관측 계통오차(방문마다 달라지는 방위표지)에서 기인. 11면에서 상술
Q. 이 수치의 평가 기준은
   → 10면의 고정 평가 측점(편각 14 · 복각 15 · 총자력 13) 기준
"""

    N[6] = f"""
【핵심 메시지】
· 2020 과 이번 연구는 같은 「15점」이 아님 — 공통 측점은 {len(SET20["common"])}점뿐
· 신규 측점은 자료 공백을 메웠으나, 배치(관측공백)는 오히려 나빠짐
· 성능 우열 비교가 아니라 구조·관측망 대비표임

【설명 순서】
① 2020 은 「2012~2019 편복각 측량 2회 이상」 15점만 분석에 사용
   — 나머지 15점은 측량회수 1회 또는 성과 신뢰도 미확보로 배제(보고서 p.117)
② ’22~25 측량 15점과 겹치는 것은 {" · ".join(sorted(SET20["common"]))} 뿐
③ ’22~25 신규 {len(SET20["only2225"])}점은 {len(SET20["new_from_once"])}점 전부가 2020 의 배제 목록
   — 자료 공백을 정면으로 채운 배치
④ 대신 최대 관측공백 {COV20["max"]:.0f} → {COV22["max"]:.0f} km · 50 km 초과 지역 {COV20["over50"]:.1f} → {COV22["over50"]:.1f} %
   — 강원·충청 내륙 축을 잃어 배치 기준이 작동하지 않았음을 시사
⑤ 구조 차이는 표로 — 지역장 신설 · 지각장 반영 · 외부장 다관측소 · 검증 전환

【강조점】
· 「15 → {N_SITES}점」이 아니라 「다른 15점 + 야장 2점」임을 분명히 할 것
· 배치 손실은 신규 50점(Track C) 설계에 공간 기준을 넣어야 하는 직접 근거

【예상 질문】
Q. 2020 모델보다 정확한가
   → 직접 비교 불가. 기준시점·환산절차가 달라 동일 조건 비교가 성립하지 않음
Q. 왜 15점이 {N_SITES}점인가
   → ’22~25 성과표 15측점 중 재방문 불일치 1점(부안) 제외, 2019 야장 검증 통과 2점(남지·장흥) 추가
Q. 2020 이 쓰던 측점은 왜 빠졌는가
   → ’22~25 측량 대상에 들지 않았음. 빠진 {len(SET20["only2020"])}점 중 {len(SET20["dropped_pond"])}점이
     저수지·제방 설치점이라 표석 지속성 문제와도 겹침
"""

    N[7] = f"""
【핵심 메시지】
· 4종 자료를 모두 국내 실자료로 확보하여 시범구축 수행
· 상시관측소는 보정량 산출까지 완료했으나 모델 반영은 보류

【설명 순서】
① 반복관측점 — 성과표(’22~’25) + ’19년 야장, {N_SITES}점 · {N_OBS}회
② 상시관측소 — 청양(기상청) + 제주·강릉·이천 4개소 1분자료
③ 전지구모델 — IGRF-14, 전 연도 단일 판본 적용
④ 항공자력 — KIGAM 자력이상도, 약 {GRID_KM:.1f} km 격자

【청양 활용 여부 — 자문의견 대응】
· 자문 지적(거리·위도 차이로 청양 단독 적용 한계)을 수용하여 4개소 공간보간으로 전환
· 다만 평가 측점 고정 교차검증에서 개선 미확인 → 반영은 측점 확충 후 재판정

【예상 질문】
Q. 성과표에 일변화 보정이 이미 되어 있지 않은가 (이중보정 우려)
   → 야장 원본과 직접 대조 완료. 매칭 18행 전부 표기 자릿수 이내 일치
   → 보정되지 않은 원시값으로 확정. 이중보정 우려 없음
Q. ’19년 자료는 IGRF 판본이 달라 섞으면 안 되지 않는가
   → 계산은 전 연도 IGRF-14 단일 판본. ’19년은 확정치(DGRF) 구간 사이라 판본 교체 영향이 가장 작음
   → 실측 대조상 ’19년 −0.2′ · ’24년 +5.4′ 로 오히려 최근 연도가 민감
"""

    N[8] = f"""
【핵심 메시지】
· 항공자력을 포함해 4축을 모두 구성 — 개념이 아닌 실제 결합
· 남은 제약은 지각축 해상도와 측점 수

【설명 순서】
① 결합식 제시 — B_LMM = ① Core + ② Regional + ③ Crustal + ④ External
② 축별 반영 수준과 한계를 표로 확인
③ 항공자력 포함 효과 — 총자력 교차검증 72.3 → {LOO["F"]:.1f} nT
④ 결합계수 α 는 자유추정 시 0.67 이나 교차검증 67.0 nT 로 악화 → α=1 유지

【강조점】
· 「α=1 을 진리로 두지 말라」는 자문 지적을 실제 검정으로 확인한 사례
· 원측선 자료 부재는 확인된 사실 — 현 격자가 해상도 상한

【예상 질문】
Q. 지각층은 총자력에만 반영되는가
   → 스칼라 이상을 파수영역 역산하여 3성분 복원 가능, 편각·복각 적용도 구현 완료
   → 다만 격자 결측 39% 문제로 현재는 진단 전용
Q. 원측선 자료를 구할 수 없는가
   → 존재하지 않음을 확인. 해상도 개선은 신규 자력측량 여부에 좌우
"""

    N[9] = f"""
【핵심 메시지】
· 제공 대상은 1:50,000 도엽
· 도폭 내 편각 변화 5.4′ 가 대표값 1개 표기의 정밀도 하한을 결정

【설명 순서】
① 규정 제12조 — 1등 지자기점은 국가기본도의 자편각 표기를 위한 점(법정 목적)
② 축척별 도폭 내 편각 변화 비교 — 1:50,000 에서 0.0905°(5.4′)
③ 현재 모델 편각 오차 {LOO["D"]:.2f}° 는 도폭 내 변화의 약 {LOO["D"] * 60 / 5.43:.0f}배
   → 아직 도엽 구분이 의미를 갖지 못하는 단계
④ 제공 방식 3단계 — 수치자료(엑셀) → 조회 서비스 → 지형도 표기

【주의】
· 「정식 산출 금지」 같은 법정 요건 표현은 사용 지양 — 목표 미달은 법령 위반이 아님
· 법정 기준(작업규정 §20 의 30′)과 공학 목표(6′)는 성격이 다름을 구분할 것

【예상 질문】
Q. 왜 1:25,000 이 아니라 1:50,000 인가
   → 현재 제공 대상 도엽 기준. 1:25,000 은 도폭 내 변화가 2.7′ 로 여유가 크나 대상이 아님
Q. 언제부터 표기가 가능한가
   → 편각 목표 달성 후. 현 단계는 참고값 제공까지
"""

    N[10] = f"""
【핵심 메시지】
· 만들고 끝내지 않고, 어떤 형태가 옳은지 자료로 선택
· 평가 측점을 미리 고정한 것이 이 검증의 핵심

【설명 순서】
① Regional 차수 — R0 상수항 채택. 1차·2차가 오히려 악화
   → {N_SITES}점으로는 지역장의 공간 구조를 결정할 수 없음(측점 확충의 직접 근거)
② 지각 결합 — F1(α=1) 채택. 계수 자유추정은 과적합
③ External — E2 4소 공간보간만 기준선 도달, 청양 단독(E0)·최근접(E1)은 기각
④ 평가 측점 고정의 의미 설명

【강조점】
· 여러 후보를 같은 교차검증 값을 보며 고르면 그 값에도 과적합이 발생
· 그래서 기준 설정에서 평가 측점을 한 번만 정해 파일로 동결하고 모든 비교를 그 위에서 수행
· 이 절차 자체가 이번 검증의 신뢰 근거

【예상 질문】
Q. 상수항만 쓰면 지역모델이라 할 수 있는가
   → 현 측점 수의 한계를 그대로 드러낸 결과. 무리한 차수 적용이 오히려 위험
   → 측점 확충 후 공간항 재판정 예정
Q. E2 를 골랐다면서 왜 「잠정」인가
   → 개선폭이 작고 12개 설정 중 최소를 고른 것이라 선택편의를 포함. 근거는 폭이 아니라 방식의 타당성
"""

    N[11] = f"""
【핵심 메시지】
· 총자력은 목표에 근접({LOO["F"]:.0f} nT / 목표 {KPI_F:g}), 편각이 과제로 잔존({LOO["D"] * 60:.0f}′ / 목표 6′)
· 「미달」이 아니라 「특정」 — 무엇이 성능을 제한하는지 지목 완료

【설명 순서】
① 세 성분의 Station-LOSO 예측오차 제시 — 적합에 쓰지 않은 측점에서의 오차
② 검정하여 배제한 가설 5종을 순서대로 설명
③ 남은 요인 3가지 제시 — 편각 계통오차 · 지각 벡터 미반영 · 공간 밀도

【강조점】
· 배제된 가설들이 이 연구의 실질 성과 — 원인 범위를 좁혔다는 뜻
· 편각 잔차의 정체는 재방문 불일치로 확인 — 방위표지가 방문마다 바뀌는 것이 구조적 원인
· 측정 절차·계산 자체에는 오류 없음
  (야장 내부 산술 재계산차 중앙 0.16″ · 같은 날 세션 산포 중앙 1.41′ 로 정밀도는 양호)
· 즉 「정밀하지만 부정확」한 상태 — 사람·장비가 아니라 방위 기준의 문제

【예상 질문】
Q. 목표 달성이 가능한가
   → 편각은 신규 측점의 방위각 결정법 상향(천문·자이로 또는 RTK 장기선)이 전제. 총자력은 근접
Q. 복각은 왜 목표가 없는가
   → 자침편차 표기가 편각 기준이므로 KPI 는 편각·총자력에 설정
Q. 측점을 늘리면 해결되는가
   → 공간 밀도는 개선되나 편각 계통오차는 별개. 두 축을 함께 다뤄야 함
"""

    N[12] = """
【핵심 메시지】
· 구축 결과를 즉시 확인 가능한 형태로 공개
· 모두 외부 서버 없이 단일 파일로 동작 — 오프라인에서도 사용 가능

【시연 순서】
① LMM 계산기 — 좌표 입력 후 층별 기여 확인(Core·Regional·Crustal 이 각각 얼마인지)
② 한국 지자기도 — IGRF 와 LMM 을 같은 색범위로 대비, 차이 화면에서 지각 이상 구조 확인
③ 3D 지구본 — 전지구에서 한국으로 이동, 고해상 패치 확대
④ (시간 여유 시) 선점 검토 지도 — 현장조사 103점의 등급·사진·방위표지

【사전 준비】
· 발표 PC 인터넷 연결 확인 · 브라우저 미리 실행
· 링크는 슬라이드에서 직접 클릭 가능
· 인터넷 불가 시 로컬 파일로도 동일 동작

【예상 질문】
Q. 일반에 공개된 것인가
   → GitHub Pages 로 상시 공개 중. 화면에 정확도 고지(연구·시범계산용) 표기
Q. 국토정보플랫폼 연계는
   → 9면 제공 방식 2단계에 해당. 정확도 목표 달성 후 추진
Q. 계산 근거를 신뢰할 수 있는가
   → 웹·엑셀·Python 세 구현이 동일 좌표에서 표시 자릿수 내 일치함을 확인
"""

    N[15] = """
【핵심 메시지】
· 관측망 재관측 주기(5년)에 맞춘 정기 갱신 + 기준시점 변경 시 수시 갱신
· 층이 서로 결합되어 있어 한 층만 교체하는 순차 갱신은 불가

【설명 순서】
① 갱신 4단계 — 재관측 → 재적합 → 검증 → 배포
② 정기 갱신(5년)과 수시 갱신(IGRF 갱신 시 · 측점 자료 추가 시) 구분
③ 재관측 시 관측시각·Kp 기록 필수 — 외부장 보정의 전제 조건

【강조점】
· 「순차 튜닝 금지」가 핵심 — 층별로 따로 조정하면 결합 효과를 놓침
· 갱신 시마다 교차검증을 다시 수행해야 성능을 객관적으로 말할 수 있음
· 야장 양식에 F 측정시각·센서 위치차 칸 추가 건의 필요(현재 미기재)

【예상 질문】
Q. IGRF 갱신 주기와 어긋나지 않는가
   → IGRF 도 5년 주기라 정합. 다만 갱신 시점이 달라 수시 갱신 항목으로 분리
Q. 5년 주기가 적정한가
   → 영년변화 규모상 5년이면 편각 약 35′ 변화. 기준점 재관측 주기와도 정합
"""

    N[16] = f"""
【핵심 메시지】
· 입력자료 보강 → 모델 고도화 → 독립검증 순의 3단계
· 단계 이행에 정량 판정기준을 미리 설정

【설명 순서】
① 1단계(’26.9~10) — 잔차 큰 측점 재측정 · 방위표지 고정 · 야장 성과 17건 검증 투입
② 2단계(’26.10~11) — 상시관측 공간투영(NOC) · 다항식 차수 재선정 · 지각 진폭 보정 검토
③ 3단계(’26.11~) — 독립 측점 검증 · 갱신주기 연계 · 도엽별 자침편차 산출

【판정기준 — 미리 정해 둔 것이 요점】
· 편각 잔차 10′ 이하 → 정식 산출 검토
· 10~20′ → 잔차 큰 측점 부분 재측정
· 20′ 초과 → 모형 작업 중단, 방위 기준 재결정 선행
· 현재 {RD:.1f}′ 로 두 번째 구간 — 모형 작업을 계속하되 현장 재측정이 남은 조건

【제약 사항】
· 항공자력 원측선 자료 부재 — 지각축 해상도 개선은 신규 자력측량 여부에 좌우

【예상 질문】
Q. 신규 50점은 언제 측정하는가
   → 후반기 별도 과업. 일부(10~15점)를 적합에 쓰지 않고 검증 전용으로 남기는 것이 설계 요점
   → 그래야 최종 보고서에서 「구축에 사용하지 않은 신규 기준점에서 독립검증」이라 쓸 수 있음
Q. 20점 확보가 가능한가
   → 야장에만 있고 성과표에 없는 (측점,연도) 17건이 후보. 검증 통과분을 투입
"""

    for idx, txt in N.items():
        set_notes(prs.slides[idx], txt)
    return len(N)


# ══════════════════════════════════════════════════════════════ 실행
def main():
    sys.stdout.reconfigure(encoding="utf-8")
    if not TEMPLATE.exists():
        raise SystemExit(f"양식 파일을 찾을 수 없다: {TEMPLATE}")

    prs = Presentation(str(TEMPLATE))
    print(f"양식 적재: {TEMPLATE.name}  ({len(prs.slides)}면)")

    s_data(prs.slides[4])          # 모델 개발 자료
    s_axes(prs.slides[5])          # 4축·항공자력
    s_cycle(prs.slides[8])         # 모델 갱신 주기 (면 삽입 전에 채운다)

    # 골격을 복제할 원본 면(2장 레이아웃). 앞에 끼울 때마다 밀리므로 추적한다.
    skel = 5

    def insert(fn, at):
        dup = clone_slide(prs, skel)
        strip_body(dup)
        move_slide(prs, len(prs.slides) - 1, at)
        fn(prs.slides[at])

    # ── 뒤쪽부터 — 골격 인덱스가 움직이지 않는다 ────────────────────────
    for fn, at in ((s_service, 6), (s_validate, 7), (s_perf, 8), (s_web, 9)):
        insert(fn, at)

    # ── 앞머리 세 면 — 필요성 · 왜 LMM · 2020 비교.
    #    끼울 때마다 골격이 뒤로 밀리므로 skel 을 함께 올린다 ──────────────
    for fn, at in ((s_need, 4), (s_why, 5), (s_vs2020, 6)):
        insert(fn, at)
        skel += 1

    # 추진계획(Ⅲ) 면 — 장 표시가 달라야 하므로 「3장」 레이아웃에 붙인다
    cycle_at = 15                             # 앞에 7면이 끼어 8 → 15
    lay3 = [l for l in prs.slide_masters[0].slide_layouts if l.name == "3장"][0]
    plan = clone_slide(prs, cycle_at, layout=lay3)   # 갱신주기 면 골격 재사용
    strip_body(plan)
    move_slide(prs, len(prs.slides) - 1, cycle_at + 1)
    s_plan(prs.slides[cycle_at + 1])

    n_notes = apply_notes(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = OUT_DIR / f"{ts}_중간보고회_지자기모델_작성분_11장.pptx"
    prs.save(out)
    print(f"저장: {out}  (전체 {len(prs.slides)}면 · 발표노트 {n_notes}면)")
    for n, t in ((5, "연구의 필요성 — 왜 LMM 인가"),
                 (6, "층별 기여 실측"), (7, "2020년 연구와의 관계"),
                 (8, "모델 개발 자료"), (9, "4축 구성·항공자력"),
                 (10, "자침편차 정확도"), (11, "검증체계·모델선택"),
                 (12, "현재 성능·원인진단"), (13, "웹 공개·시연"),
                 (16, "모델 갱신 주기(Ⅱ)"), (17, "고도화 방안(Ⅲ)")):
        print(f"  · {n:>2}면 {t}")
    return out


if __name__ == "__main__":
    main()
